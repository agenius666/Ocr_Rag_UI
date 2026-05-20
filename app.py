import os

os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")
os.environ.setdefault("HF_HUB_DISABLE_PROGRESS_BARS", "1")
os.environ.setdefault("DISABLE_SAFETENSORS_CONVERSION", "true")

import gc
import hashlib
import io
import json
import platform
import re
import shutil
import shlex
import sqlite3
import subprocess
import sys
import tempfile
import time
import uuid
import zipfile
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple

import fitz
import streamlit as st
from docx import Document
from dotenv import dotenv_values, load_dotenv
from openai import OpenAI
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from rag_utils import (
    keyword_rank_documents,
    parse_markdown_table,
    reciprocal_rank_merge,
    split_semantic_chunks,
)

# =========================
# 基础配置
# =========================
load_dotenv()

ENV_FILE = ".env"
APP_DB_FILE = "app_state.sqlite3"
UPLOAD_DIR = "uploads"
QDRANT_DIR = "qdrant_db"
COLLECTION_NAME = "ocr_rag_docs"
VECTOR_SIZE = 1024
EMBEDDING_MODEL_NAME = "BAAI/bge-m3"
RERANKER_MODEL_NAME = "BAAI/bge-reranker-v2-m3"
EXTRACTED_IMAGE_DIR = os.path.join(UPLOAD_DIR, "extracted_images")
CONVERTED_DIR = os.path.join(UPLOAD_DIR, "converted")
DEFAULT_MODEL_CACHE_ROOT = os.path.abspath("model_cache")
DEFAULT_PADDLEOCR_CACHE_DIR = os.path.join(DEFAULT_MODEL_CACHE_ROOT, "paddlex")
DEFAULT_BGE_CACHE_DIR = os.path.join(DEFAULT_MODEL_CACHE_ROOT, "bge-m3")
DEFAULT_RERANKER_CACHE_DIR = os.path.join(DEFAULT_MODEL_CACHE_ROOT, "bge-reranker-v2-m3")
DEFAULT_SOFFICE_BINARY_PATH = ""

DEFAULT_LLM_BASE_URL = "http://127.0.0.1:27292/v1"
DEFAULT_LLM_API_KEY = "EMPTY"
DEFAULT_LLM_MODEL = "local-model"
DEFAULT_LLM_EXTRA_BODY = "{}"
LLM_MODE_OPTIONS = {
    "快速": "fast",
    "思考": "thinking",
}
PDF_OCR_MODE_OPTIONS = {
    "智能 OCR（推荐，低内存）": "smart",
    "强制每页 OCR（最全但高内存）": "force",
    "仅提取 PDF 文字（最低内存）": "text",
}
PADDLEOCR_MODEL_OPTIONS = {
    "Server（高精度，占用更高）": {
        "det": "PP-OCRv5_server_det",
        "rec": "PP-OCRv5_server_rec",
    },
    "Mobile（低内存，速度更快）": {
        "det": "PP-OCRv5_mobile_det",
        "rec": "PP-OCRv5_mobile_rec",
    },
}
CHAT_PANEL_HEIGHT = 560
DEFAULT_CONTEXT_TURNS = 6
DEFAULT_RAG_TOP_K = 3
DEFAULT_COMPLIANCE_REGULATION_TOP_K = 6
DEFAULT_COMPLIANCE_ENTERPRISE_TOP_K = 8
DEFAULT_COMPLIANCE_MIN_REGULATION_EVIDENCE = 3
DEFAULT_COMPLIANCE_MIN_ENTERPRISE_EVIDENCE = 3
DEFAULT_VECTOR_MAX_DISTANCE = 0.85
DEFAULT_USE_HYBRID_SEARCH = True
DEFAULT_USE_RERANKER = False
DEFAULT_RETRIEVAL_FETCH_K = 20
DEFAULT_QUERY_DECOMPOSE = True
DEFAULT_BACKGROUND_INGEST = True
EMBEDDING_BATCH_SIZE = 8
VECTOR_ADD_BATCH_SIZE = 32
OCR_PDF_DPI = 120
OCR_MAX_PAGE_SIDE_PIXELS = 1800
DEFAULT_PADDLEOCR_MODEL_LABEL = "Server（高精度，占用更高）"
DEFAULT_UI_LANGUAGE = "en"

MODERN_SUPPORTED_TYPES = [
    "pdf",
    "png",
    "jpg",
    "jpeg",
    "webp",
    "bmp",
    "docx",
    "pptx",
    "xlsx",
]
LEGACY_OFFICE_TYPES = ["doc", "ppt", "xls"]
SUPPORTED_TYPES = MODERN_SUPPORTED_TYPES + LEGACY_OFFICE_TYPES
SUPPORTED_TYPE_LABEL = ", ".join(ext.upper() for ext in SUPPORTED_TYPES)

DOC_CATEGORY_OPTIONS = {
    "监管要求 / 规章制度": "regulation",
    "企业资料": "enterprise",
    "其他资料": "general",
}
DOC_CATEGORY_NAMES = {value: key for key, value in DOC_CATEGORY_OPTIONS.items()}
OCR_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}

os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(QDRANT_DIR, exist_ok=True)
os.makedirs(EXTRACTED_IMAGE_DIR, exist_ok=True)
os.makedirs(CONVERTED_DIR, exist_ok=True)
os.makedirs(DEFAULT_MODEL_CACHE_ROOT, exist_ok=True)

# =========================
# Streamlit 页面
# =========================
st.set_page_config(page_title="OCR RAG UI", layout="wide")

st.markdown(
    """
    <style>
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    [data-testid="stToolbar"] {visibility: hidden; height: 0; position: fixed;}
    [data-testid="stDecoration"] {display: none;}
    </style>
    """,
    unsafe_allow_html=True,
)

if "model_events" not in st.session_state:
    st.session_state["model_events"] = []


# =========================
# 应用状态数据库
# =========================
def current_timestamp() -> float:
    return time.time()


@st.cache_resource
def load_app_db() -> sqlite3.Connection:
    conn = sqlite3.connect(APP_DB_FILE, check_same_thread=False)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA foreign_keys = ON")
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS app_config (
            key TEXT PRIMARY KEY,
            value TEXT NOT NULL,
            updated_at REAL NOT NULL
        )
        """
    )
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS chat_sessions (
            id TEXT PRIMARY KEY,
            session_type TEXT NOT NULL,
            title TEXT NOT NULL,
            created_at REAL NOT NULL,
            updated_at REAL NOT NULL
        )
        """
    )
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS chat_messages (
            id TEXT PRIMARY KEY,
            session_id TEXT NOT NULL,
            role TEXT NOT NULL,
            content TEXT NOT NULL,
            payload TEXT NOT NULL,
            created_at REAL NOT NULL,
            FOREIGN KEY(session_id) REFERENCES chat_sessions(id) ON DELETE CASCADE
        )
        """
    )
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS ingested_files (
            sha256 TEXT PRIMARY KEY,
            file_name TEXT NOT NULL,
            doc_category TEXT NOT NULL,
            doc_label TEXT NOT NULL,
            chunk_count INTEGER NOT NULL,
            created_at REAL NOT NULL,
            updated_at REAL NOT NULL
        )
        """
    )
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS ingest_tasks (
            id TEXT PRIMARY KEY,
            status TEXT NOT NULL,
            total_files INTEGER NOT NULL,
            processed_files INTEGER NOT NULL,
            success_count INTEGER NOT NULL,
            duplicate_count INTEGER NOT NULL,
            skipped_count INTEGER NOT NULL,
            failed_count INTEGER NOT NULL,
            current_file TEXT NOT NULL,
            message TEXT NOT NULL,
            created_at REAL NOT NULL,
            updated_at REAL NOT NULL
        )
        """
    )
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS ingest_task_items (
            id TEXT PRIMARY KEY,
            task_id TEXT NOT NULL,
            file_name TEXT NOT NULL,
            status TEXT NOT NULL,
            message TEXT NOT NULL,
            chunk_count INTEGER NOT NULL,
            sha256 TEXT NOT NULL,
            created_at REAL NOT NULL,
            updated_at REAL NOT NULL,
            FOREIGN KEY(task_id) REFERENCES ingest_tasks(id) ON DELETE CASCADE
        )
        """
    )
    conn.execute(
        "CREATE INDEX IF NOT EXISTS idx_chat_sessions_type_updated "
        "ON chat_sessions(session_type, updated_at DESC)"
    )
    conn.execute(
        "CREATE INDEX IF NOT EXISTS idx_chat_messages_session_created "
        "ON chat_messages(session_id, created_at ASC)"
    )
    conn.execute(
        "CREATE INDEX IF NOT EXISTS idx_ingested_files_updated "
        "ON ingested_files(updated_at DESC)"
    )
    conn.execute(
        "CREATE INDEX IF NOT EXISTS idx_ingest_tasks_updated "
        "ON ingest_tasks(updated_at DESC)"
    )
    conn.execute(
        "CREATE INDEX IF NOT EXISTS idx_ingest_task_items_task "
        "ON ingest_task_items(task_id, updated_at DESC)"
    )
    conn.commit()
    return conn


app_db = load_app_db()


def get_config_value(key: str, default: str = "") -> str:
    row = app_db.execute(
        "SELECT value FROM app_config WHERE key = ?",
        (key,),
    ).fetchone()
    return str(row["value"]) if row else default


def set_config_value(key: str, value: Any) -> None:
    app_db.execute(
        """
        INSERT INTO app_config (key, value, updated_at)
        VALUES (?, ?, ?)
        ON CONFLICT(key) DO UPDATE SET
            value = excluded.value,
            updated_at = excluded.updated_at
        """,
        (key, str(value), current_timestamp()),
    )
    app_db.commit()


def get_int_config(key: str, default: int) -> int:
    try:
        return int(get_config_value(key, str(default)))
    except (TypeError, ValueError):
        return default


def get_float_config(key: str, default: float) -> float:
    try:
        return float(get_config_value(key, str(default)))
    except (TypeError, ValueError):
        return default


def get_bool_config(key: str, default: bool) -> bool:
    raw_value = get_config_value(key, "true" if default else "false").lower()
    return raw_value in {"1", "true", "yes", "on"}


def set_bool_config(key: str, value: bool) -> None:
    set_config_value(key, "true" if value else "false")


LANGUAGE_OPTIONS = {
    "English": "en",
    "简体中文": "zh_CN",
    "繁體中文": "zh_TW",
}
LANGUAGE_LABEL_BY_CODE = {value: key for key, value in LANGUAGE_OPTIONS.items()}

TRANSLATIONS = {
    "OCR + BGE-M3 + Qdrant + 本地大模型": {
        "en": "OCR + BGE-M3 + Qdrant + Local LLM",
        "zh_TW": "OCR + BGE-M3 + Qdrant + 本地大模型",
    },
    "上传制度、监管要求和企业资料，解析后写入 Qdrant 向量库，再调用 OpenAI 兼容接口做问答和合规差距分析。": {
        "en": "Upload policies, regulatory requirements, and enterprise materials; parse them into Qdrant, then answer questions and analyze compliance through an OpenAI-compatible local LLM endpoint.",
        "zh_TW": "上傳制度、監管要求和企業資料，解析後寫入 Qdrant 向量庫，再調用 OpenAI 相容接口做問答和合規差距分析。",
    },
    "上传入库": {"en": "Ingest", "zh_TW": "上傳入庫"},
    "检索问答": {"en": "RAG Chat", "zh_TW": "檢索問答"},
    "合规分析": {"en": "Compliance", "zh_TW": "合規分析"},
    "模型配置": {"en": "Settings", "zh_TW": "模型配置"},
    "模型状态": {"en": "Model Status", "zh_TW": "模型狀態"},
    "文档库管理": {"en": "Library", "zh_TW": "文件庫管理"},
    "上传文件并写入向量库": {"en": "Upload And Ingest Files", "zh_TW": "上傳文件並寫入向量庫"},
    "上传方式": {"en": "Upload Mode", "zh_TW": "上傳方式"},
    "文件 / 多文件": {"en": "File / Multiple Files", "zh_TW": "文件 / 多文件"},
    "文件夹（含子文件夹）": {"en": "Folder Including Subfolders", "zh_TW": "資料夾（含子資料夾）"},
    "选择文件夹": {"en": "Choose Folder", "zh_TW": "選擇資料夾"},
    "选择文件": {"en": "Choose Files", "zh_TW": "選擇文件"},
    "上传文件或文件夹": {"en": "Upload Files Or Folder", "zh_TW": "上傳文件或資料夾"},
    "资料类型": {"en": "Document Type", "zh_TW": "資料類型"},
    "监管要求 / 规章制度": {"en": "Regulations / Policies", "zh_TW": "監管要求 / 規章制度"},
    "企业资料": {"en": "Enterprise Materials", "zh_TW": "企業資料"},
    "其他资料": {"en": "Other Materials", "zh_TW": "其他資料"},
    "资料名称 / 备注": {"en": "Document Name / Note", "zh_TW": "資料名稱 / 備註"},
    "启用 Office 内嵌图片 OCR": {"en": "Enable OCR For Office Embedded Images", "zh_TW": "啟用 Office 內嵌圖片 OCR"},
    "入库完成后自动释放 OCR / BGE-M3 模型缓存": {"en": "Release OCR / BGE-M3 Model Cache After Ingestion", "zh_TW": "入庫完成後自動釋放 OCR / BGE-M3 模型快取"},
    "同名文件变更时替换旧版本": {"en": "Replace Old Version When Same-Name File Changes", "zh_TW": "同名文件變更時替換舊版本"},
    "后台入库队列": {"en": "Background Ingestion Queue", "zh_TW": "後台入庫隊列"},
    "PDF OCR 模式": {"en": "PDF OCR Mode", "zh_TW": "PDF OCR 模式"},
    "智能 OCR（推荐，低内存）": {"en": "Smart OCR (Recommended, Low Memory)", "zh_TW": "智慧 OCR（推薦，低記憶體）"},
    "强制每页 OCR（最全但高内存）": {"en": "Force OCR On Every Page (Complete, High Memory)", "zh_TW": "強制每頁 OCR（最完整但高記憶體）"},
    "仅提取 PDF 文字（最低内存）": {"en": "Extract PDF Text Only (Lowest Memory)", "zh_TW": "僅提取 PDF 文字（最低記憶體）"},
    "缺少 LibreOffice 时自动下载安装转换工具": {"en": "Auto Install LibreOffice If Missing", "zh_TW": "缺少 LibreOffice 時自動下載安裝轉換工具"},
    "Chunk 大小": {"en": "Chunk Size", "zh_TW": "Chunk 大小"},
    "Chunk 重叠": {"en": "Chunk Overlap", "zh_TW": "Chunk 重疊"},
    "最近入库任务": {"en": "Recent Ingestion Tasks", "zh_TW": "最近入庫任務"},
    "任务控制": {"en": "Task Controls", "zh_TW": "任務控制"},
    "任务列表": {"en": "Task List", "zh_TW": "任務列表"},
    "最近任务文件明细": {"en": "Latest Task File Details", "zh_TW": "最近任務文件明細"},
    "开始导入文件": {"en": "Start Import", "zh_TW": "開始導入文件"},
    "暂停": {"en": "Pause", "zh_TW": "暫停"},
    "继续": {"en": "Resume", "zh_TW": "繼續"},
    "终止": {"en": "Stop", "zh_TW": "終止"},
    "多轮检索问答": {"en": "Multi-Turn RAG Chat", "zh_TW": "多輪檢索問答"},
    "对话与检索设置": {"en": "Chat And Retrieval Settings", "zh_TW": "對話與檢索設定"},
    "检索范围": {"en": "Search Scope", "zh_TW": "檢索範圍"},
    "全部资料": {"en": "All Materials", "zh_TW": "全部資料"},
    "召回片段数量": {"en": "Top-K Chunks", "zh_TW": "召回片段數量"},
    "回答模式": {"en": "Answer Mode", "zh_TW": "回答模式"},
    "快速": {"en": "Fast", "zh_TW": "快速"},
    "思考": {"en": "Thinking", "zh_TW": "思考"},
    "上下文轮数": {"en": "Context Turns", "zh_TW": "上下文輪數"},
    "未检索到资料时使用本地大模型普通聊天": {"en": "Use Local LLM Chat When No Evidence Is Retrieved", "zh_TW": "未檢索到資料時使用本地大模型普通聊天"},
    "追问补全成完整检索问题": {"en": "Rewrite Follow-Ups Into Complete Search Queries", "zh_TW": "追問補全成完整檢索問題"},
    "复杂问题拆解后分别检索": {"en": "Decompose Complex Questions Before Search", "zh_TW": "複雜問題拆解後分別檢索"},
    "启用向量距离阈值过滤": {"en": "Enable Vector Distance Threshold", "zh_TW": "啟用向量距離閾值過濾"},
    "最大距离": {"en": "Max Distance", "zh_TW": "最大距離"},
    "启用混合检索": {"en": "Enable Hybrid Search", "zh_TW": "啟用混合檢索"},
    "启用重排模型": {"en": "Enable Reranker", "zh_TW": "啟用重排模型"},
    "候选召回数": {"en": "Candidate Count", "zh_TW": "候選召回數"},
    "清空当前对话": {"en": "Clear Current Chat", "zh_TW": "清空當前對話"},
    "输入问题": {"en": "Enter Question", "zh_TW": "輸入問題"},
    "发送": {"en": "Send", "zh_TW": "發送"},
    "请输入问题。": {"en": "Please enter a question.", "zh_TW": "請輸入問題。"},
    "多轮合规差距分析": {"en": "Multi-Turn Compliance Gap Analysis", "zh_TW": "多輪合規差距分析"},
    "分析与检索设置": {"en": "Analysis And Retrieval Settings", "zh_TW": "分析與檢索設定"},
    "召回监管 / 规章片段数量": {"en": "Regulation Top-K Chunks", "zh_TW": "召回監管 / 規章片段數量"},
    "召回企业资料片段数量": {"en": "Enterprise Top-K Chunks", "zh_TW": "召回企業資料片段數量"},
    "分析模式": {"en": "Analysis Mode", "zh_TW": "分析模式"},
    "监管最少证据数": {"en": "Minimum Regulation Evidence", "zh_TW": "監管最少證據數"},
    "企业最少证据数": {"en": "Minimum Enterprise Evidence", "zh_TW": "企業最少證據數"},
    "按监管条款逐条对照": {"en": "Compare Clause By Clause", "zh_TW": "按監管條款逐條對照"},
    "输出资料不足清单": {"en": "Output Missing Materials List", "zh_TW": "輸出資料不足清單"},
    "清空合规分析对话": {"en": "Clear Compliance Chat", "zh_TW": "清空合規分析對話"},
    "输入合规分析问题": {"en": "Enter Compliance Analysis Question", "zh_TW": "輸入合規分析問題"},
    "请输入分析主题。": {"en": "Please enter an analysis topic.", "zh_TW": "請輸入分析主題。"},
    "配置中心": {"en": "Settings", "zh_TW": "配置中心"},
    "模型与路径": {"en": "Models And Paths", "zh_TW": "模型與路徑"},
    "本地大模型": {"en": "Local LLM", "zh_TW": "本地大模型"},
    "初始化": {"en": "Reset", "zh_TW": "初始化"},
    "OCR 模型": {"en": "OCR Model", "zh_TW": "OCR 模型"},
    "PaddleOCR 模型": {"en": "PaddleOCR Model", "zh_TW": "PaddleOCR 模型"},
    "Server（高精度，占用更高）": {"en": "Server (Higher Accuracy, More Memory)", "zh_TW": "Server（高精度，佔用更高）"},
    "Mobile（低内存，速度更快）": {"en": "Mobile (Lower Memory, Faster)", "zh_TW": "Mobile（低記憶體，速度更快）"},
    "模型保存路径": {"en": "Model Storage Paths", "zh_TW": "模型保存路徑"},
    "默认模型根目录": {"en": "Default Model Root", "zh_TW": "預設模型根目錄"},
    "PaddleOCR 模型目录": {"en": "PaddleOCR Model Directory", "zh_TW": "PaddleOCR 模型目錄"},
    "BAAI/bge-m3 模型目录": {"en": "BAAI/bge-m3 Model Directory", "zh_TW": "BAAI/bge-m3 模型目錄"},
    "BAAI/bge-reranker-v2-m3 模型目录": {"en": "BAAI/bge-reranker-v2-m3 Model Directory", "zh_TW": "BAAI/bge-reranker-v2-m3 模型目錄"},
    "LibreOffice / soffice 路径": {"en": "LibreOffice / soffice Path", "zh_TW": "LibreOffice / soffice 路徑"},
    "保存模型路径": {"en": "Save Model Paths", "zh_TW": "保存模型路徑"},
    "恢复默认路径": {"en": "Restore Default Paths", "zh_TW": "恢復預設路徑"},
    "当前路径": {"en": "Current Paths", "zh_TW": "當前路徑"},
    "接口": {"en": "Endpoint", "zh_TW": "接口"},
    "OpenAI 兼容接口 Base URL": {"en": "OpenAI-Compatible Base URL", "zh_TW": "OpenAI 相容接口 Base URL"},
    "模型名称": {"en": "Model Name", "zh_TW": "模型名稱"},
    "默认模型名称": {"en": "Default Model Name", "zh_TW": "預設模型名稱"},
    "模式": {"en": "Modes", "zh_TW": "模式"},
    "快速模式模型名": {"en": "Fast Mode Model", "zh_TW": "快速模式模型名"},
    "思考模式模型名": {"en": "Thinking Mode Model", "zh_TW": "思考模式模型名"},
    "快速模式 extra_body JSON": {"en": "Fast Mode extra_body JSON", "zh_TW": "快速模式 extra_body JSON"},
    "思考模式 extra_body JSON": {"en": "Thinking Mode extra_body JSON", "zh_TW": "思考模式 extra_body JSON"},
    "保存配置": {"en": "Save Settings", "zh_TW": "保存配置"},
    "恢复默认值": {"en": "Restore Defaults", "zh_TW": "恢復預設值"},
    "测试模式": {"en": "Test Mode", "zh_TW": "測試模式"},
    "测试当前配置": {"en": "Test Current Settings", "zh_TW": "測試當前配置"},
    "当前生效配置": {"en": "Active Settings", "zh_TW": "當前生效配置"},
    "初始化可配置数据": {"en": "Reset Configurable Data", "zh_TW": "初始化可配置資料"},
    "我确认要初始化可配置数据和历史会话": {"en": "I confirm resetting configurable data and chat history", "zh_TW": "我確認要初始化可配置資料和歷史會話"},
    "模型状态 / 下载查询": {"en": "Model Status / Download Check", "zh_TW": "模型狀態 / 下載查詢"},
    "操作": {"en": "Actions", "zh_TW": "操作"},
    "内存管理": {"en": "Memory Management", "zh_TW": "記憶體管理"},
    "释放 OCR / BGE-M3 / Reranker 模型缓存": {"en": "Release OCR / BGE-M3 / Reranker Cache", "zh_TW": "釋放 OCR / BGE-M3 / Reranker 模型快取"},
    "预加载 PaddleOCR": {"en": "Preload PaddleOCR", "zh_TW": "預載 PaddleOCR"},
    "预加载 BGE-M3": {"en": "Preload BGE-M3", "zh_TW": "預載 BGE-M3"},
    "预加载 Reranker": {"en": "Preload Reranker", "zh_TW": "預載 Reranker"},
    "Office 老格式转换": {"en": "Legacy Office Conversion", "zh_TW": "Office 舊格式轉換"},
    "测试 LibreOffice": {"en": "Test LibreOffice", "zh_TW": "測試 LibreOffice"},
    "自动安装 LibreOffice": {"en": "Auto Install LibreOffice", "zh_TW": "自動安裝 LibreOffice"},
    "测试本地大模型": {"en": "Test Local LLM", "zh_TW": "測試本地大模型"},
    "最近模型事件": {"en": "Recent Model Events", "zh_TW": "最近模型事件"},
    "文档库管理": {"en": "Document Library", "zh_TW": "文件庫管理"},
    "刷新文件摘要": {"en": "Refresh File Summary", "zh_TW": "刷新文件摘要"},
    "查看去重记录": {"en": "View Deduplication Records", "zh_TW": "查看去重記錄"},
    "向量库备份 / 导入 / 导出": {"en": "Vector Store Backup / Import / Export", "zh_TW": "向量庫備份 / 導入 / 導出"},
    "导出文档库备份": {"en": "Export Library Backup", "zh_TW": "導出文件庫備份"},
    "导入备份 ZIP": {"en": "Import Backup ZIP", "zh_TW": "導入備份 ZIP"},
    "我确认导入备份并覆盖当前文档库": {"en": "I confirm importing backup and overwriting the current library", "zh_TW": "我確認導入備份並覆蓋當前文件庫"},
    "导入并覆盖当前文档库": {"en": "Import And Overwrite Current Library", "zh_TW": "導入並覆蓋當前文件庫"},
    "清空 Qdrant 向量库": {"en": "Clear Qdrant Vector Store", "zh_TW": "清空 Qdrant 向量庫"},
    "总 chunk": {"en": "Total Chunks", "zh_TW": "總 chunk"},
    "监管/规章": {"en": "Regulations", "zh_TW": "監管/規章"},
    "新建会话": {"en": "New Chat", "zh_TW": "新建會話"},
    "删除会话": {"en": "Delete Chat", "zh_TW": "刪除會話"},
    "本轮检索资料": {"en": "Retrieved Materials", "zh_TW": "本輪檢索資料"},
    "本轮证据": {"en": "Evidence", "zh_TW": "本輪證據"},
    "监管 / 规章": {"en": "Regulations / Policies", "zh_TW": "監管 / 規章"},
    "无监管证据": {"en": "No regulation evidence", "zh_TW": "無監管證據"},
    "无企业证据": {"en": "No enterprise evidence", "zh_TW": "無企業證據"},
    "结构化合规分析表": {"en": "Structured Compliance Analysis Table", "zh_TW": "結構化合規分析表"},
    "导出本轮合规分析 Excel": {"en": "Export This Analysis To Excel", "zh_TW": "導出本輪合規分析 Excel"},
    "当前 Qdrant Collection：": {"en": "Current Qdrant Collection:", "zh_TW": "當前 Qdrant Collection："},
    "当前文档库为空。": {"en": "The document library is empty.", "zh_TW": "當前文件庫為空。"},
    "暂无去重记录。": {"en": "No deduplication records yet.", "zh_TW": "暫無去重記錄。"},
    "文件": {"en": "File", "zh_TW": "文件"},
    "状态": {"en": "Status", "zh_TW": "狀態"},
    "进度": {"en": "Progress", "zh_TW": "進度"},
    "成功": {"en": "Succeeded", "zh_TW": "成功"},
    "重复": {"en": "Duplicate", "zh_TW": "重複"},
    "跳过": {"en": "Skipped", "zh_TW": "跳過"},
    "失败": {"en": "Failed", "zh_TW": "失敗"},
    "当前文件": {"en": "Current File", "zh_TW": "當前文件"},
    "说明": {"en": "Message", "zh_TW": "說明"},
    "更新时间": {"en": "Updated At", "zh_TW": "更新時間"},
    "原因": {"en": "Reason", "zh_TW": "原因"},
    "已入库文件": {"en": "Ingested File", "zh_TW": "已入庫文件"},
    "chunk 数": {"en": "Chunk Count", "zh_TW": "chunk 數"},
    "来源格式": {"en": "Source Format", "zh_TW": "來源格式"},
    "资料名称": {"en": "Document Name", "zh_TW": "資料名稱"},
    "入库时间": {"en": "Ingested At", "zh_TW": "入庫時間"},
    "组件": {"en": "Component", "zh_TW": "組件"},
    "用途": {"en": "Purpose", "zh_TW": "用途"},
    "time": {"en": "Time", "zh_TW": "時間"},
    "component": {"en": "Component", "zh_TW": "組件"},
    "status": {"en": "Status", "zh_TW": "狀態"},
    "detail": {"en": "Detail", "zh_TW": "詳情"},
    "运行中": {"en": "Running", "zh_TW": "運行中"},
    "暂停中": {"en": "Pausing", "zh_TW": "暫停中"},
    "已暂停": {"en": "Paused", "zh_TW": "已暫停"},
    "终止中": {"en": "Stopping", "zh_TW": "終止中"},
    "已终止": {"en": "Stopped", "zh_TW": "已終止"},
    "已完成": {"en": "Completed", "zh_TW": "已完成"},
    "完成": {"en": "Completed", "zh_TW": "完成"},
    "未知文件": {"en": "Unknown File", "zh_TW": "未知文件"},
    "未知类型": {"en": "Unknown Type", "zh_TW": "未知類型"},
    "未知": {"en": "Unknown", "zh_TW": "未知"},
    "自动搜索系统路径": {"en": "Auto-detect system path", "zh_TW": "自動搜尋系統路徑"},
    "未检测到": {"en": "Not detected", "zh_TW": "未檢測到"},
    "语言设置": {"en": "Language", "zh_TW": "語言設定"},
    "界面语言": {"en": "UI Language", "zh_TW": "介面語言"},
    "选择界面语言": {"en": "Choose UI Language", "zh_TW": "選擇介面語言"},
    "语言设置已保存。": {"en": "Language setting saved.", "zh_TW": "語言設定已保存。"},
    "语言偏好会保存到 app_state.sqlite3，下次打开会自动生效。": {
        "en": "Language preference is saved to app_state.sqlite3 and will be applied automatically next time.",
        "zh_TW": "語言偏好會保存到 app_state.sqlite3，下次打開會自動生效。",
    },
    "当前会话": {"en": "Current Chat", "zh_TW": "當前會話"},
    "输入问题后会保留上下文；每一轮都会按当前检索设置重新召回资料。": {
        "en": "Context is preserved after you ask a question; each turn retrieves materials using the current retrieval settings.",
        "zh_TW": "輸入問題後會保留上下文；每一輪都會按當前檢索設定重新召回資料。",
    },
    "输入合规分析问题后会保留上下文；每一轮都会分别检索监管资料和企业资料。": {
        "en": "Context is preserved after you ask a compliance question; each turn retrieves regulation and enterprise materials separately.",
        "zh_TW": "輸入合規分析問題後會保留上下文；每一輪都會分別檢索監管資料和企業資料。",
    },
    "批量上传时留空会使用各自文件名": {
        "en": "Leave empty to use each file name for batch upload",
        "zh_TW": "批量上傳時留空會使用各自文件名",
    },
    "没有鉴权时可填 EMPTY": {"en": "Use EMPTY if no authentication is required", "zh_TW": "沒有鑑權時可填 EMPTY"},
    "填写 OLMX / Ollama / LM Studio 中显示的真实模型名": {
        "en": "Use the actual model name shown by OLMX / Ollama / LM Studio",
        "zh_TW": "填寫 OLMX / Ollama / LM Studio 中顯示的真實模型名",
    },
    "留空则使用默认模型名称": {"en": "Leave empty to use the default model name", "zh_TW": "留空則使用預設模型名稱"},
    "如果后端用不同模型区分快慢，可在这里填思考模型名": {
        "en": "If your backend uses a different model for thinking mode, enter it here",
        "zh_TW": "如果後端用不同模型區分快慢，可在這裡填思考模型名",
    },
    "留空则自动搜索系统路径": {"en": "Leave empty to auto-detect system path", "zh_TW": "留空則自動搜尋系統路徑"},
    "当前检索范围没有任何入库 chunk，请先上传资料。": {
        "en": "The current search scope has no ingested chunks. Please upload materials first.",
        "zh_TW": "當前檢索範圍沒有任何入庫 chunk，請先上傳資料。",
    },
    "没有检索到满足当前距离阈值的相关内容。可以在检索设置里调大“最大距离”或暂时关闭阈值过滤。": {
        "en": "No relevant content matched the current distance threshold. Increase Max Distance or temporarily disable threshold filtering in retrieval settings.",
        "zh_TW": "沒有檢索到滿足當前距離閾值的相關內容。可以在檢索設定裡調大「最大距離」或暫時關閉閾值過濾。",
    },
    "没有检索到相关内容。": {"en": "No relevant content was retrieved.", "zh_TW": "沒有檢索到相關內容。"},
    "请先上传并入库监管要求或规章制度。": {
        "en": "Please upload and ingest regulations or policy documents first.",
        "zh_TW": "請先上傳並入庫監管要求或規章制度。",
    },
    "请先上传并入库企业资料。": {
        "en": "Please upload and ingest enterprise materials first.",
        "zh_TW": "請先上傳並入庫企業資料。",
    },
    "没有检索到相关监管要求。": {"en": "No relevant regulation requirements were retrieved.", "zh_TW": "沒有檢索到相關監管要求。"},
    "没有检索到相关企业资料。": {"en": "No relevant enterprise materials were retrieved.", "zh_TW": "沒有檢索到相關企業資料。"},
    "本次没有文件成功入库。": {"en": "No files were ingested successfully this time.", "zh_TW": "本次沒有文件成功入庫。"},
    "当前向量库为空。": {"en": "The current vector store is empty.", "zh_TW": "當前向量庫為空。"},
}

PHRASE_TRANSLATIONS = {
    "支持格式：": {"en": "Supported formats: ", "zh_TW": "支援格式："},
    "文件夹模式会包含子文件夹中的文件。": {
        "en": "Folder mode includes files in subfolders.",
        "zh_TW": "資料夾模式會包含子資料夾中的文件。",
    },
    "不支持的文件会在批量处理结果中列出。": {
        "en": "Unsupported files will be listed in the batch results.",
        "zh_TW": "不支援的文件會在批量處理結果中列出。",
    },
    "当前会话共 ": {"en": "Current chat has ", "zh_TW": "當前會話共 "},
    " 条消息": {"en": " messages", "zh_TW": " 條訊息"},
    "模式：": {"en": "Mode: ", "zh_TW": "模式："},
    "本轮检索问题：": {"en": "Search query: ", "zh_TW": "本輪檢索問題："},
    "拆解子问题：": {"en": "Sub-queries: ", "zh_TW": "拆解子問題："},
    "问题改写失败，已使用原问题检索：": {
        "en": "Query rewrite failed; using the original question: ",
        "zh_TW": "問題改寫失敗，已使用原問題檢索：",
    },
    "问题拆解失败，已使用启发式拆解：": {
        "en": "Query decomposition failed; using heuristic decomposition: ",
        "zh_TW": "問題拆解失敗，已使用啟發式拆解：",
    },
    "资料 ": {"en": "Material ", "zh_TW": "資料 "},
    "片段 ": {"en": "Chunk ", "zh_TW": "片段 "},
    "距离 ": {"en": "Distance ", "zh_TW": "距離 "},
    "来源 ": {"en": "Source ", "zh_TW": "來源 "},
    "重排 ": {"en": "Rerank ", "zh_TW": "重排 "},
    "覆盖补充": {"en": "coverage supplement", "zh_TW": "覆蓋補充"},
    "合并 ": {"en": "Merged ", "zh_TW": "合併 "},
    "已选择 ": {"en": "Selected ", "zh_TW": "已選擇 "},
    " 个文件，其中当前可处理 ": {"en": " files; processable now: ", "zh_TW": " 個文件，其中當前可處理 "},
    " 个，待跳过/不支持 ": {"en": "; to skip/unsupported: ", "zh_TW": " 個，待跳過/不支援 "},
    " 个。开始处理后会按 SHA256 自动跳过重复文件。": {
        "en": ". Duplicate files will be skipped automatically by SHA256.",
        "zh_TW": " 個。開始處理後會按 SHA256 自動跳過重複文件。",
    },
    "检测到老版 Office 文件，将使用 LibreOffice 转换：": {
        "en": "Legacy Office files detected; LibreOffice will convert them: ",
        "zh_TW": "檢測到舊版 Office 文件，將使用 LibreOffice 轉換：",
    },
    "当前 PaddleOCR 模型：": {"en": "Current PaddleOCR model: ", "zh_TW": "當前 PaddleOCR 模型："},
    "当前系统：": {"en": "Current system: ", "zh_TW": "當前系統："},
    "安装命令：": {"en": "Install command: ", "zh_TW": "安裝命令："},
    "正在处理 ": {"en": "Processing ", "zh_TW": "正在處理 "},
    "正在处理：": {"en": "Processing: ", "zh_TW": "正在處理："},
    "已跳过：": {"en": "Skipped: ", "zh_TW": "已跳過："},
    "已完成：": {"en": "Completed: ", "zh_TW": "已完成："},
    "已跳过重复文件：": {"en": "Skipped duplicate file: ", "zh_TW": "已跳過重複文件："},
    "已跳过不支持文件：": {"en": "Skipped unsupported file: ", "zh_TW": "已跳過不支援文件："},
    "已跳过无法转换文件：": {"en": "Skipped unconvertible file: ", "zh_TW": "已跳過無法轉換文件："},
    "没有解析到有效文字": {"en": "No valid text was extracted", "zh_TW": "沒有解析到有效文字"},
    "解析成功但没有可入库 chunk": {"en": "Parsed successfully, but no ingestible chunks were produced", "zh_TW": "解析成功但沒有可入庫 chunk"},
    "入库成功，写入 ": {"en": "Ingested successfully; wrote ", "zh_TW": "入庫成功，寫入 "},
    " 个 chunk": {"en": " chunks", "zh_TW": " 個 chunk"},
    "成功入库 ": {"en": "Successfully ingested ", "zh_TW": "成功入庫 "},
    "不支持 ": {"en": "Unsupported files: ", "zh_TW": "不支援 "},
    "重复文件 ": {"en": "Duplicate files: ", "zh_TW": "重複文件 "},
    "跳过 ": {"en": "Skipped files: ", "zh_TW": "跳過 "},
    "处理失败 ": {"en": "Failed files: ", "zh_TW": "處理失敗 "},
    " 个文件。": {"en": " files.", "zh_TW": " 個文件。"},
    "已跳过。": {"en": " skipped.", "zh_TW": "已跳過。"},
    "接口地址：": {"en": "Endpoint: ", "zh_TW": "接口地址："},
    "模型名称：": {"en": "Model name: ", "zh_TW": "模型名稱："},
    "接口可用：": {"en": "Endpoint available: ", "zh_TW": "接口可用："},
    "本地大模型可用：": {"en": "Local LLM available: ", "zh_TW": "本地大模型可用："},
    "已删除 ": {"en": "Deleted ", "zh_TW": "已刪除 "},
    " 个 chunk。": {"en": " chunks.", "zh_TW": " 個 chunk。"},
}


def get_ui_language() -> str:
    language = get_config_value("ui_language", DEFAULT_UI_LANGUAGE)
    if language not in LANGUAGE_LABEL_BY_CODE:
        language = DEFAULT_UI_LANGUAGE
    return language


def translate_text(text: Any) -> Any:
    if not isinstance(text, str):
        return text
    language = get_ui_language()
    if language == "zh_CN":
        return text
    heading_match = re.match(r"^(#{1,6}\s+)(.+)$", text)
    if heading_match:
        return heading_match.group(1) + translate_text(heading_match.group(2))
    translated = TRANSLATIONS.get(text, {}).get(language)
    if translated:
        return translated
    if any("\u4e00" <= char <= "\u9fff" for char in text):
        translated_text = text
        for source, targets in sorted(PHRASE_TRANSLATIONS.items(), key=lambda item: len(item[0]), reverse=True):
            translated_text = translated_text.replace(source, targets.get(language, source))
        return translated_text
    return text


def translate_options(options):
    return [translate_text(option) for option in options]


def translate_kwargs(kwargs: Dict[str, Any]) -> Dict[str, Any]:
    translated = dict(kwargs)
    for key in ["help", "placeholder", "text", "label"]:
        if key in translated:
            translated[key] = translate_text(translated[key])
    return translated


def translate_table_data(data: Any) -> Any:
    if isinstance(data, list):
        translated_rows = []
        for row in data:
            if isinstance(row, dict):
                translated_rows.append(
                    {
                        translate_text(key): translate_text(value)
                        for key, value in row.items()
                    }
                )
            else:
                translated_rows.append(row)
        return translated_rows
    if isinstance(data, dict):
        return {
            translate_text(key): translate_text(value)
            for key, value in data.items()
        }
    return data


def patch_streamlit_i18n() -> None:
    if getattr(st, "_ocr_rag_i18n_patched", False):
        return

    st._ocr_rag_originals = {
        name: getattr(st, name)
        for name in [
            "title",
            "caption",
            "subheader",
            "markdown",
            "button",
            "checkbox",
            "radio",
            "selectbox",
            "text_input",
            "text_area",
            "file_uploader",
            "form_submit_button",
            "download_button",
            "info",
            "warning",
            "error",
            "success",
            "metric",
            "expander",
            "tabs",
            "slider",
            "status",
            "spinner",
            "progress",
            "dataframe",
            "table",
        ]
    }

    def wrap_label_method(name):
        original = st._ocr_rag_originals[name]

        def wrapped(label, *args, **kwargs):
            return original(translate_text(label), *args, **translate_kwargs(kwargs))

        return wrapped

    for method_name in [
        "title",
        "caption",
        "subheader",
        "markdown",
        "button",
        "checkbox",
        "text_input",
        "text_area",
        "file_uploader",
        "form_submit_button",
        "info",
        "warning",
        "error",
        "success",
        "metric",
        "expander",
        "slider",
        "spinner",
    ]:
        setattr(st, method_name, wrap_label_method(method_name))

    def wrapped_status(label, *args, **kwargs):
        status_object = st._ocr_rag_originals["status"](translate_text(label), *args, **translate_kwargs(kwargs))
        original_update = status_object.update

        def translated_update(*update_args, **update_kwargs):
            if "label" in update_kwargs:
                update_kwargs["label"] = translate_text(update_kwargs["label"])
            elif update_args:
                update_args = (translate_text(update_args[0]), *update_args[1:])
            return original_update(*update_args, **update_kwargs)

        status_object.update = translated_update
        return status_object

    def wrapped_progress(value, *args, **kwargs):
        kwargs = translate_kwargs(kwargs)
        progress_object = st._ocr_rag_originals["progress"](value, *args, **kwargs)
        original_progress = progress_object.progress

        def translated_progress(progress_value, *progress_args, **progress_kwargs):
            progress_kwargs = translate_kwargs(progress_kwargs)
            return original_progress(progress_value, *progress_args, **progress_kwargs)

        progress_object.progress = translated_progress
        return progress_object

    def wrapped_download_button(label, *args, **kwargs):
        return st._ocr_rag_originals["download_button"](translate_text(label), *args, **translate_kwargs(kwargs))

    def wrapped_radio(label, options, *args, **kwargs):
        kwargs = translate_kwargs(kwargs)
        existing_format_func = kwargs.get("format_func")
        if existing_format_func is None:
            kwargs["format_func"] = lambda option: translate_text(option)
        else:
            kwargs["format_func"] = lambda option: translate_text(existing_format_func(option))
        return st._ocr_rag_originals["radio"](translate_text(label), options, *args, **kwargs)

    def wrapped_selectbox(label, options, *args, **kwargs):
        kwargs = translate_kwargs(kwargs)
        existing_format_func = kwargs.get("format_func")
        if existing_format_func is None:
            kwargs["format_func"] = lambda option: translate_text(option)
        else:
            kwargs["format_func"] = lambda option: translate_text(existing_format_func(option))
        return st._ocr_rag_originals["selectbox"](translate_text(label), options, *args, **kwargs)

    def wrapped_tabs(tabs, *args, **kwargs):
        return st._ocr_rag_originals["tabs"](translate_options(tabs), *args, **kwargs)

    def wrapped_dataframe(data=None, *args, **kwargs):
        return st._ocr_rag_originals["dataframe"](translate_table_data(data), *args, **kwargs)

    def wrapped_table(data=None, *args, **kwargs):
        return st._ocr_rag_originals["table"](translate_table_data(data), *args, **kwargs)

    st.download_button = wrapped_download_button
    st.status = wrapped_status
    st.progress = wrapped_progress
    st.radio = wrapped_radio
    st.selectbox = wrapped_selectbox
    st.tabs = wrapped_tabs
    st.dataframe = wrapped_dataframe
    st.table = wrapped_table
    st._ocr_rag_i18n_patched = True


if not get_config_value("ui_language", ""):
    set_config_value("ui_language", DEFAULT_UI_LANGUAGE)

patch_streamlit_i18n()

st.title("OCR + BGE-M3 + Qdrant + 本地大模型")
st.caption("上传制度、监管要求和企业资料，解析后写入 Qdrant 向量库，再调用 OpenAI 兼容接口做问答和合规差距分析。")


def normalize_local_path(path_value: str, default: str = "") -> str:
    path_value = (path_value or "").strip()
    if not path_value:
        path_value = default
    if not path_value:
        return ""
    return os.path.abspath(os.path.expanduser(path_value))


def delete_all_config_values() -> None:
    app_db.execute("DELETE FROM app_config")
    app_db.commit()


def get_ingested_file(file_sha256: str) -> Optional[Dict[str, Any]]:
    row = app_db.execute(
        """
        SELECT sha256, file_name, doc_category, doc_label, chunk_count, created_at, updated_at
        FROM ingested_files
        WHERE sha256 = ?
        """,
        (file_sha256,),
    ).fetchone()
    return dict(row) if row else None


def list_ingested_files_by_name(file_name: str) -> List[Dict[str, Any]]:
    rows = app_db.execute(
        """
        SELECT sha256, file_name, doc_category, doc_label, chunk_count, created_at, updated_at
        FROM ingested_files
        WHERE file_name = ?
        ORDER BY updated_at DESC
        """,
        (file_name,),
    ).fetchall()
    return [dict(row) for row in rows]


def delete_ingested_file_records_by_name(file_name: str) -> int:
    cursor = app_db.execute("DELETE FROM ingested_files WHERE file_name = ?", (file_name,))
    app_db.commit()
    return cursor.rowcount


def record_ingested_file(
    file_sha256: str,
    file_name: str,
    doc_category: str,
    doc_label: str,
    chunk_count: int,
) -> None:
    now = current_timestamp()
    app_db.execute(
        """
        INSERT INTO ingested_files
            (sha256, file_name, doc_category, doc_label, chunk_count, created_at, updated_at)
        VALUES (?, ?, ?, ?, ?, ?, ?)
        ON CONFLICT(sha256) DO UPDATE SET
            file_name = excluded.file_name,
            doc_category = excluded.doc_category,
            doc_label = excluded.doc_label,
            chunk_count = excluded.chunk_count,
            updated_at = excluded.updated_at
        """,
        (
            file_sha256,
            file_name,
            doc_category,
            doc_label,
            chunk_count,
            now,
            now,
        ),
    )
    app_db.commit()


def list_ingested_files() -> List[Dict[str, Any]]:
    rows = app_db.execute(
        """
        SELECT sha256, file_name, doc_category, doc_label, chunk_count, created_at, updated_at
        FROM ingested_files
        ORDER BY updated_at DESC
        """
    ).fetchall()
    return [dict(row) for row in rows]


def default_session_title(session_type: str) -> str:
    label = "检索问答" if session_type == "rag" else "合规分析"
    return f"新建{label} {time.strftime('%m-%d %H:%M')}"


def create_chat_session(session_type: str, title: Optional[str] = None) -> str:
    session_id = uuid.uuid4().hex
    now = current_timestamp()
    app_db.execute(
        """
        INSERT INTO chat_sessions (id, session_type, title, created_at, updated_at)
        VALUES (?, ?, ?, ?, ?)
        """,
        (session_id, session_type, title or default_session_title(session_type), now, now),
    )
    app_db.commit()
    return session_id


def list_chat_sessions(session_type: str) -> List[Dict[str, Any]]:
    rows = app_db.execute(
        """
        SELECT id, session_type, title, created_at, updated_at
        FROM chat_sessions
        WHERE session_type = ?
        ORDER BY updated_at DESC
        """,
        (session_type,),
    ).fetchall()
    return [dict(row) for row in rows]


def get_active_session_id(session_type: str) -> str:
    state_key = f"active_{session_type}_session_id"
    sessions = list_chat_sessions(session_type)
    session_ids = {session["id"] for session in sessions}
    active_id = st.session_state.get(state_key)
    if active_id in session_ids:
        return active_id

    if sessions:
        active_id = sessions[0]["id"]
    else:
        active_id = create_chat_session(session_type)

    st.session_state[state_key] = active_id
    return active_id


def set_active_session_id(session_type: str, session_id: str) -> None:
    st.session_state[f"active_{session_type}_session_id"] = session_id


def delete_chat_session(session_id: str) -> None:
    app_db.execute("DELETE FROM chat_sessions WHERE id = ?", (session_id,))
    app_db.commit()


def clear_chat_session(session_id: str) -> None:
    now = current_timestamp()
    app_db.execute("DELETE FROM chat_messages WHERE session_id = ?", (session_id,))
    app_db.execute(
        "UPDATE chat_sessions SET updated_at = ? WHERE id = ?",
        (now, session_id),
    )
    app_db.commit()


def get_chat_messages(session_id: str) -> List[Dict[str, Any]]:
    rows = app_db.execute(
        """
        SELECT role, content, payload, created_at
        FROM chat_messages
        WHERE session_id = ?
        ORDER BY created_at ASC
        """,
        (session_id,),
    ).fetchall()
    messages = []
    for row in rows:
        try:
            payload = json.loads(row["payload"]) if row["payload"] else {}
        except json.JSONDecodeError:
            payload = {}
        messages.append(
            {
                "role": row["role"],
                "content": row["content"],
                **payload,
            }
        )
    return messages


def append_chat_message(
    session_id: str,
    role: str,
    content: str,
    payload: Optional[Dict[str, Any]] = None,
) -> None:
    now = current_timestamp()
    app_db.execute(
        """
        INSERT INTO chat_messages (id, session_id, role, content, payload, created_at)
        VALUES (?, ?, ?, ?, ?, ?)
        """,
        (
            uuid.uuid4().hex,
            session_id,
            role,
            content,
            json.dumps(payload or {}, ensure_ascii=False, default=str),
            now,
        ),
    )
    app_db.execute(
        "UPDATE chat_sessions SET updated_at = ? WHERE id = ?",
        (now, session_id),
    )
    app_db.commit()


def maybe_update_session_title(session_id: str, first_user_message: str) -> None:
    row = app_db.execute(
        "SELECT title FROM chat_sessions WHERE id = ?",
        (session_id,),
    ).fetchone()
    if not row:
        return

    title = str(row["title"])
    if not title.startswith("新建"):
        return

    new_title = first_user_message.strip().replace("\n", " ")[:32]
    if not new_title:
        return

    app_db.execute(
        "UPDATE chat_sessions SET title = ?, updated_at = ? WHERE id = ?",
        (new_title, current_timestamp(), session_id),
    )
    app_db.commit()


def delete_all_chat_sessions() -> None:
    app_db.execute("DELETE FROM chat_messages")
    app_db.execute("DELETE FROM chat_sessions")
    app_db.commit()


def delete_all_ingested_file_records() -> None:
    app_db.execute("DELETE FROM ingested_files")
    app_db.commit()


def create_ingest_task(total_files: int) -> str:
    task_id = uuid.uuid4().hex
    now = current_timestamp()
    app_db.execute(
        """
        INSERT INTO ingest_tasks (
            id, status, total_files, processed_files, success_count, duplicate_count,
            skipped_count, failed_count, current_file, message, created_at, updated_at
        )
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        """,
        (task_id, "running", total_files, 0, 0, 0, 0, 0, "", "准备处理文件", now, now),
    )
    app_db.commit()
    return task_id


def update_ingest_task(task_id: str, **kwargs: Any) -> None:
    if not kwargs:
        return
    kwargs["updated_at"] = current_timestamp()
    assignments = ", ".join(f"{key} = ?" for key in kwargs)
    values = list(kwargs.values())
    values.append(task_id)
    app_db.execute(
        f"UPDATE ingest_tasks SET {assignments} WHERE id = ?",
        values,
    )
    app_db.commit()


def record_ingest_task_item(
    task_id: str,
    file_name: str,
    status: str,
    message: str,
    chunk_count: int = 0,
    file_sha256: str = "",
) -> None:
    now = current_timestamp()
    app_db.execute(
        """
        INSERT INTO ingest_task_items (
            id, task_id, file_name, status, message, chunk_count, sha256, created_at, updated_at
        )
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
        """,
        (
            uuid.uuid4().hex,
            task_id,
            file_name,
            status,
            message,
            int(chunk_count or 0),
            file_sha256 or "",
            now,
            now,
        ),
    )
    app_db.commit()


def list_ingest_tasks(limit: int = 10) -> List[Dict[str, Any]]:
    rows = app_db.execute(
        """
        SELECT id, status, total_files, processed_files, success_count, duplicate_count,
               skipped_count, failed_count, current_file, message, created_at, updated_at
        FROM ingest_tasks
        ORDER BY updated_at DESC
        LIMIT ?
        """,
        (limit,),
    ).fetchall()
    return [dict(row) for row in rows]


def get_ingest_task(task_id: str) -> Optional[Dict[str, Any]]:
    row = app_db.execute(
        """
        SELECT id, status, total_files, processed_files, success_count, duplicate_count,
               skipped_count, failed_count, current_file, message, created_at, updated_at
        FROM ingest_tasks
        WHERE id = ?
        """,
        (task_id,),
    ).fetchone()
    return dict(row) if row else None


def request_pause_ingest_task(task_id: str) -> None:
    update_ingest_task(task_id, status="pause_requested", message="已请求暂停，当前文件处理到安全点后暂停")


def resume_ingest_task(task_id: str) -> None:
    update_ingest_task(task_id, status="running", message="已继续后台入库任务")


def request_cancel_ingest_task(task_id: str) -> None:
    update_ingest_task(task_id, status="cancel_requested", message="已请求终止，当前步骤结束后停止")


class IngestTaskCancelled(RuntimeError):
    pass


def wait_if_task_paused_or_cancelled(task_id: str) -> None:
    while True:
        task = get_ingest_task(task_id)
        status = task.get("status") if task else ""
        if status == "cancel_requested":
            update_ingest_task(task_id, status="cancelled", message="入库任务已终止")
            raise IngestTaskCancelled("入库任务已终止")
        if status in {"pause_requested", "paused"}:
            if status == "pause_requested":
                update_ingest_task(task_id, status="paused", message="入库任务已暂停")
            time.sleep(1)
            continue
        return


def list_ingest_task_items(task_id: str, limit: int = 200) -> List[Dict[str, Any]]:
    rows = app_db.execute(
        """
        SELECT file_name, status, message, chunk_count, sha256, created_at, updated_at
        FROM ingest_task_items
        WHERE task_id = ?
        ORDER BY updated_at DESC
        LIMIT ?
        """,
        (task_id, limit),
    ).fetchall()
    return [dict(row) for row in rows]


def reset_app_state_database() -> None:
    delete_all_config_values()
    delete_all_chat_sessions()
    set_config_value("ui_language", DEFAULT_UI_LANGUAGE)
    save_llm_config(
        DEFAULT_LLM_BASE_URL,
        DEFAULT_LLM_API_KEY,
        DEFAULT_LLM_MODEL,
        fast_model=DEFAULT_LLM_MODEL,
        thinking_model=DEFAULT_LLM_MODEL,
        fast_extra_body=DEFAULT_LLM_EXTRA_BODY,
        thinking_extra_body=DEFAULT_LLM_EXTRA_BODY,
    )
    set_config_value("rag_context_turns", DEFAULT_CONTEXT_TURNS)
    set_config_value("compliance_context_turns", DEFAULT_CONTEXT_TURNS)
    set_config_value("rag_top_k", DEFAULT_RAG_TOP_K)
    set_bool_config("rag_query_rewrite", True)
    set_bool_config("rag_query_decompose", DEFAULT_QUERY_DECOMPOSE)
    set_bool_config("rag_use_distance_threshold", True)
    set_config_value("rag_max_distance", DEFAULT_VECTOR_MAX_DISTANCE)
    set_bool_config("rag_use_hybrid", DEFAULT_USE_HYBRID_SEARCH)
    set_bool_config("rag_use_reranker", DEFAULT_USE_RERANKER)
    set_config_value("rag_fetch_k", DEFAULT_RETRIEVAL_FETCH_K)
    set_config_value("compliance_regulation_top_k", DEFAULT_COMPLIANCE_REGULATION_TOP_K)
    set_config_value("compliance_enterprise_top_k", DEFAULT_COMPLIANCE_ENTERPRISE_TOP_K)
    set_bool_config("compliance_query_rewrite", True)
    set_bool_config("compliance_query_decompose", DEFAULT_QUERY_DECOMPOSE)
    set_config_value("compliance_min_regulation_evidence", DEFAULT_COMPLIANCE_MIN_REGULATION_EVIDENCE)
    set_config_value("compliance_min_enterprise_evidence", DEFAULT_COMPLIANCE_MIN_ENTERPRISE_EVIDENCE)
    set_bool_config("compliance_clause_by_clause", False)
    set_bool_config("compliance_include_missing_list", True)
    set_bool_config("compliance_use_distance_threshold", True)
    set_config_value("compliance_max_distance", DEFAULT_VECTOR_MAX_DISTANCE)
    set_bool_config("compliance_use_hybrid", DEFAULT_USE_HYBRID_SEARCH)
    set_bool_config("compliance_use_reranker", DEFAULT_USE_RERANKER)
    set_config_value("compliance_fetch_k", DEFAULT_RETRIEVAL_FETCH_K)
    set_config_value("model_cache_root", DEFAULT_MODEL_CACHE_ROOT)
    set_config_value("paddleocr_cache_dir", "")
    set_config_value("bge_cache_dir", "")
    set_config_value("reranker_cache_dir", "")
    set_config_value("soffice_binary_path", DEFAULT_SOFFICE_BINARY_PATH)
    set_bool_config("replace_changed_same_name", True)
    set_bool_config("background_ingest", DEFAULT_BACKGROUND_INGEST)
    apply_model_cache_environment()
    for key in [
        "llm_config",
        "active_rag_session_id",
        "active_compliance_session_id",
        "upload_mode_label",
        "upload_doc_category_label",
        "upload_ocr_enhance",
        "auto_unload_models_after_ingest",
        "replace_changed_same_name_input",
        "background_ingest_input",
        "upload_pdf_ocr_mode_label",
        "paddleocr_model_label",
        "auto_install_libreoffice",
        "upload_chunk_size",
        "upload_overlap",
        "rag_session_select",
        "compliance_session_select",
        "rag_search_scope_label",
        "search_top_k",
        "chat_llm_mode_label",
        "rag_allow_general_fallback_input",
        "rag_query_rewrite_input",
        "rag_query_decompose_input",
        "rag_use_distance_threshold_input",
        "rag_max_distance_input",
        "rag_use_hybrid_input",
        "rag_use_reranker_input",
        "rag_fetch_k_input",
        "rag_context_turns_input",
        "regulation_top_k",
        "enterprise_top_k",
        "compliance_query_rewrite_input",
        "compliance_query_decompose_input",
        "compliance_min_regulation_evidence_input",
        "compliance_min_enterprise_evidence_input",
        "compliance_clause_by_clause_input",
        "compliance_include_missing_list_input",
        "compliance_use_distance_threshold_input",
        "compliance_max_distance_input",
        "compliance_use_hybrid_input",
        "compliance_use_reranker_input",
        "compliance_fetch_k_input",
        "compliance_llm_mode_label",
        "compliance_context_turns_input",
        "config_test_llm_mode_label",
        "config_paddleocr_model_label",
        "model_cache_root_input",
        "paddleocr_cache_dir_input",
        "bge_cache_dir_input",
        "reranker_cache_dir_input",
        "soffice_binary_path_input",
        "ui_language_selector",
        "model_status_test_llm_mode_label",
    ]:
        st.session_state.pop(key, None)


def get_model_cache_root() -> str:
    return normalize_local_path(get_config_value("model_cache_root", DEFAULT_MODEL_CACHE_ROOT), DEFAULT_MODEL_CACHE_ROOT)


def get_paddleocr_cache_dir() -> str:
    default_path = os.path.join(get_model_cache_root(), "paddlex")
    return normalize_local_path(get_config_value("paddleocr_cache_dir", ""), default_path)


def get_bge_cache_dir() -> str:
    default_path = os.path.join(get_model_cache_root(), "bge-m3")
    return normalize_local_path(get_config_value("bge_cache_dir", ""), default_path)


def get_reranker_cache_dir() -> str:
    default_path = os.path.join(get_model_cache_root(), "bge-reranker-v2-m3")
    return normalize_local_path(get_config_value("reranker_cache_dir", ""), default_path)


def get_configured_soffice_path() -> str:
    return normalize_local_path(get_config_value("soffice_binary_path", DEFAULT_SOFFICE_BINARY_PATH), "")


def ensure_model_cache_dirs() -> None:
    for path in [
        get_model_cache_root(),
        get_paddleocr_cache_dir(),
        get_bge_cache_dir(),
        get_reranker_cache_dir(),
    ]:
        if path:
            os.makedirs(path, exist_ok=True)


def refresh_paddlex_cache_runtime() -> None:
    cache_module = sys.modules.get("paddlex.utils.cache")
    if not cache_module:
        return

    cache_dir = get_paddleocr_cache_dir()
    cache_module.CACHE_DIR = cache_dir
    cache_module.FUNC_CACHE_DIR = os.path.join(cache_dir, "func_ret")
    cache_module.FILE_LOCK_DIR = os.path.join(cache_dir, "locks")
    cache_module.TEMP_DIR = os.path.join(cache_dir, "temp")
    for path in [
        cache_module.CACHE_DIR,
        cache_module.FUNC_CACHE_DIR,
        cache_module.FILE_LOCK_DIR,
        cache_module.TEMP_DIR,
    ]:
        os.makedirs(path, exist_ok=True)


def apply_model_cache_environment() -> None:
    ensure_model_cache_dirs()
    os.environ["PADDLE_PDX_CACHE_HOME"] = get_paddleocr_cache_dir()
    os.environ["HF_HOME"] = os.path.join(get_model_cache_root(), "huggingface")
    os.environ["SENTENCE_TRANSFORMERS_HOME"] = get_model_cache_root()
    refresh_paddlex_cache_runtime()


def save_model_cache_config(
    model_cache_root: str,
    paddleocr_cache_dir: str,
    bge_cache_dir: str,
    reranker_cache_dir: str,
    soffice_binary_path: str,
) -> None:
    model_cache_root = normalize_local_path(model_cache_root, DEFAULT_MODEL_CACHE_ROOT)
    paddleocr_cache_dir = normalize_local_path(paddleocr_cache_dir, "") if paddleocr_cache_dir.strip() else ""
    bge_cache_dir = normalize_local_path(bge_cache_dir, "") if bge_cache_dir.strip() else ""
    reranker_cache_dir = normalize_local_path(reranker_cache_dir, "") if reranker_cache_dir.strip() else ""
    soffice_binary_path = normalize_local_path(soffice_binary_path, "")

    set_config_value("model_cache_root", model_cache_root)
    set_config_value("paddleocr_cache_dir", paddleocr_cache_dir)
    set_config_value("bge_cache_dir", bge_cache_dir)
    set_config_value("reranker_cache_dir", reranker_cache_dir)
    set_config_value("soffice_binary_path", soffice_binary_path)
    apply_model_cache_environment()
    try:
        load_ocr_model.clear()
        load_embedding_model.clear()
        load_reranker_model.clear()
    except Exception:
        pass
    release_memory_after_file()


apply_model_cache_environment()


def get_paddleocr_model_label() -> str:
    label = get_config_value("paddleocr_model_label", DEFAULT_PADDLEOCR_MODEL_LABEL)
    if label not in PADDLEOCR_MODEL_OPTIONS:
        return DEFAULT_PADDLEOCR_MODEL_LABEL
    return label


def get_paddleocr_model_config() -> Dict[str, str]:
    return PADDLEOCR_MODEL_OPTIONS[get_paddleocr_model_label()]


def save_paddleocr_model_label(label: str) -> None:
    if label not in PADDLEOCR_MODEL_OPTIONS:
        raise ValueError("未知 PaddleOCR 模型配置")
    old_label = get_paddleocr_model_label()
    set_config_value("paddleocr_model_label", label)
    if old_label != label:
        try:
            load_ocr_model.clear()
        except Exception:
            pass
        release_memory_after_file()


def is_bge_model_cached() -> bool:
    try:
        if os.path.exists(os.path.join(get_bge_cache_dir(), "modules.json")):
            return True
        from huggingface_hub import try_to_load_from_cache

        cached_path = try_to_load_from_cache(
            EMBEDDING_MODEL_NAME,
            "modules.json",
            cache_dir=get_bge_cache_dir(),
        )
        return isinstance(cached_path, str) and os.path.exists(cached_path)
    except Exception:
        return False


def is_reranker_model_cached() -> bool:
    try:
        if os.path.exists(os.path.join(get_reranker_cache_dir(), "config.json")):
            return True
        from huggingface_hub import try_to_load_from_cache

        cached_path = try_to_load_from_cache(
            RERANKER_MODEL_NAME,
            "config.json",
            cache_dir=get_reranker_cache_dir(),
        )
        return isinstance(cached_path, str) and os.path.exists(cached_path)
    except Exception:
        return False


# =========================
# 加载模型和数据库
# =========================
@st.cache_resource
def load_ocr_model():
    """
    PaddleOCR 中文模型。
    use_textline_orientation=True 用于处理文字方向。
    lang='ch' 适合中文，也能识别一部分英文。
    """
    os.environ["PADDLE_PDX_CACHE_HOME"] = get_paddleocr_cache_dir()
    from paddleocr import PaddleOCR

    ocr_config = get_paddleocr_model_config()
    return PaddleOCR(
        use_doc_orientation_classify=False,
        use_doc_unwarping=False,
        use_textline_orientation=True,
        text_detection_model_name=ocr_config["det"],
        text_recognition_model_name=ocr_config["rec"],
        text_recognition_batch_size=4,
        text_det_limit_side_len=960,
        lang="ch",
    )


@st.cache_resource
def load_embedding_model():
    """
    BAAI/bge-m3 embedding 模型。
    第一次运行会下载模型，速度取决于网络。
    """
    from sentence_transformers import SentenceTransformer

    model_source = get_bge_cache_dir() if os.path.exists(os.path.join(get_bge_cache_dir(), "modules.json")) else EMBEDDING_MODEL_NAME
    return SentenceTransformer(
        model_source,
        cache_folder=get_bge_cache_dir(),
        local_files_only=is_bge_model_cached(),
    )


@st.cache_resource
def load_reranker_model():
    from sentence_transformers import CrossEncoder

    model_source = get_reranker_cache_dir() if os.path.exists(os.path.join(get_reranker_cache_dir(), "config.json")) else RERANKER_MODEL_NAME
    return CrossEncoder(
        model_source,
        max_length=512,
        cache_folder=get_reranker_cache_dir(),
        local_files_only=is_reranker_model_cached(),
    )


def import_qdrant_models():
    from qdrant_client import models

    return models


def build_qdrant_filter(where: Optional[Dict[str, Any]] = None):
    if not where:
        return None
    models = import_qdrant_models()
    return models.Filter(
        must=[
            models.FieldCondition(
                key=key,
                match=models.MatchValue(value=value),
            )
            for key, value in where.items()
        ]
    )


def ensure_qdrant_collection(client) -> None:
    models = import_qdrant_models()
    collection_exists = False
    try:
        collection_exists = client.collection_exists(COLLECTION_NAME)
    except Exception:
        try:
            client.get_collection(COLLECTION_NAME)
            collection_exists = True
        except Exception:
            collection_exists = False

    if not collection_exists:
        client.create_collection(
            collection_name=COLLECTION_NAME,
            vectors_config=models.VectorParams(size=VECTOR_SIZE, distance=models.Distance.COSINE),
        )


def recreate_qdrant_collection() -> None:
    try:
        vector_client.delete_collection(COLLECTION_NAME)
    except Exception:
        pass
    ensure_qdrant_collection(vector_client)


@st.cache_resource
def load_qdrant_client():
    from qdrant_client import QdrantClient

    client = QdrantClient(path=QDRANT_DIR)
    ensure_qdrant_collection(client)
    return client


@st.cache_resource
def load_llm_client(base_url: str, api_key: str):
    return OpenAI(base_url=base_url, api_key=api_key, timeout=60)


vector_client = load_qdrant_client()


@st.cache_resource
def load_ingest_executor():
    return {
        "executor": ThreadPoolExecutor(max_workers=1, thread_name_prefix="ocr-rag-ingest"),
        "futures": {},
    }


def qdrant_scroll_points(
    where: Optional[Dict[str, Any]] = None,
    limit: int = 10000,
    with_payload: bool = True,
) -> List[Any]:
    points = []
    next_page = None
    scroll_filter = build_qdrant_filter(where)
    while True:
        batch, next_page = vector_client.scroll(
            collection_name=COLLECTION_NAME,
            scroll_filter=scroll_filter,
            limit=min(limit, 256),
            offset=next_page,
            with_payload=with_payload,
            with_vectors=False,
        )
        points.extend(batch)
        if next_page is None or len(points) >= limit:
            break
    return points[:limit]


def point_payload(point: Any) -> Dict[str, Any]:
    return dict(getattr(point, "payload", None) or {})


def payload_to_result(point: Any, score: Optional[float] = None, retrieval_source: str = "vector") -> Dict[str, Any]:
    payload = point_payload(point)
    content = payload.pop("document", "")
    distance = None
    if score is not None:
        distance = max(0.0, 1.0 - float(score))
    return {
        "id": str(getattr(point, "id", uuid.uuid4())),
        "content": content,
        "metadata": payload,
        "distance": distance,
        "retrieval_source": retrieval_source,
    }


def delete_vector_chunks_by_where(where: Dict[str, Any]) -> int:
    points = qdrant_scroll_points(where=where, limit=100000)
    point_ids = [point.id for point in points]
    if point_ids:
        models = import_qdrant_models()
        vector_client.delete(
            collection_name=COLLECTION_NAME,
            points_selector=models.PointIdsList(points=point_ids),
        )
    return len(point_ids)


def delete_vector_chunks_by_file_name(file_name: str) -> int:
    if not file_name:
        return 0
    return delete_vector_chunks_by_where({"file_name": file_name})


def count_chunks_by_file_sha256(file_sha256: str) -> int:
    if not file_sha256:
        return 0
    try:
        result = vector_client.count(
            collection_name=COLLECTION_NAME,
            count_filter=build_qdrant_filter({"file_sha256": file_sha256}),
            exact=True,
        )
        return int(result.count)
    except Exception:
        return 0


def get_duplicate_ingested_file(file_sha256: str) -> Optional[Dict[str, Any]]:
    existing_file = get_ingested_file(file_sha256)
    if existing_file and count_chunks_by_file_sha256(file_sha256) > 0:
        return existing_file
    return None


def replace_existing_same_name_if_needed(file_name: str, file_sha256: str, enabled: bool) -> Tuple[int, int]:
    if not enabled:
        return 0, 0
    same_name_records = [
        record for record in list_ingested_files_by_name(file_name) if record.get("sha256") != file_sha256
    ]
    existing_points = qdrant_scroll_points(where={"file_name": file_name}, limit=100000)
    existing_point_ids = [point.id for point in existing_points]
    if not same_name_records and not existing_point_ids:
        return 0, 0

    if existing_point_ids:
        models = import_qdrant_models()
        vector_client.delete(
            collection_name=COLLECTION_NAME,
            points_selector=models.PointIdsList(points=existing_point_ids),
        )
    deleted_chunks = len(existing_point_ids)
    deleted_records = delete_ingested_file_records_by_name(file_name)
    return deleted_chunks, deleted_records


def load_llm_config_from_env() -> Dict[str, str]:
    env_values = dotenv_values(ENV_FILE) if os.path.exists(ENV_FILE) else {}
    model = (
        env_values.get("LLM_MODEL")
        or os.getenv("LLM_MODEL")
        or DEFAULT_LLM_MODEL
    )
    return {
        "base_url": (
            env_values.get("LLM_BASE_URL")
            or os.getenv("LLM_BASE_URL")
            or DEFAULT_LLM_BASE_URL
        ),
        "api_key": (
            env_values.get("LLM_API_KEY")
            or os.getenv("LLM_API_KEY")
            or DEFAULT_LLM_API_KEY
        ),
        "model": model,
        "fast_model": (
            env_values.get("LLM_FAST_MODEL")
            or os.getenv("LLM_FAST_MODEL")
            or model
        ),
        "thinking_model": (
            env_values.get("LLM_THINKING_MODEL")
            or os.getenv("LLM_THINKING_MODEL")
            or model
        ),
        "fast_extra_body": (
            env_values.get("LLM_FAST_EXTRA_BODY")
            or os.getenv("LLM_FAST_EXTRA_BODY")
            or DEFAULT_LLM_EXTRA_BODY
        ),
        "thinking_extra_body": (
            env_values.get("LLM_THINKING_EXTRA_BODY")
            or os.getenv("LLM_THINKING_EXTRA_BODY")
            or DEFAULT_LLM_EXTRA_BODY
        ),
    }


def load_llm_config_from_db() -> Optional[Dict[str, str]]:
    base_url = get_config_value("LLM_BASE_URL", "")
    if not base_url:
        return None

    model = get_config_value("LLM_MODEL", DEFAULT_LLM_MODEL)
    return {
        "base_url": base_url,
        "api_key": get_config_value("LLM_API_KEY", DEFAULT_LLM_API_KEY),
        "model": model,
        "fast_model": get_config_value("LLM_FAST_MODEL", model),
        "thinking_model": get_config_value("LLM_THINKING_MODEL", model),
        "fast_extra_body": get_config_value("LLM_FAST_EXTRA_BODY", DEFAULT_LLM_EXTRA_BODY),
        "thinking_extra_body": get_config_value("LLM_THINKING_EXTRA_BODY", DEFAULT_LLM_EXTRA_BODY),
    }


def persist_llm_config(config: Dict[str, str]) -> None:
    set_config_value("LLM_BASE_URL", config["base_url"])
    set_config_value("LLM_API_KEY", config["api_key"])
    set_config_value("LLM_MODEL", config["model"])
    set_config_value("LLM_FAST_MODEL", config["fast_model"])
    set_config_value("LLM_THINKING_MODEL", config["thinking_model"])
    set_config_value("LLM_FAST_EXTRA_BODY", config["fast_extra_body"])
    set_config_value("LLM_THINKING_EXTRA_BODY", config["thinking_extra_body"])


def get_llm_config() -> Dict[str, str]:
    if "llm_config" not in st.session_state:
        config = load_llm_config_from_db()
        if config is None:
            config = load_llm_config_from_env()
            persist_llm_config(config)
        st.session_state["llm_config"] = config
    return st.session_state["llm_config"]


def save_llm_config(
    base_url: str,
    api_key: str,
    model: str,
    fast_model: str = "",
    thinking_model: str = "",
    fast_extra_body: str = DEFAULT_LLM_EXTRA_BODY,
    thinking_extra_body: str = DEFAULT_LLM_EXTRA_BODY,
) -> None:
    base_url = base_url.strip() or DEFAULT_LLM_BASE_URL
    api_key = api_key.strip() or DEFAULT_LLM_API_KEY
    model = model.strip() or DEFAULT_LLM_MODEL
    fast_model = fast_model.strip() or model
    thinking_model = thinking_model.strip() or model
    fast_extra_body = fast_extra_body.strip() or DEFAULT_LLM_EXTRA_BODY
    thinking_extra_body = thinking_extra_body.strip() or DEFAULT_LLM_EXTRA_BODY

    parse_extra_body(fast_extra_body)
    parse_extra_body(thinking_extra_body)

    config = {
        "base_url": base_url,
        "api_key": api_key,
        "model": model,
        "fast_model": fast_model,
        "thinking_model": thinking_model,
        "fast_extra_body": fast_extra_body,
        "thinking_extra_body": thinking_extra_body,
    }
    persist_llm_config(config)

    os.environ["LLM_BASE_URL"] = base_url
    os.environ["LLM_API_KEY"] = api_key
    os.environ["LLM_MODEL"] = model
    os.environ["LLM_FAST_MODEL"] = fast_model
    os.environ["LLM_THINKING_MODEL"] = thinking_model
    os.environ["LLM_FAST_EXTRA_BODY"] = fast_extra_body
    os.environ["LLM_THINKING_EXTRA_BODY"] = thinking_extra_body
    st.session_state["llm_config"] = config


def get_llm_client() -> OpenAI:
    config = get_llm_config()
    return load_llm_client(config["base_url"], config["api_key"])


def parse_extra_body(raw_extra_body: str) -> Dict[str, Any]:
    raw_extra_body = (raw_extra_body or "").strip()
    if not raw_extra_body:
        return {}
    parsed = json.loads(raw_extra_body)
    if not isinstance(parsed, dict):
        raise ValueError("extra_body 必须是 JSON 对象，例如 {} 或 {\"enable_thinking\": true}")
    return parsed


def get_llm_mode_config(mode: str) -> Tuple[str, Dict[str, Any]]:
    config = get_llm_config()
    if mode == "thinking":
        model = config.get("thinking_model") or config["model"]
        extra_body = parse_extra_body(config.get("thinking_extra_body", DEFAULT_LLM_EXTRA_BODY))
    else:
        model = config.get("fast_model") or config["model"]
        extra_body = parse_extra_body(config.get("fast_extra_body", DEFAULT_LLM_EXTRA_BODY))
    return model, extra_body


def create_llm_chat_completion(
    messages: List[Dict[str, str]],
    temperature: float,
    mode: str = "fast",
):
    model, extra_body = get_llm_mode_config(mode)
    kwargs = {
        "model": model,
        "messages": messages,
        "temperature": temperature,
    }
    if extra_body:
        kwargs["extra_body"] = extra_body
    return get_llm_client().chat.completions.create(**kwargs)


# =========================
# 模型状态
# =========================
def record_model_event(component: str, status: str, detail: str) -> None:
    st.session_state["model_events"].append(
        {
            "time": time.strftime("%H:%M:%S"),
            "component": component,
            "status": status,
            "detail": detail,
        }
    )


def get_bge_cache_status() -> str:
    try:
        if is_bge_model_cached():
            return f"已发现本地缓存：{get_bge_cache_dir()}"
        return f"未发现完整缓存，首次使用会下载到：{get_bge_cache_dir()}"
    except Exception as e:
        return f"无法检查缓存：{e}"


def get_reranker_cache_status() -> str:
    try:
        if is_reranker_model_cached():
            return f"已发现本地缓存：{get_reranker_cache_dir()}"
        return f"未发现完整缓存，首次使用会下载到：{get_reranker_cache_dir()}"
    except Exception as e:
        return f"无法检查缓存：{e}"


def get_paddle_cache_status() -> str:
    cache_roots = [
        Path(get_paddleocr_cache_dir()) / "official_models",
        Path.home() / ".paddlex" / "official_models",
        Path.home() / ".paddleocr",
        Path.home() / ".cache" / "paddle",
    ]
    existing = []
    for root in cache_roots:
        if root.exists():
            children = [child.name for child in root.iterdir()]
            if children:
                existing.append(f"{root} ({len(children)} 项)")

    if existing:
        return "已发现缓存：" + "；".join(existing[:2])
    return "未发现明显缓存，首次 OCR 会下载"


def test_llm_connection(mode: str = "fast") -> str:
    response = create_llm_chat_completion(
        messages=[
            {"role": "system", "content": "你是连通性测试助手，只回复 OK。"},
            {"role": "user", "content": "测试"},
        ],
        temperature=0,
        mode=mode,
    )
    return response.choices[0].message.content or ""


# =========================
# 文件处理函数
# =========================
def sanitize_uploaded_relative_name(file_name: str) -> str:
    raw_name = (file_name or "uploaded_file").replace("\\", "/")
    safe_parts = []
    invalid_chars = '<>:"|?*'
    translation = str.maketrans({char: "_" for char in invalid_chars})
    for part in raw_name.split("/"):
        cleaned = part.strip().translate(translation).rstrip(" .")
        if cleaned and cleaned not in {".", ".."}:
            safe_parts.append(cleaned)
    if not safe_parts:
        safe_parts = ["uploaded_file"]
    return os.path.join(*safe_parts)


def get_uploaded_relative_name(uploaded_file) -> str:
    return sanitize_uploaded_relative_name(getattr(uploaded_file, "name", "uploaded_file"))


def calculate_uploaded_file_sha256(uploaded_file) -> str:
    digest = hashlib.sha256()
    digest.update(uploaded_file.getbuffer())
    return digest.hexdigest()


def get_file_extension_from_name(file_name: str) -> str:
    return Path(file_name).suffix.lower().lstrip(".")


def is_legacy_office_file(file_name: str) -> bool:
    return get_file_extension_from_name(file_name) in LEGACY_OFFICE_TYPES


def is_supported_upload(file_name: str) -> bool:
    return get_file_extension_from_name(file_name) in SUPPORTED_TYPES


def is_soffice_executable(candidate: str) -> bool:
    return bool(candidate and os.path.exists(candidate) and not os.path.isdir(candidate))


def expand_soffice_candidates(path_value: str) -> List[str]:
    if not path_value:
        return []
    path_value = normalize_local_path(path_value, "")
    if os.path.isdir(path_value):
        return [
            os.path.join(path_value, "soffice"),
            os.path.join(path_value, "libreoffice"),
            os.path.join(path_value, "program", "soffice.exe"),
            os.path.join(path_value, "Contents", "MacOS", "soffice"),
        ]
    return [path_value]


def find_soffice_binary() -> Optional[str]:
    program_files = os.environ.get("PROGRAMFILES", r"C:\Program Files")
    program_files_x86 = os.environ.get("PROGRAMFILES(X86)", r"C:\Program Files (x86)")
    configured_candidates = expand_soffice_candidates(get_configured_soffice_path())
    for candidate in [
        *configured_candidates,
        shutil.which("soffice"),
        shutil.which("libreoffice"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/usr/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/libreoffice",
        "/usr/local/bin/libreoffice",
        "/snap/bin/libreoffice",
        os.path.join(program_files, "LibreOffice", "program", "soffice.exe"),
        os.path.join(program_files_x86, "LibreOffice", "program", "soffice.exe"),
    ]:
        if candidate and is_soffice_executable(candidate):
            return candidate
    return None


def find_textutil_binary() -> Optional[str]:
    return shutil.which("textutil")


def quote_command(command: List[str]) -> str:
    if platform.system() == "Windows":
        return subprocess.list2cmdline(command)
    return " ".join(shlex.quote(part) for part in command)


def linux_privileged_command(command: List[str]) -> List[str]:
    if platform.system() != "Linux":
        return command
    geteuid = getattr(os, "geteuid", None)
    if callable(geteuid) and geteuid() == 0:
        return command
    return ["sudo", "-n", *command]


def get_libreoffice_install_plan() -> Dict[str, Any]:
    system_name = platform.system()

    if system_name == "Darwin":
        brew = shutil.which("brew")
        if brew:
            command = [brew, "install", "--cask", "libreoffice"]
            return {
                "platform": "macOS",
                "commands": [command],
                "manual": quote_command(command),
            }
        return {
            "platform": "macOS",
            "commands": [],
            "manual": "先安装 Homebrew，然后执行：brew install --cask libreoffice",
        }

    if system_name == "Windows":
        winget = shutil.which("winget")
        if winget:
            command = [
                winget,
                "install",
                "--id",
                "TheDocumentFoundation.LibreOffice",
                "-e",
                "--accept-source-agreements",
                "--accept-package-agreements",
            ]
            return {
                "platform": "Windows",
                "commands": [command],
                "manual": quote_command(command),
            }

        choco = shutil.which("choco")
        if choco:
            command = [choco, "install", "libreoffice-fresh", "-y"]
            return {
                "platform": "Windows",
                "commands": [command],
                "manual": quote_command(command),
            }

        return {
            "platform": "Windows",
            "commands": [],
            "manual": "安装 winget 或 Chocolatey 后执行：winget install --id TheDocumentFoundation.LibreOffice -e",
        }

    if system_name == "Linux":
        if shutil.which("apt-get"):
            commands = [
                linux_privileged_command(["apt-get", "update"]),
                linux_privileged_command(["apt-get", "install", "-y", "libreoffice"]),
            ]
            return {
                "platform": "Linux",
                "commands": commands,
                "manual": "sudo apt-get update && sudo apt-get install -y libreoffice",
            }

        if shutil.which("dnf"):
            command = linux_privileged_command(["dnf", "install", "-y", "libreoffice"])
            return {
                "platform": "Linux",
                "commands": [command],
                "manual": "sudo dnf install -y libreoffice",
            }

        if shutil.which("yum"):
            command = linux_privileged_command(["yum", "install", "-y", "libreoffice"])
            return {
                "platform": "Linux",
                "commands": [command],
                "manual": "sudo yum install -y libreoffice",
            }

        if shutil.which("zypper"):
            command = linux_privileged_command(["zypper", "--non-interactive", "install", "libreoffice"])
            return {
                "platform": "Linux",
                "commands": [command],
                "manual": "sudo zypper --non-interactive install libreoffice",
            }

        if shutil.which("pacman"):
            command = linux_privileged_command(["pacman", "-S", "--noconfirm", "libreoffice-still"])
            return {
                "platform": "Linux",
                "commands": [command],
                "manual": "sudo pacman -S --noconfirm libreoffice-still",
            }

        if shutil.which("snap"):
            command = linux_privileged_command(["snap", "install", "libreoffice"])
            return {
                "platform": "Linux",
                "commands": [command],
                "manual": "sudo snap install libreoffice",
            }

        return {
            "platform": "Linux",
            "commands": [],
            "manual": "请使用当前发行版包管理器安装 LibreOffice，例如：sudo apt-get install -y libreoffice",
        }

    return {
        "platform": system_name or "未知系统",
        "commands": [],
        "manual": "当前系统未内置自动安装方案，请手动安装 LibreOffice。",
    }


def install_libreoffice_automatically() -> Tuple[bool, str]:
    existing_binary = find_soffice_binary()
    if existing_binary:
        return True, f"已检测到 LibreOffice：{existing_binary}"

    plan = get_libreoffice_install_plan()
    commands = plan.get("commands", [])
    manual_command = plan.get("manual", "")
    if not commands:
        return False, f"未找到可用的自动安装方式。请手动执行：{manual_command}"

    logs = []
    for command in commands:
        result = run_subprocess(command, timeout=1800)
        command_text = quote_command(command)
        if result.returncode != 0:
            output = (result.stderr or result.stdout or "").strip()
            return (
                False,
                f"自动安装失败。\n命令：{command_text}\n原因：{output or '未知错误'}\n可手动执行：{manual_command}",
            )
        logs.append(f"完成：{command_text}")

    installed_binary = find_soffice_binary()
    if installed_binary:
        return True, f"LibreOffice 安装完成：{installed_binary}"

    return (
        False,
        "安装命令已执行，但当前进程还没有检测到 soffice。"
        "请重启应用后再试；如仍失败，可手动执行："
        f"{manual_command}",
    )


def legacy_conversion_target_ext(ext: str) -> str:
    return {
        "doc": "docx",
        "ppt": "pptx",
        "xls": "xlsx",
    }[ext]


def get_legacy_conversion_status(ext: str) -> Tuple[bool, str]:
    ext = ext.lower().lstrip(".")
    if ext not in LEGACY_OFFICE_TYPES:
        return True, "无需转换"

    if find_soffice_binary():
        return True, "将使用 LibreOffice 转换"

    if ext == "doc" and find_textutil_binary():
        return True, "将使用 macOS textutil 转换 DOC"

    return False, "需要安装 LibreOffice 才能转换 .doc/.ppt/.xls 老格式"


def is_processable_upload(file_name: str) -> bool:
    if not is_supported_upload(file_name):
        return False
    if not is_legacy_office_file(file_name):
        return True

    ext = get_file_extension_from_name(file_name)
    can_convert, _message = get_legacy_conversion_status(ext)
    return can_convert


def save_uploaded_file(uploaded_file, batch_id: Optional[str] = None) -> str:
    relative_name = get_uploaded_relative_name(uploaded_file)
    if batch_id:
        file_path = os.path.join(UPLOAD_DIR, batch_id, relative_name)
    else:
        safe_name = Path(relative_name).name
        file_path = os.path.join(UPLOAD_DIR, f"{uuid.uuid4().hex}_{safe_name}")
    os.makedirs(os.path.dirname(file_path), exist_ok=True)
    with open(file_path, "wb") as f:
        f.write(uploaded_file.getbuffer())
    return file_path


def run_subprocess(command: List[str], timeout: int = 180) -> subprocess.CompletedProcess:
    return subprocess.run(
        command,
        check=False,
        capture_output=True,
        text=True,
        timeout=timeout,
    )


def convert_legacy_office_file(file_path: str) -> str:
    ext = get_file_extension(file_path)
    if ext not in LEGACY_OFFICE_TYPES:
        return file_path

    target_ext = legacy_conversion_target_ext(ext)
    output_dir = os.path.join(CONVERTED_DIR, uuid.uuid4().hex)
    os.makedirs(output_dir, exist_ok=True)
    converted_path = os.path.join(output_dir, f"{Path(file_path).stem}.{target_ext}")

    soffice_binary = find_soffice_binary()
    if soffice_binary:
        result = run_subprocess(
            [
                soffice_binary,
                "--headless",
                "--convert-to",
                target_ext,
                "--outdir",
                output_dir,
                file_path,
            ],
            timeout=240,
        )
        if result.returncode != 0:
            detail = (result.stderr or result.stdout or "").strip()
            raise RuntimeError(f"LibreOffice 转换失败：{detail or '未知错误'}")

        if os.path.exists(converted_path):
            return converted_path

        converted_candidates = list(Path(output_dir).glob(f"*.{target_ext}"))
        if converted_candidates:
            return str(converted_candidates[0])

        detail = (result.stderr or result.stdout or "").strip()
        raise RuntimeError(f"LibreOffice 未生成 {target_ext} 文件：{detail or '无输出'}")

    if ext == "doc":
        textutil_binary = find_textutil_binary()
        if textutil_binary:
            result = run_subprocess(
                [
                    textutil_binary,
                    "-convert",
                    "docx",
                    "-output",
                    converted_path,
                    file_path,
                ],
                timeout=180,
            )
            if result.returncode == 0 and os.path.exists(converted_path):
                return converted_path
            detail = (result.stderr or result.stdout or "").strip()
            raise RuntimeError(f"textutil 转换 DOC 失败：{detail or '未知错误'}")

    raise RuntimeError("需要安装 LibreOffice 才能转换 .doc/.ppt/.xls 老格式文件")


def get_file_extension(file_path: str) -> str:
    return Path(file_path).suffix.lower().lstrip(".")


def get_source_type(file_path: str) -> str:
    ext = get_file_extension(file_path)
    if ext in ["png", "jpg", "jpeg", "webp", "bmp"]:
        return "image"
    return ext


def make_json_safe(value: Any) -> Any:
    if hasattr(value, "tolist"):
        return value.tolist()
    if isinstance(value, (list, tuple)):
        return [make_json_safe(item) for item in value]
    if isinstance(value, dict):
        return {str(key): make_json_safe(item) for key, item in value.items()}
    if isinstance(value, (str, int, float, bool)) or value is None:
        return value
    return str(value)


def get_indexed_ocr_geometry(page_result: Dict[str, Any], index: int) -> Any:
    for key in ["rec_boxes", "rec_polys", "dt_polys", "det_polys"]:
        values = page_result.get(key)
        if values is not None and index < len(values):
            return make_json_safe(values[index])
    return None


def ocr_image_with_boxes(image_path: str) -> Tuple[str, List[Dict[str, Any]]]:
    """
    对单张图片做 OCR，并保留可用于后续定位原文的坐标信息。
    """
    ocr_model = load_ocr_model()
    result = ocr_model.predict(image_path)
    lines = []
    boxes = []
    for page_result in result:
        if not page_result:
            continue

        rec_texts = page_result.get("rec_texts", [])
        rec_scores = page_result.get("rec_scores", [])
        for index, (text, score) in enumerate(zip(rec_texts, rec_scores)):
            if text and score >= 0.5:
                lines.append(text)
                boxes.append(
                    {
                        "text": text,
                        "score": float(score),
                        "box": get_indexed_ocr_geometry(page_result, index),
                    }
                )
    return "\n".join(lines), boxes


def ocr_image(image_path: str) -> str:
    text, _boxes = ocr_image_with_boxes(image_path)
    return text


def save_extracted_image(source_file: str, image_name: str, image_bytes: bytes) -> str:
    image_ext = Path(image_name).suffix.lower()
    if image_ext not in OCR_IMAGE_EXTENSIONS:
        raise ValueError(f"不支持的图片格式：{image_ext}")

    safe_source = Path(source_file).stem[:40]
    output_path = os.path.join(
        EXTRACTED_IMAGE_DIR,
        f"{safe_source}_{uuid.uuid4().hex}{image_ext}",
    )
    with open(output_path, "wb") as f:
        f.write(image_bytes)
    return output_path


def build_labeled_text(parts: List[Tuple[str, str]]) -> str:
    """
    合并直接解析文本和 OCR 文本，并去掉完全重复的行。
    """
    output = []
    seen_lines = set()

    for label, text in parts:
        lines = []
        for line in text.splitlines():
            normalized = line.strip()
            if not normalized:
                continue
            if normalized in seen_lines:
                continue
            seen_lines.add(normalized)
            lines.append(normalized)

        if lines:
            output.append(f"[{label}]\n" + "\n".join(lines))

    return "\n\n".join(output)


def extract_office_embedded_image_sections(
    file_path: str,
    source_type: str,
    media_prefix: str,
) -> List[Dict[str, Any]]:
    """
    从 Office Open XML 包里提取内嵌图片并 OCR。
    DOCX/PPTX/XLSX 本质都是 zip 包，图片通常放在 word/media、ppt/media、xl/media。
    """
    sections = []
    try:
        with zipfile.ZipFile(file_path) as archive:
            media_names = [
                name
                for name in archive.namelist()
                if name.startswith(media_prefix) and Path(name).suffix.lower() in OCR_IMAGE_EXTENSIONS
            ]

            for image_index, media_name in enumerate(sorted(media_names), start=1):
                try:
                    image_path = save_extracted_image(
                        source_file=file_path,
                        image_name=media_name,
                        image_bytes=archive.read(media_name),
                    )
                    text, ocr_boxes = ocr_image_with_boxes(image_path)
                except Exception as e:
                    text = f"OCR 失败：{e}"
                    ocr_boxes = []

                section = make_section(
                    f"[内嵌图片 OCR {image_index}：{Path(media_name).name}]\n{text}",
                    {
                        "source_type": source_type,
                        "section_type": "embedded_image",
                        "image_index": image_index,
                        "image_name": Path(media_name).name,
                        "extract_method": "paddleocr_embedded_image",
                        "ocr_boxes": json.dumps(ocr_boxes, ensure_ascii=False),
                    },
                )
                if section:
                    sections.append(section)
    except zipfile.BadZipFile:
        pass

    return sections


def render_pdf_page_to_image(doc: fitz.Document, pdf_path: str, page_index: int, dpi: int = OCR_PDF_DPI) -> str:
    page = doc[page_index]
    base_scale = dpi / 72
    max_side = max(page.rect.width * base_scale, page.rect.height * base_scale)
    scale = base_scale
    if max_side > OCR_MAX_PAGE_SIDE_PIXELS:
        scale = OCR_MAX_PAGE_SIDE_PIXELS / max(page.rect.width, page.rect.height)

    matrix = fitz.Matrix(scale, scale)
    pix = page.get_pixmap(matrix=matrix, alpha=False, colorspace=fitz.csGRAY)
    image_path = os.path.join(
        UPLOAD_DIR,
        f"{Path(pdf_path).stem}_page_{page_index + 1}.png",
    )
    pix.save(image_path)
    return image_path


def remove_file_quietly(file_path: str) -> None:
    try:
        if file_path and os.path.exists(file_path):
            os.remove(file_path)
    except OSError:
        pass


def release_memory_after_file() -> None:
    gc.collect()
    try:
        import torch

        if hasattr(torch, "cuda") and torch.cuda.is_available():
            torch.cuda.empty_cache()
        if hasattr(torch, "mps") and hasattr(torch.mps, "empty_cache"):
            torch.mps.empty_cache()
    except Exception:
        pass


def make_section(text: str, metadata: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    text = text.strip()
    if not text:
        return None
    return {"text": text, "metadata": metadata}


def normalize_pdf_line(line: str) -> str:
    return " ".join((line or "").strip().split())


def is_probable_pdf_noise_line(text: str) -> bool:
    compact = normalize_pdf_line(text)
    if not compact:
        return True
    if len(compact) <= 6 and (
        compact.isdigit()
        or compact in {"目录", "目 录"}
        or compact.strip("-—– ").isdigit()
    ):
        return True
    if re.match(r"^[\-.·•\s]+$", compact):
        return True
    if re.search(r"\.{4,}\s*\d+\s*$", compact):
        return True
    return False


def collect_pdf_layout_lines(page: fitz.Page) -> List[Dict[str, Any]]:
    page_dict = page.get_text("dict")
    lines = []
    for block in page_dict.get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            spans = line.get("spans", [])
            text = normalize_pdf_line("".join(span.get("text", "") for span in spans))
            if not text:
                continue
            sizes = [float(span.get("size", 0) or 0) for span in spans if span.get("text", "").strip()]
            bbox = line.get("bbox") or [0, 0, 0, 0]
            lines.append(
                {
                    "text": text,
                    "size": max(sizes) if sizes else 0.0,
                    "bbox": make_json_safe(bbox),
                    "y0": float(bbox[1]) if len(bbox) > 1 else 0.0,
                    "y1": float(bbox[3]) if len(bbox) > 3 else 0.0,
                    "page_height": float(page.rect.height),
                }
            )
    return lines


def find_repeated_pdf_margin_lines(layout_pages: List[List[Dict[str, Any]]]) -> set:
    page_count = len(layout_pages)
    if page_count < 3:
        return set()

    margin_counts = {}
    for lines in layout_pages:
        seen_on_page = set()
        for line in lines:
            page_height = line.get("page_height") or 1
            in_margin = line.get("y0", 0) < page_height * 0.10 or line.get("y1", 0) > page_height * 0.90
            text = normalize_pdf_line(line.get("text", ""))
            if in_margin and text and len(text) <= 80:
                seen_on_page.add(text)
        for text in seen_on_page:
            margin_counts[text] = margin_counts.get(text, 0) + 1

    threshold = max(3, int(page_count * 0.35))
    return {text for text, count in margin_counts.items() if count >= threshold}


def extract_pdf_table_text(page: fitz.Page) -> str:
    try:
        if not hasattr(page, "find_tables"):
            return ""
        tables = page.find_tables()
        table_parts = []
        for table_index, table in enumerate(tables, start=1):
            rows = table.extract()
            formatted_rows = []
            for row in rows:
                cells = [normalize_pdf_line(str(cell or "")) for cell in row]
                if any(cells):
                    formatted_rows.append(" | ".join(cells))
            if formatted_rows:
                table_parts.append(f"[表格 {table_index}]\n" + "\n".join(formatted_rows))
        return "\n\n".join(table_parts)
    except Exception:
        return ""


def build_pdf_layout_text(
    page: fitz.Page,
    lines: List[Dict[str, Any]],
    repeated_margin_lines: set,
) -> str:
    sizes = [line.get("size", 0.0) for line in lines if line.get("size", 0.0) > 0]
    avg_size = sum(sizes) / len(sizes) if sizes else 0.0
    text_lines = []
    for line in lines:
        text = normalize_pdf_line(line.get("text", ""))
        if not text or text in repeated_margin_lines or is_probable_pdf_noise_line(text):
            continue
        is_title = avg_size > 0 and line.get("size", 0.0) >= avg_size * 1.25 and len(text) <= 100
        text_lines.append(f"[标题] {text}" if is_title else text)

    table_text = extract_pdf_table_text(page)
    if table_text:
        text_lines.append(table_text)
    return "\n".join(text_lines).strip()


def extract_pdf_sections(
    pdf_path: str,
    ocr_threshold: int = 40,
    pdf_ocr_mode: str = "smart",
    progress_callback: Optional[Callable[[str], None]] = None,
) -> List[Dict[str, Any]]:
    """
    PDF 混合解析：
    smart：优先提取 PDF 文字，只有文字很少的页才 OCR。
    force：每页都做整页 OCR，覆盖“文字 + 图片”混排，但内存和耗时最高。
    text：只提取 PDF 内置文字，不做 OCR。
    """
    sections = []
    pdf_ocr_mode = pdf_ocr_mode if pdf_ocr_mode in {"smart", "force", "text"} else "smart"
    doc = fitz.open(pdf_path)
    try:
        layout_pages = [collect_pdf_layout_lines(doc[index]) for index in range(len(doc))]
        repeated_margin_lines = find_repeated_pdf_margin_lines(layout_pages)
        for page_index in range(len(doc)):
            page = doc[page_index]
            page_no = page_index + 1
            if progress_callback:
                progress_callback(f"正在解析 PDF 第 {page_no}/{len(doc)} 页")
            direct_text = build_pdf_layout_text(page, layout_pages[page_index], repeated_margin_lines)
            ocr_text = ""
            ocr_boxes = []
            extract_method = "pymupdf_text"
            compact_text_length = len(direct_text.replace("\n", "").strip())

            if pdf_ocr_mode == "force":
                image_path = render_pdf_page_to_image(doc, pdf_path, page_index)
                try:
                    if progress_callback:
                        progress_callback(f"正在 OCR PDF 第 {page_no}/{len(doc)} 页")
                    ocr_text, ocr_boxes = ocr_image_with_boxes(image_path)
                finally:
                    remove_file_quietly(image_path)
                text = build_labeled_text(
                    [
                        ("直接提取文本", direct_text),
                        ("整页 OCR 文本", ocr_text),
                    ]
                )
                extract_method = "pymupdf_text+paddleocr_page"
            elif pdf_ocr_mode == "smart" and compact_text_length < ocr_threshold:
                image_path = render_pdf_page_to_image(doc, pdf_path, page_index)
                try:
                    if progress_callback:
                        progress_callback(f"正在 OCR PDF 第 {page_no}/{len(doc)} 页")
                    text, ocr_boxes = ocr_image_with_boxes(image_path)
                finally:
                    remove_file_quietly(image_path)
                extract_method = "paddleocr_fallback"
            else:
                text = direct_text

            section = make_section(
                f"[第 {page_no} 页]\n{text}",
                {
                    "source_type": "pdf",
                    "page": page_no,
                    "extract_method": extract_method,
                    "ocr_boxes": json.dumps(ocr_boxes, ensure_ascii=False) if ocr_boxes else "",
                },
            )
            if section:
                sections.append(section)
    finally:
        doc.close()
    return sections


def extract_image_sections(image_path: str) -> List[Dict[str, Any]]:
    text, ocr_boxes = ocr_image_with_boxes(image_path)
    section = make_section(
        text,
        {
            "source_type": "image",
            "extract_method": "paddleocr",
            "ocr_boxes": json.dumps(ocr_boxes, ensure_ascii=False),
        },
    )
    return [section] if section else []


def docx_row_cells(row) -> List[str]:
    return trim_empty_tail([cell.text.strip().replace("\n", " ") for cell in row.cells])


def build_docx_row_text(table_index: int, row_number: int, cells: List[str], header_cells: Optional[List[str]]) -> str:
    parts = [f"[Word 表格 {table_index}，行 {row_number}]"]
    if header_cells:
        max_column_count = max(len(cells), len(header_cells))
        for column_index in range(max_column_count):
            value = cells[column_index].strip() if column_index < len(cells) else ""
            if not value:
                continue
            header = header_cells[column_index].strip() if column_index < len(header_cells) else ""
            parts.append(f"{header or f'第 {column_index + 1} 列'}：{value}")
    else:
        for column_index, value in enumerate(cells, start=1):
            if value.strip():
                parts.append(f"第 {column_index} 列：{value.strip()}")
    return "\n".join(parts)


def extract_docx_sections(
    docx_path: str,
    ocr_enhance: bool = True,
    progress_callback: Optional[Callable[[str], None]] = None,
) -> List[Dict[str, Any]]:
    doc = Document(docx_path)
    sections = []
    if progress_callback:
        progress_callback("正在解析 Word 段落")

    paragraphs = [paragraph.text.strip() for paragraph in doc.paragraphs if paragraph.text.strip()]
    paragraph_section = make_section(
        "\n".join(paragraphs),
        {
            "source_type": "docx",
            "section_type": "paragraphs",
            "extract_method": "python-docx",
        },
    )
    if paragraph_section:
        sections.append(paragraph_section)

    for table_index, table in enumerate(doc.tables, start=1):
        if progress_callback:
            progress_callback(f"正在解析 Word 表格 {table_index}/{len(doc.tables)}")
        header_cells = None
        for row_number, row in enumerate(table.rows, start=1):
            cells = docx_row_cells(row)
            if not any(cells):
                continue
            if header_cells is None and looks_like_xlsx_header(cells):
                header_cells = cells
                header_section = make_section(
                    f"[Word 表格 {table_index}，表头行 {row_number}]\n"
                    + "\n".join(f"第 {index + 1} 列：{cell}" for index, cell in enumerate(cells) if cell),
                    {
                        "source_type": "docx",
                        "section_type": "table_header",
                        "table_index": table_index,
                        "row_number": row_number,
                        "row_range": str(row_number),
                        "extract_method": "python-docx",
                    },
                )
                if header_section:
                    sections.append(header_section)
                continue

            row_section = make_section(
                build_docx_row_text(table_index, row_number, cells, header_cells),
                {
                    "source_type": "docx",
                    "section_type": "table_row",
                    "table_index": table_index,
                    "row_number": row_number,
                    "row_range": str(row_number),
                    "extract_method": "python-docx",
                },
            )
            if row_section:
                sections.append(row_section)

    if ocr_enhance:
        if progress_callback:
            progress_callback("正在 OCR Word 内嵌图片")
        sections.extend(
            extract_office_embedded_image_sections(
                file_path=docx_path,
                source_type="docx",
                media_prefix="word/media/",
            )
        )

    return sections


def shape_position_key(shape) -> Tuple[int, int]:
    return (
        int(getattr(shape, "top", 0) or 0),
        int(getattr(shape, "left", 0) or 0),
    )


def collect_ppt_shape_text(shape) -> List[str]:
    texts = []
    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip()
        if text:
            texts.append(text)

    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            texts.append("\n".join(rows))

    if hasattr(shape, "shapes"):
        for child_shape in sorted(shape.shapes, key=shape_position_key):
            texts.extend(collect_ppt_shape_text(child_shape))

    return texts


def extract_pptx_sections(
    pptx_path: str,
    ocr_enhance: bool = True,
    progress_callback: Optional[Callable[[str], None]] = None,
) -> List[Dict[str, Any]]:
    presentation = Presentation(pptx_path)
    sections = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        if progress_callback:
            progress_callback(f"正在解析 PPT 第 {slide_index}/{len(presentation.slides)} 页")
        slide_texts = []
        ordered_shapes = sorted(slide.shapes, key=shape_position_key)
        for shape in ordered_shapes:
            slide_texts.extend(collect_ppt_shape_text(shape))

        if slide_texts:
            first_text = slide_texts[0].strip()
            if first_text and len(first_text.replace("\n", "")) <= 100:
                slide_texts[0] = f"[标题]\n{first_text}"

        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append("备注：\n" + notes)
        except Exception:
            pass

        section = make_section(
            f"[第 {slide_index} 页幻灯片]\n" + "\n\n".join(slide_texts),
            {
                "source_type": "pptx",
                "slide": slide_index,
                "extract_method": "python-pptx",
            },
        )
        if section:
            sections.append(section)

    if ocr_enhance:
        if progress_callback:
            progress_callback("正在 OCR PPT 内嵌图片")
        sections.extend(
            extract_office_embedded_image_sections(
                file_path=pptx_path,
                source_type="pptx",
                media_prefix="ppt/media/",
            )
        )

    return sections


def format_cell_value(value: Any) -> str:
    if value is None:
        return ""
    if hasattr(value, "isoformat"):
        return value.isoformat()
    return str(value)


def trim_empty_tail(cells: List[str]) -> List[str]:
    cells = list(cells)
    while cells and not cells[-1]:
        cells.pop()
    return cells


def looks_like_xlsx_header(cells: List[str]) -> bool:
    non_empty_cells = [cell for cell in cells if cell.strip()]
    return len(non_empty_cells) >= 2


def build_xlsx_row_text(
    sheet_title: str,
    row_number: int,
    cells: List[str],
    header_cells: Optional[List[str]],
) -> str:
    parts = [f"[工作表：{sheet_title}，行 {row_number}]"]
    if header_cells:
        max_column_count = max(len(cells), len(header_cells))
        for column_index in range(max_column_count):
            value = cells[column_index].strip() if column_index < len(cells) else ""
            if not value:
                continue
            header = header_cells[column_index].strip() if column_index < len(header_cells) else ""
            label = header or f"第 {column_index + 1} 列"
            parts.append(f"{label}：{value}")
    else:
        for column_index, value in enumerate(cells, start=1):
            value = value.strip()
            if value:
                parts.append(f"第 {column_index} 列：{value}")
    return "\n".join(parts)


def extract_xlsx_sections(
    xlsx_path: str,
    rows_per_section: int = 80,
    ocr_enhance: bool = True,
    progress_callback: Optional[Callable[[str], None]] = None,
) -> List[Dict[str, Any]]:
    workbook = load_workbook(xlsx_path, data_only=True, read_only=True)
    sections = []

    try:
        for sheet_index, worksheet in enumerate(workbook.worksheets, start=1):
            if progress_callback:
                progress_callback(f"正在解析 Excel 工作表 {sheet_index}/{len(workbook.worksheets)}：{worksheet.title}")
            header_cells = None

            for row_number, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
                cells = trim_empty_tail([format_cell_value(value).strip() for value in row])
                if not any(cells):
                    continue

                if header_cells is None and looks_like_xlsx_header(cells):
                    header_cells = cells
                    header_section = make_section(
                        f"[工作表：{worksheet.title}，表头行 {row_number}]\n"
                        + "\n".join(f"第 {index + 1} 列：{cell}" for index, cell in enumerate(cells) if cell),
                        {
                            "source_type": "xlsx",
                            "sheet": worksheet.title,
                            "row_number": row_number,
                            "row_range": str(row_number),
                            "section_type": "table_header",
                            "extract_method": "openpyxl",
                        },
                    )
                    if header_section:
                        sections.append(header_section)
                    continue

                section = make_section(
                    build_xlsx_row_text(worksheet.title, row_number, cells, header_cells),
                    {
                        "source_type": "xlsx",
                        "sheet": worksheet.title,
                        "row_number": row_number,
                        "row_range": str(row_number),
                        "section_type": "table_row",
                        "extract_method": "openpyxl",
                    },
                )
                if section:
                    sections.append(section)
    finally:
        workbook.close()

    if ocr_enhance:
        if progress_callback:
            progress_callback("正在 OCR Excel 内嵌图片")
        sections.extend(
            extract_office_embedded_image_sections(
                file_path=xlsx_path,
                source_type="xlsx",
                media_prefix="xl/media/",
            )
        )

    return sections


def extract_document_sections(
    file_path: str,
    ocr_enhance: bool = True,
    pdf_ocr_mode: str = "smart",
    progress_callback: Optional[Callable[[str], None]] = None,
) -> List[Dict[str, Any]]:
    file_path = convert_legacy_office_file(file_path)
    ext = get_file_extension(file_path)
    if ext == "pdf":
        return extract_pdf_sections(file_path, pdf_ocr_mode=pdf_ocr_mode, progress_callback=progress_callback)
    if ext in ["png", "jpg", "jpeg", "webp", "bmp"]:
        if progress_callback:
            progress_callback("正在 OCR 图片")
        return extract_image_sections(file_path)
    if ext == "docx":
        return extract_docx_sections(file_path, ocr_enhance=ocr_enhance, progress_callback=progress_callback)
    if ext == "pptx":
        return extract_pptx_sections(file_path, ocr_enhance=ocr_enhance, progress_callback=progress_callback)
    if ext == "xlsx":
        return extract_xlsx_sections(file_path, ocr_enhance=ocr_enhance, progress_callback=progress_callback)
    raise ValueError("当前支持 PDF、图片、DOCX、PPTX、XLSX；旧版 DOC、PPT、XLS 会先尝试转换为新版格式。")


def sections_to_text(sections: List[Dict[str, Any]]) -> str:
    return "\n\n".join(section["text"] for section in sections)


def sections_have_text(sections: List[Dict[str, Any]]) -> bool:
    return any(section.get("text", "").strip() for section in sections)


# =========================
# 文本切分与向量入库
# =========================
def iter_text_chunks(text: str, chunk_size: int = 600, overlap: int = 100):
    yield from split_semantic_chunks(text, chunk_size=chunk_size, overlap=overlap)


def split_text(text: str, chunk_size: int = 600, overlap: int = 100) -> List[str]:
    """
    优先按标题、段落、条款和句子边界切分；超长文本再按字符兜底。
    """
    return list(iter_text_chunks(text, chunk_size=chunk_size, overlap=overlap))


def embed_texts(texts: List[str]) -> List[List[float]]:
    """
    用 BGE-M3 生成向量。
    normalize_embeddings=True 一般更适合向量相似度检索。
    """
    embedding_model = load_embedding_model()
    embeddings = embedding_model.encode(
        texts,
        normalize_embeddings=True,
        show_progress_bar=False,
        batch_size=EMBEDDING_BATCH_SIZE,
        convert_to_numpy=True,
    )
    if hasattr(embeddings, "tolist"):
        return embeddings.tolist()
    return [list(vector) for vector in embeddings]


def clean_metadata(metadata: Dict[str, Any]) -> Dict[str, Any]:
    cleaned = {}
    for key, value in metadata.items():
        if value is None:
            continue
        if isinstance(value, (str, int, float, bool)):
            cleaned[key] = value
        else:
            cleaned[key] = str(value)
    return cleaned


def split_sections(
    sections: List[Dict[str, Any]],
    chunk_size: int,
    overlap: int,
) -> List[Tuple[str, Dict[str, Any]]]:
    chunk_items = []
    for section_index, section in enumerate(sections):
        chunks = split_text(section["text"], chunk_size=chunk_size, overlap=overlap)
        for section_chunk_index, chunk in enumerate(chunks):
            metadata = {
                **section["metadata"],
                "section_index": section_index,
                "section_chunk_index": section_chunk_index,
            }
            chunk_items.append((chunk, metadata))
    return chunk_items


def iter_section_chunks(
    sections: List[Dict[str, Any]],
    chunk_size: int,
    overlap: int,
):
    for section_index, section in enumerate(sections):
        for section_chunk_index, chunk in enumerate(
            iter_text_chunks(section["text"], chunk_size=chunk_size, overlap=overlap)
        ):
            yield chunk, {
                **section["metadata"],
                "section_index": section_index,
                "section_chunk_index": section_chunk_index,
            }


def flush_vector_batch(
    chunks: List[str],
    ids: List[str],
    metadatas: List[Dict[str, Any]],
) -> int:
    if not chunks:
        return 0

    embeddings = embed_texts(chunks)
    models = import_qdrant_models()
    points = [
        models.PointStruct(
            id=ids[index],
            vector=embeddings[index],
            payload={
                "document": chunks[index],
                **metadatas[index],
            },
        )
        for index in range(len(chunks))
    ]
    vector_client.upsert(
        collection_name=COLLECTION_NAME,
        points=points,
    )
    written_count = len(chunks)
    del embeddings, points
    chunks.clear()
    ids.clear()
    metadatas.clear()
    gc.collect()
    return written_count


def add_document_to_vector_store(
    file_name: str,
    file_sha256: str,
    sections: List[Dict[str, Any]],
    chunk_size: int,
    overlap: int,
    doc_category: str,
    doc_label: str,
) -> int:
    doc_id = str(uuid.uuid4())
    chunks = []
    ids = []
    metadatas = []
    total_count = 0

    for index, (chunk, section_metadata) in enumerate(
        iter_section_chunks(sections, chunk_size=chunk_size, overlap=overlap)
    ):
        chunks.append(chunk)
        ids.append(str(uuid.uuid4()))
        metadatas.append(
            clean_metadata(
                {
                    **section_metadata,
                    "doc_id": doc_id,
                    "file_sha256": file_sha256,
                    "file_name": file_name,
                    "doc_label": doc_label or file_name,
                    "doc_category": doc_category,
                    "doc_category_name": DOC_CATEGORY_NAMES.get(doc_category, doc_category),
                    "chunk_index": index,
                    "embedding_model": EMBEDDING_MODEL_NAME,
                }
            )
        )

        if len(chunks) >= VECTOR_ADD_BATCH_SIZE:
            total_count += flush_vector_batch(chunks, ids, metadatas)

    total_count += flush_vector_batch(chunks, ids, metadatas)
    return total_count


def run_background_ingest_task(
    task_id: str,
    file_specs: List[Dict[str, Any]],
    settings: Dict[str, Any],
) -> None:
    success_count = 0
    duplicate_count = 0
    skipped_count = 0
    failed_count = 0
    total_files = len(file_specs)

    def sync_task(index: int, current_file: str, message: str, status: str = "running") -> None:
        if status == "running":
            wait_if_task_paused_or_cancelled(task_id)
        update_ingest_task(
            task_id,
            status=status,
            processed_files=index,
            success_count=success_count,
            duplicate_count=duplicate_count,
            skipped_count=skipped_count,
            failed_count=failed_count,
            current_file=current_file,
            message=message,
        )

    try:
        for file_index, spec in enumerate(file_specs, start=1):
            wait_if_task_paused_or_cancelled(task_id)
            relative_name = spec["relative_name"]
            file_path = spec.get("file_path", "")
            file_sha256 = spec.get("sha256", "")
            sync_task(file_index - 1, relative_name, f"正在处理：{relative_name}")

            try:
                wait_if_task_paused_or_cancelled(task_id)
                if not is_supported_upload(relative_name):
                    skipped_count += 1
                    message = f"不支持的文件类型：{Path(relative_name).suffix.lower() or '无扩展名'}"
                    record_ingest_task_item(task_id, relative_name, "unsupported", message, file_sha256=file_sha256)
                    sync_task(file_index, relative_name, message)
                    continue

                existing_file = get_duplicate_ingested_file(file_sha256)
                if existing_file:
                    duplicate_count += 1
                    record_ingest_task_item(
                        task_id,
                        relative_name,
                        "duplicate",
                        "已跳过重复文件",
                        chunk_count=int(existing_file.get("chunk_count", 0) or 0),
                        file_sha256=file_sha256,
                    )
                    sync_task(file_index, relative_name, f"已跳过重复文件：{relative_name}")
                    continue

                if is_legacy_office_file(relative_name):
                    ext = get_file_extension_from_name(relative_name)
                    can_convert, conversion_message = get_legacy_conversion_status(ext)
                    if not can_convert:
                        skipped_count += 1
                        record_ingest_task_item(task_id, relative_name, "unsupported", conversion_message, file_sha256=file_sha256)
                        sync_task(file_index, relative_name, conversion_message)
                        continue

                def extraction_progress(message: str) -> None:
                    sync_task(file_index - 1, relative_name, message)

                sections = extract_document_sections(
                    file_path,
                    ocr_enhance=bool(settings.get("ocr_enhance", True)),
                    pdf_ocr_mode=str(settings.get("pdf_ocr_mode", "smart")),
                    progress_callback=extraction_progress,
                )

                wait_if_task_paused_or_cancelled(task_id)
                if not sections_have_text(sections):
                    skipped_count += 1
                    record_ingest_task_item(task_id, relative_name, "skipped", "没有解析到有效文字", file_sha256=file_sha256)
                    sync_task(file_index, relative_name, "没有解析到有效文字")
                    continue

                deleted_chunks, _deleted_records = replace_existing_same_name_if_needed(
                    relative_name,
                    file_sha256,
                    enabled=bool(settings.get("replace_changed_same_name", True)),
                )
                chunk_count = add_document_to_vector_store(
                    file_name=relative_name,
                    file_sha256=file_sha256,
                    sections=sections,
                    chunk_size=int(settings.get("chunk_size", 600)),
                    overlap=int(settings.get("overlap", 100)),
                    doc_category=str(settings.get("doc_category", "general")),
                    doc_label=str(settings.get("doc_label") or relative_name),
                )
                if chunk_count > 0:
                    success_count += 1
                    record_ingested_file(
                        file_sha256=file_sha256,
                        file_name=relative_name,
                        doc_category=str(settings.get("doc_category", "general")),
                        doc_label=str(settings.get("doc_label") or relative_name),
                        chunk_count=chunk_count,
                    )
                    replace_note = f"，已替换旧 {deleted_chunks} 个 chunk" if deleted_chunks else ""
                    record_ingest_task_item(
                        task_id,
                        relative_name,
                        "success",
                        f"入库成功，写入 {chunk_count} 个 chunk{replace_note}",
                        chunk_count=chunk_count,
                        file_sha256=file_sha256,
                    )
                    sync_task(file_index, relative_name, f"入库成功：{relative_name}")
                else:
                    skipped_count += 1
                    record_ingest_task_item(task_id, relative_name, "skipped", "解析成功但没有可入库 chunk", file_sha256=file_sha256)
                    sync_task(file_index, relative_name, "解析成功但没有可入库 chunk")
            except IngestTaskCancelled:
                raise
            except Exception as e:
                failed_count += 1
                record_ingest_task_item(task_id, relative_name, "failed", str(e), file_sha256=file_sha256)
                sync_task(file_index, relative_name, f"处理失败：{e}")
            finally:
                release_memory_after_file()
    except IngestTaskCancelled:
        update_ingest_task(
            task_id,
            status="cancelled",
            success_count=success_count,
            duplicate_count=duplicate_count,
            skipped_count=skipped_count,
            failed_count=failed_count,
            message="入库任务已终止",
        )
        return

    if settings.get("auto_unload_models_after_ingest"):
        try:
            load_ocr_model.clear()
            load_embedding_model.clear()
            load_reranker_model.clear()
        except Exception:
            pass
        release_memory_after_file()

    sync_task(total_files, "", "后台入库任务完成", status="completed")


# =========================
# 检索与本地大模型回答
# =========================
def count_chunks(doc_category: Optional[str] = None) -> int:
    result = vector_client.count(
        collection_name=COLLECTION_NAME,
        count_filter=build_qdrant_filter({"doc_category": doc_category}) if doc_category else None,
        exact=True,
    )
    return int(result.count)


def default_fetch_k(top_k: int) -> int:
    return max(top_k, min(50, max(DEFAULT_RETRIEVAL_FETCH_K, top_k * 4)))


def query_qdrant_points(
    query_embedding: List[float],
    top_k: int,
    doc_category: Optional[str] = None,
) -> List[Dict[str, Any]]:
    query_filter = build_qdrant_filter({"doc_category": doc_category}) if doc_category else None
    try:
        result = vector_client.query_points(
            collection_name=COLLECTION_NAME,
            query=query_embedding,
            query_filter=query_filter,
            limit=top_k,
            with_payload=True,
            with_vectors=False,
        )
        points = getattr(result, "points", result)
    except Exception:
        points = vector_client.search(
            collection_name=COLLECTION_NAME,
            query_vector=query_embedding,
            query_filter=query_filter,
            limit=top_k,
            with_payload=True,
            with_vectors=False,
        )
    return [payload_to_result(point, score=getattr(point, "score", None)) for point in points]


def keyword_search_vector_store(
    query: str,
    top_k: int,
    doc_category: Optional[str] = None,
) -> List[Dict[str, Any]]:
    points = qdrant_scroll_points(where={"doc_category": doc_category} if doc_category else None, limit=100000)
    ids = [str(point.id) for point in points]
    payloads = [point_payload(point) for point in points]
    docs = [payload.get("document", "") for payload in payloads]
    metadatas = []
    for payload in payloads:
        metadata = dict(payload)
        metadata.pop("document", None)
        metadatas.append(metadata)
    ranked = keyword_rank_documents(query, docs, top_k=top_k)

    output = []
    for doc_index, score in ranked:
        output.append(
            {
                "id": ids[doc_index] if doc_index < len(ids) else str(uuid.uuid4()),
                "content": docs[doc_index],
                "metadata": metadatas[doc_index] if doc_index < len(metadatas) else {},
                "distance": None,
                "keyword_score": score,
                "retrieval_source": "keyword",
            }
        )
    return output


def merge_retrieval_results(
    vector_results: List[Dict[str, Any]],
    keyword_results: List[Dict[str, Any]],
) -> List[Dict[str, Any]]:
    if not keyword_results:
        return vector_results
    if not vector_results:
        return keyword_results

    vector_ranked = [(item["id"], item) for item in vector_results]
    keyword_ranked = [(item["id"], item) for item in keyword_results]
    merged = reciprocal_rank_merge([vector_ranked, keyword_ranked])

    vector_by_id = {item["id"]: item for item in vector_results}
    keyword_by_id = {item["id"]: item for item in keyword_results}
    for item in merged:
        item_id = item["id"]
        if item_id in vector_by_id and item_id in keyword_by_id:
            item["retrieval_source"] = "vector+keyword"
            item["keyword_score"] = keyword_by_id[item_id].get("keyword_score")
            item["distance"] = vector_by_id[item_id].get("distance")
    return merged


def rerank_search_results(query: str, candidates: List[Dict[str, Any]], top_k: int) -> List[Dict[str, Any]]:
    if not candidates:
        return []
    reranker = load_reranker_model()
    pairs = [[query, item["content"]] for item in candidates]
    scores = reranker.predict(pairs, show_progress_bar=False)
    for item, score in zip(candidates, scores):
        item["rerank_score"] = float(score)
    candidates.sort(key=lambda item: item.get("rerank_score", 0.0), reverse=True)
    return candidates[:top_k]


def int_metadata_value(metadata: Dict[str, Any], key: str) -> Optional[int]:
    value = metadata.get(key)
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def retrieval_dedupe_key(item: Dict[str, Any]) -> Tuple[Any, ...]:
    metadata = item.get("metadata") or {}
    return (
        metadata.get("file_name", ""),
        metadata.get("source_type", ""),
        metadata.get("page", ""),
        metadata.get("slide", ""),
        metadata.get("sheet", ""),
        metadata.get("row_range", ""),
        metadata.get("table_index", ""),
        metadata.get("image_index", ""),
        metadata.get("section_index", ""),
    )


def merge_retrieval_item_content(existing: Dict[str, Any], item: Dict[str, Any]) -> None:
    existing_content = existing.get("content", "")
    new_content = item.get("content", "")
    if new_content and new_content not in existing_content and len(existing_content) < 3200:
        existing["content"] = f"{existing_content}\n\n--- 同来源相邻片段 ---\n\n{new_content}".strip()

    existing_metadata = existing.setdefault("metadata", {})
    new_metadata = item.get("metadata") or {}
    existing_chunks = str(existing_metadata.get("merged_chunk_indices", existing_metadata.get("chunk_index", "")))
    new_chunk = str(new_metadata.get("chunk_index", ""))
    chunk_values = [value for value in [existing_chunks, new_chunk] if value]
    existing_metadata["merged_chunk_indices"] = ",".join(dict.fromkeys(",".join(chunk_values).split(",")))
    existing_metadata["merged_count"] = int(existing_metadata.get("merged_count", 1)) + 1

    if existing.get("distance") is None or (
        item.get("distance") is not None and item.get("distance") < existing.get("distance")
    ):
        existing["distance"] = item.get("distance")


def dedupe_search_results(results: List[Dict[str, Any]], top_k: int) -> List[Dict[str, Any]]:
    output = []
    by_key: Dict[Tuple[Any, ...], Dict[str, Any]] = {}
    for item in results:
        key = retrieval_dedupe_key(item)
        existing = by_key.get(key)
        if existing is None:
            if len(output) >= top_k:
                continue
            copied = {
                **item,
                "metadata": dict(item.get("metadata") or {}),
            }
            by_key[key] = copied
            output.append(copied)
        else:
            existing_chunk = int_metadata_value(existing.get("metadata") or {}, "chunk_index")
            new_chunk = int_metadata_value(item.get("metadata") or {}, "chunk_index")
            if existing_chunk is None or new_chunk is None or abs(existing_chunk - new_chunk) <= 2:
                merge_retrieval_item_content(existing, item)
            elif len(output) < top_k:
                copied = {
                    **item,
                    "metadata": dict(item.get("metadata") or {}),
                }
                output.append(copied)
    return output[:top_k]


def search_vector_store(
    query: str,
    top_k: int = 5,
    doc_category: Optional[str] = None,
    max_distance: Optional[float] = None,
    fetch_k: Optional[int] = None,
    use_hybrid: bool = DEFAULT_USE_HYBRID_SEARCH,
    use_reranker: bool = DEFAULT_USE_RERANKER,
) -> List[Dict[str, Any]]:
    available_count = count_chunks(doc_category)
    if available_count <= 0:
        return []

    query_embedding = embed_texts([query])[0]
    n_results = fetch_k if fetch_k is not None else default_fetch_k(top_k)
    vector_results = query_qdrant_points(
        query_embedding,
        top_k=min(max(top_k, n_results), available_count),
        doc_category=doc_category,
    )
    keyword_results = keyword_search_vector_store(query, top_k=n_results, doc_category=doc_category) if use_hybrid else []
    candidates = merge_retrieval_results(vector_results, keyword_results)

    filtered = []
    for item in candidates:
        distance = item.get("distance")
        if max_distance is not None and distance is not None and distance > max_distance:
            if not item.get("keyword_score"):
                continue
        filtered.append(item)

    if use_reranker:
        return dedupe_search_results(rerank_search_results(query, filtered[:n_results], top_k=max(top_k, n_results)), top_k)

    output = []
    for item in filtered:
        output.append(item)
        if len(output) >= n_results:
            break
    return dedupe_search_results(output, top_k)


def is_complex_retrieval_question(question: str) -> bool:
    compact = question.strip()
    if len(compact) >= 60:
        return True
    return any(marker in compact for marker in ["分别", "逐条", "对照", "以及", "并且", "同时", "、", "；", ";"])


def heuristic_decompose_question(question: str, max_queries: int = 4) -> List[str]:
    parts = re.split(r"[；;]|以及|并且|同时|分别|、", question)
    queries = [part.strip(" ，,。？?") for part in parts if part.strip(" ，,。？?")]
    if len(queries) <= 1:
        return [question]
    return queries[:max_queries]


def decompose_retrieval_query(
    query: str,
    enabled: bool,
    mode: str,
    purpose: str,
    max_queries: int = 4,
) -> Tuple[List[str], Optional[str]]:
    if not enabled or not is_complex_retrieval_question(query):
        return [query], None

    system_prompt = """
你是一个检索查询拆解器。
你的任务是把复杂问题拆成 2 到 4 个可以分别检索的中文子问题。
只输出子问题，每行一个；不要回答，不要解释。
如果问题不需要拆解，只输出原问题。
"""
    purpose_hint = "合规分析" if purpose == "compliance" else "本地资料问答"
    user_prompt = f"""
用途：{purpose_hint}
原始检索问题：
{query}

请拆成多个独立、具体、便于检索的子问题。
"""
    try:
        response = create_llm_chat_completion(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            temperature=0,
            mode=mode,
        )
        raw_text = (response.choices[0].message.content or "").strip()
        queries = []
        for line in raw_text.splitlines():
            cleaned = re.sub(r"^\s*[-*\d.、）)]+", "", line).strip()
            if cleaned and cleaned not in queries:
                queries.append(cleaned)
        if not queries:
            queries = heuristic_decompose_question(query, max_queries=max_queries)
        return queries[:max_queries], None
    except Exception as e:
        return heuristic_decompose_question(query, max_queries=max_queries), str(e)


def search_vector_store_multi_query(
    queries: List[str],
    top_k: int = 5,
    doc_category: Optional[str] = None,
    max_distance: Optional[float] = None,
    fetch_k: Optional[int] = None,
    use_hybrid: bool = DEFAULT_USE_HYBRID_SEARCH,
    use_reranker: bool = DEFAULT_USE_RERANKER,
) -> List[Dict[str, Any]]:
    cleaned_queries = [query.strip() for query in queries if query and query.strip()]
    if not cleaned_queries:
        return []
    if len(cleaned_queries) == 1:
        return search_vector_store(
            cleaned_queries[0],
            top_k=top_k,
            doc_category=doc_category,
            max_distance=max_distance,
            fetch_k=fetch_k,
            use_hybrid=use_hybrid,
            use_reranker=use_reranker,
        )

    ranked_lists = []
    results_by_id: Dict[str, Dict[str, Any]] = {}
    per_query_top_k = max(top_k, min(fetch_k or default_fetch_k(top_k), 12))
    for sub_query in cleaned_queries:
        sub_results = search_vector_store(
            sub_query,
            top_k=per_query_top_k,
            doc_category=doc_category,
            max_distance=max_distance,
            fetch_k=fetch_k,
            use_hybrid=use_hybrid,
            use_reranker=use_reranker,
        )
        ranked = []
        for item in sub_results:
            item_id = item.get("id", str(uuid.uuid4()))
            copied = dict(item)
            copied["matched_query"] = sub_query
            results_by_id[item_id] = copied
            ranked.append((item_id, copied))
        ranked_lists.append(ranked)

    merged = reciprocal_rank_merge(ranked_lists)
    return dedupe_search_results(merged, top_k)


def search_with_min_coverage(
    queries: List[str],
    top_k: int,
    min_results: int,
    doc_category: str,
    max_distance: Optional[float],
    fetch_k: Optional[int],
    use_hybrid: bool,
    use_reranker: bool,
) -> List[Dict[str, Any]]:
    min_results = max(0, min(min_results, top_k))
    results = search_vector_store_multi_query(
        queries,
        top_k=top_k,
        doc_category=doc_category,
        max_distance=max_distance,
        fetch_k=fetch_k,
        use_hybrid=use_hybrid,
        use_reranker=use_reranker,
    )
    if len(results) >= min_results or max_distance is None:
        return results

    relaxed_results = search_vector_store_multi_query(
        queries,
        top_k=top_k,
        doc_category=doc_category,
        max_distance=None,
        fetch_k=fetch_k,
        use_hybrid=use_hybrid,
        use_reranker=use_reranker,
    )
    existing_ids = {item.get("id") for item in results}
    for item in relaxed_results:
        if len(results) >= min_results:
            break
        if item.get("id") in existing_ids:
            continue
        item["coverage_relaxed"] = True
        results.append(item)
    return dedupe_search_results(results, top_k)


def describe_source(metadata: Dict[str, Any]) -> str:
    parts = []
    if metadata.get("page"):
        parts.append(f"页码：{metadata.get('page')}")
    if metadata.get("slide"):
        parts.append(f"幻灯片：{metadata.get('slide')}")
    if metadata.get("sheet"):
        parts.append(f"工作表：{metadata.get('sheet')}")
    if metadata.get("row_range"):
        parts.append(f"行：{metadata.get('row_range')}")
    if metadata.get("table_index"):
        parts.append(f"表格：{metadata.get('table_index')}")
    if metadata.get("image_index"):
        image_label = f"图片：{metadata.get('image_index')}"
        if metadata.get("image_name"):
            image_label += f"（{metadata.get('image_name')}）"
        parts.append(image_label)
    return "；".join(parts) if parts else "无"


def build_context(search_results: List[Dict[str, Any]], title: str = "检索资料") -> str:
    context_parts = [f"## {title}"]
    for i, item in enumerate(search_results, start=1):
        metadata = item["metadata"]
        content = item["content"]
        file_name = metadata.get("file_name", "未知文件")
        chunk_index = metadata.get("chunk_index", "未知片段")
        doc_category_name = metadata.get("doc_category_name", "未知类型")
        context_parts.append(
            f"【资料 {i}】\n"
            f"资料类型：{doc_category_name}\n"
            f"来源文件：{file_name}\n"
            f"片段编号：{chunk_index}\n"
            f"来源位置：{describe_source(metadata)}\n"
            f"内容：\n{content}\n"
        )
    return "\n\n".join(context_parts)


def build_chat_history(
    chat_messages: List[Dict[str, Any]],
    max_messages: int = DEFAULT_CONTEXT_TURNS * 2,
) -> str:
    if not chat_messages or max_messages <= 0:
        return "无"

    history_parts = []
    for message in chat_messages[-max_messages:]:
        role_name = "用户" if message.get("role") == "user" else "助手"
        content = str(message.get("content", "")).strip()
        if content:
            history_parts.append(f"{role_name}：{content}")
    return "\n".join(history_parts) if history_parts else "无"


def rewrite_retrieval_query(
    question: str,
    chat_history: Optional[List[Dict[str, Any]]] = None,
    mode: str = "fast",
    context_turns: int = DEFAULT_CONTEXT_TURNS,
    purpose: str = "rag",
) -> str:
    history = build_chat_history(chat_history or [], max_messages=context_turns * 2)
    if history == "无":
        return question

    purpose_hint = (
        "用于合规差距分析，需要同时覆盖监管要求、规章制度、企业资料、问题点和整改方向。"
        if purpose == "compliance"
        else "用于本地文档问答，需要保留用户真正要检索的实体、制度、事项和限定条件。"
    )
    system_prompt = """
你是一个检索问题改写器。
你的任务是把用户当前问题结合对话历史，改写成一个适合向量数据库检索的完整中文问题。
只输出改写后的问题，不要回答，不要解释，不要添加列表。
如果当前问题已经完整，原样或轻微补全即可。
"""
    user_prompt = f"""
用途：{purpose_hint}

对话历史：
{history}

当前问题：
{question}

请输出一个完整、明确、适合检索的中文问题。
"""
    response = create_llm_chat_completion(
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=0,
        mode=mode,
    )
    rewritten_query = (response.choices[0].message.content or "").strip()
    return rewritten_query or question


def maybe_rewrite_retrieval_query(
    question: str,
    chat_history: Optional[List[Dict[str, Any]]],
    enabled: bool,
    mode: str,
    context_turns: int,
    purpose: str,
) -> Tuple[str, Optional[str]]:
    if not enabled:
        return question, None
    try:
        return rewrite_retrieval_query(
            question,
            chat_history=chat_history,
            mode=mode,
            context_turns=context_turns,
            purpose=purpose,
        ), None
    except Exception as e:
        return question, str(e)


def ask_llm(
    question: str,
    search_results: List[Dict[str, Any]],
    chat_history: Optional[List[Dict[str, Any]]] = None,
    mode: str = "fast",
    context_turns: int = DEFAULT_CONTEXT_TURNS,
) -> str:
    context = build_context(search_results)
    history = build_chat_history(chat_history or [], max_messages=context_turns * 2)
    system_prompt = """
你是一个严谨的本地文档问答助手。
你只能根据用户提供的【检索资料】回答问题。
如果资料中没有明确答案，请直接说“根据当前资料无法确定”，不要编造。
回答时尽量引用来源文件、片段编号和来源位置。
如果【对话历史】与【检索资料】冲突，以【检索资料】为准。
"""
    user_prompt = f"""
下面是当前对话历史：
{history}

下面是从向量数据库检索出来的资料：
{context}

当前用户问题：
{question}

请基于上述资料回答。
"""
    response = create_llm_chat_completion(
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=0.2,
        mode=mode,
    )
    return response.choices[0].message.content


def ask_llm_general(
    question: str,
    chat_history: Optional[List[Dict[str, Any]]] = None,
    mode: str = "fast",
    context_turns: int = DEFAULT_CONTEXT_TURNS,
) -> str:
    history = build_chat_history(chat_history or [], max_messages=context_turns * 2)
    system_prompt = """
你是一个可靠的中文助手。
当没有可用的本地检索资料，或用户只是普通寒暄时，可以进行普通对话。
如果问题涉及用户本地文档、企业资料或监管要求，请提醒用户先上传资料入库。
"""
    user_prompt = f"""
下面是当前对话历史：
{history}

当前用户问题：
{question}
"""
    response = create_llm_chat_completion(
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=0.4,
        mode=mode,
    )
    return response.choices[0].message.content


def is_likely_general_chat_question(question: str) -> bool:
    compact_question = "".join(question.strip().lower().split())
    return compact_question in {
        "你好",
        "您好",
        "hello",
        "hi",
        "嗨",
        "在吗",
        "谢谢",
        "感谢",
        "ok",
    }


def ask_llm_compliance(
    topic: str,
    regulation_results: List[Dict[str, Any]],
    enterprise_results: List[Dict[str, Any]],
    chat_history: Optional[List[Dict[str, Any]]] = None,
    mode: str = "fast",
    context_turns: int = DEFAULT_CONTEXT_TURNS,
    clause_by_clause: bool = False,
    include_missing_list: bool = True,
) -> str:
    history = build_chat_history(chat_history or [], max_messages=context_turns * 2)
    regulation_context = build_context(regulation_results, title="监管要求 / 规章制度")
    enterprise_context = build_context(enterprise_results, title="企业资料")
    clause_prompt = (
        "请按监管资料中的条款或要求逐条对照企业资料；每一行只分析一个监管条款或要求。"
        if clause_by_clause
        else "请围绕用户主题归纳主要合规差距。"
    )
    missing_prompt = (
        "表格后必须追加“### 资料不足清单”，列出为了进一步判断还需要补充的企业资料。"
        if include_missing_list
        else ""
    )
    system_prompt = f"""
你是一个严谨的合规差距分析助手。
你必须同时参考【监管要求 / 规章制度】和【企业资料】。
不要编造监管要求，也不要在企业资料不足时直接判定企业违规。
每个结论都必须引用监管资料和企业资料的来源文件、片段编号、来源位置。
如果企业资料没有对应证据，请标记为“资料不足，需补充”。
{clause_prompt}
请先输出 Markdown 表格，列为：问题点、对应监管要求、企业资料证据、风险等级、整改建议、引用来源。
{missing_prompt}
表格后可以再补充必要说明。
"""
    user_prompt = f"""
下面是当前合规分析对话历史：
{history}

分析主题：
{topic}

【监管要求 / 规章制度】
{regulation_context}

【企业资料】
{enterprise_context}

请基于上述资料做合规差距分析。
"""
    response = create_llm_chat_completion(
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=0.2,
        mode=mode,
    )
    return response.choices[0].message.content


# =========================
# UI 辅助
# =========================
def render_search_results(search_results: List[Dict[str, Any]]) -> None:
    for i, item in enumerate(search_results, start=1):
        metadata = item["metadata"]
        distance = item["distance"]
        retrieval_source = item.get("retrieval_source", "vector")
        rerank_score = item.get("rerank_score")
        extra = f"｜来源 {retrieval_source}"
        if rerank_score is not None:
            extra += f"｜重排 {rerank_score:.4f}"
        if item.get("coverage_relaxed"):
            extra += "｜覆盖补充"
        if metadata.get("merged_count"):
            extra += f"｜合并 {metadata.get('merged_count')} 段"
        title = (
            f"资料 {i}｜{metadata.get('doc_category_name')}｜{metadata.get('file_name')}｜"
            f"片段 {metadata.get('chunk_index')}｜距离 {distance}{extra}"
        )
        with st.expander(title):
            st.write(item["content"])
            st.json(metadata)


def extract_missing_info_items(answer: str) -> List[str]:
    if not answer:
        return []
    match = re.search(r"#{1,4}\s*资料不足清单\s*(.*)", answer, flags=re.S)
    if not match:
        return []
    section_text = match.group(1)
    next_heading = re.search(r"\n#{1,4}\s+", section_text)
    if next_heading:
        section_text = section_text[: next_heading.start()]
    items = []
    for line in section_text.splitlines():
        cleaned = re.sub(r"^\s*[-*•\d.、）)]+", "", line).strip()
        if cleaned and not cleaned.startswith("|"):
            items.append(cleaned)
    return items


def write_rows_to_sheet(sheet, rows: List[Dict[str, Any]]) -> None:
    if not rows:
        sheet.append(["无"])
        return
    headers = list(rows[0].keys())
    sheet.append(headers)
    for row in rows:
        sheet.append([row.get(header, "") for header in headers])


def evidence_rows_for_report(results: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    rows = []
    for index, item in enumerate(results, start=1):
        metadata = item.get("metadata") or {}
        rows.append(
            {
                "序号": index,
                "资料类型": metadata.get("doc_category_name", ""),
                "来源文件": metadata.get("file_name", ""),
                "来源位置": describe_source(metadata),
                "片段编号": metadata.get("chunk_index", ""),
                "合并片段": metadata.get("merged_chunk_indices", ""),
                "距离": item.get("distance", ""),
                "重排分数": item.get("rerank_score", ""),
                "内容": item.get("content", ""),
            }
        )
    return rows


def build_compliance_report_excel(
    answer: str,
    structured_rows: List[Dict[str, Any]],
    regulation_results: List[Dict[str, Any]],
    enterprise_results: List[Dict[str, Any]],
    retrieval_query: str,
) -> bytes:
    workbook = Workbook()
    summary_sheet = workbook.active
    summary_sheet.title = "分析摘要"
    summary_sheet.append(["生成时间", time.strftime("%Y-%m-%d %H:%M:%S")])
    summary_sheet.append(["检索问题", retrieval_query])
    summary_sheet.append([])
    summary_sheet.append(["回答正文"])
    for line in (answer or "").splitlines():
        summary_sheet.append([line])

    table_sheet = workbook.create_sheet("结构化分析表")
    write_rows_to_sheet(table_sheet, structured_rows)

    missing_sheet = workbook.create_sheet("资料不足清单")
    missing_items = extract_missing_info_items(answer)
    if missing_items:
        missing_sheet.append(["序号", "需补充资料"])
        for index, item in enumerate(missing_items, start=1):
            missing_sheet.append([index, item])
    else:
        missing_sheet.append(["无明确资料不足清单"])

    reg_sheet = workbook.create_sheet("监管证据")
    write_rows_to_sheet(reg_sheet, evidence_rows_for_report(regulation_results))

    ent_sheet = workbook.create_sheet("企业证据")
    write_rows_to_sheet(ent_sheet, evidence_rows_for_report(enterprise_results))

    output = io.BytesIO()
    workbook.save(output)
    return output.getvalue()


def render_rag_chat_panel(messages: List[Dict[str, Any]]) -> None:
    chat_panel = st.container(height=CHAT_PANEL_HEIGHT, border=True)
    with chat_panel:
        if not messages:
            st.info("输入问题后会保留上下文；每一轮都会按当前检索设置重新召回资料。")

        for message in messages:
            role = message.get("role", "assistant")
            with st.chat_message(role):
                search_results = message.get("search_results") or []
                st.write(message.get("content", ""))
                if role == "assistant" and message.get("mode_label"):
                    st.caption(f"模式：{message['mode_label']}")
                if role == "assistant" and message.get("retrieval_query"):
                    st.caption(f"本轮检索问题：{message['retrieval_query']}")
                if role == "assistant" and message.get("retrieval_sub_queries"):
                    st.caption("拆解子问题：" + "；".join(message["retrieval_sub_queries"]))
                if role == "assistant" and message.get("query_rewrite_error"):
                    st.caption(f"问题改写失败，已使用原问题检索：{message['query_rewrite_error']}")
                if role == "assistant" and message.get("query_decompose_error"):
                    st.caption(f"问题拆解失败，已使用启发式拆解：{message['query_decompose_error']}")
                if role == "assistant" and search_results:
                    with st.expander("本轮检索资料", expanded=False):
                        render_search_results(search_results)


def render_compliance_chat_panel(messages: List[Dict[str, Any]]) -> None:
    chat_panel = st.container(height=CHAT_PANEL_HEIGHT, border=True)
    with chat_panel:
        if not messages:
            st.info("输入合规分析问题后会保留上下文；每一轮都会分别检索监管资料和企业资料。")

        for message_index, message in enumerate(messages):
            role = message.get("role", "assistant")
            with st.chat_message(role):
                regulation_results = message.get("regulation_results") or []
                enterprise_results = message.get("enterprise_results") or []
                st.write(message.get("content", ""))
                structured_rows = message.get("structured_rows") or []
                if role == "assistant" and structured_rows:
                    with st.expander("结构化合规分析表", expanded=True):
                        st.dataframe(structured_rows, width="stretch", hide_index=True)
                if role == "assistant" and (structured_rows or regulation_results or enterprise_results):
                    report_bytes = build_compliance_report_excel(
                        answer=message.get("content", ""),
                        structured_rows=structured_rows,
                        regulation_results=regulation_results,
                        enterprise_results=enterprise_results,
                        retrieval_query=message.get("retrieval_query", ""),
                    )
                    st.download_button(
                        "导出本轮合规分析 Excel",
                        data=report_bytes,
                        file_name=f"compliance_report_{message_index + 1}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        key=f"download_compliance_report_{message_index}",
                    )
                if role == "assistant" and message.get("mode_label"):
                    st.caption(f"模式：{message['mode_label']}")
                if role == "assistant" and message.get("retrieval_query"):
                    st.caption(f"本轮检索问题：{message['retrieval_query']}")
                if role == "assistant" and message.get("retrieval_sub_queries"):
                    st.caption("拆解子问题：" + "；".join(message["retrieval_sub_queries"]))
                if role == "assistant" and message.get("query_rewrite_error"):
                    st.caption(f"问题改写失败，已使用原问题检索：{message['query_rewrite_error']}")
                if role == "assistant" and message.get("query_decompose_error"):
                    st.caption(f"问题拆解失败，已使用启发式拆解：{message['query_decompose_error']}")

                if role == "assistant" and (regulation_results or enterprise_results):
                    with st.expander("本轮证据", expanded=False):
                        evidence_col1, evidence_col2 = st.columns(2)
                        with evidence_col1:
                            st.markdown("##### 监管 / 规章")
                            if regulation_results:
                                render_search_results(regulation_results)
                            else:
                                st.info("无监管证据")
                        with evidence_col2:
                            st.markdown("##### 企业资料")
                            if enterprise_results:
                                render_search_results(enterprise_results)
                            else:
                                st.info("无企业证据")


def render_library_summary() -> None:
    total_count = count_chunks()
    regulation_count = count_chunks("regulation")
    enterprise_count = count_chunks("enterprise")
    general_count = count_chunks("general")

    col_total, col_reg, col_ent, col_general = st.columns(4)
    col_total.metric("总 chunk", total_count)
    col_reg.metric("监管/规章", regulation_count)
    col_ent.metric("企业资料", enterprise_count)
    col_general.metric("其他资料", general_count)


def render_result_dataframe(rows: List[Dict[str, Any]], max_rows: int = 200) -> None:
    if not rows:
        return

    st.dataframe(rows[:max_rows], width="stretch", hide_index=True)
    if len(rows) > max_rows:
        st.caption(f"仅展示前 {max_rows} 行，共 {len(rows)} 行。")


TASK_STATUS_LABELS = {
    "running": "运行中",
    "pause_requested": "暂停中",
    "paused": "已暂停",
    "cancel_requested": "终止中",
    "cancelled": "已终止",
    "completed": "已完成",
}


def format_task_status(status: str) -> str:
    return TASK_STATUS_LABELS.get(status, status or "未知")


def render_ingest_task_controls(task: Dict[str, Any]) -> None:
    task_id = task["id"]
    status = task["status"]
    control_cols = st.columns([2, 1, 1, 1])
    with control_cols[0]:
        st.caption(
            f"{format_task_status(status)}｜"
            f"{task['processed_files']}/{task['total_files']}｜"
            f"{task['message']}"
        )
    with control_cols[1]:
        if st.button(
            "暂停",
            key=f"pause_task_{task_id}",
            disabled=status not in {"running"},
        ):
            request_pause_ingest_task(task_id)
            st.rerun()
    with control_cols[2]:
        if st.button(
            "继续",
            key=f"resume_task_{task_id}",
            disabled=status not in {"paused", "pause_requested"},
        ):
            resume_ingest_task(task_id)
            st.rerun()
    with control_cols[3]:
        if st.button(
            "终止",
            key=f"cancel_task_{task_id}",
            disabled=status not in {"running", "pause_requested", "paused"},
        ):
            request_cancel_ingest_task(task_id)
            st.rerun()


def format_session_option(session: Dict[str, Any]) -> str:
    updated_at = time.strftime("%m-%d %H:%M", time.localtime(float(session["updated_at"])))
    return f"{session['title']} ｜ {updated_at}"


def get_session_select_key(session_type: str) -> str:
    return f"{session_type}_session_select"


def sync_selected_chat_session(session_type: str) -> None:
    selected_id = st.session_state.get(get_session_select_key(session_type))
    if selected_id:
        set_active_session_id(session_type, selected_id)


def create_and_select_chat_session(session_type: str) -> None:
    new_session_id = create_chat_session(session_type)
    set_active_session_id(session_type, new_session_id)
    st.session_state[get_session_select_key(session_type)] = new_session_id


def delete_and_select_next_chat_session(session_type: str, session_id: str) -> None:
    delete_chat_session(session_id)
    sessions = list_chat_sessions(session_type)
    next_session_id = sessions[0]["id"] if sessions else create_chat_session(session_type)
    set_active_session_id(session_type, next_session_id)
    st.session_state[get_session_select_key(session_type)] = next_session_id


def render_chat_session_controls(session_type: str) -> Tuple[str, List[Dict[str, Any]]]:
    active_id = get_active_session_id(session_type)
    sessions = list_chat_sessions(session_type)
    session_ids = [session["id"] for session in sessions]
    session_by_id = {session["id"]: session for session in sessions}
    if active_id not in session_ids:
        active_id = get_active_session_id(session_type)
        sessions = list_chat_sessions(session_type)
        session_ids = [session["id"] for session in sessions]
        session_by_id = {session["id"]: session for session in sessions}

    select_key = get_session_select_key(session_type)
    if st.session_state.get(select_key) not in session_ids:
        st.session_state[select_key] = active_id

    select_col, new_col, delete_col = st.columns([4, 1, 1])
    with select_col:
        selected_id = st.selectbox(
            "当前会话",
            session_ids,
            index=session_ids.index(active_id),
            format_func=lambda session_id: format_session_option(session_by_id[session_id]),
            key=select_key,
            on_change=sync_selected_chat_session,
            args=(session_type,),
        )
    if selected_id != active_id:
        set_active_session_id(session_type, selected_id)
        active_id = selected_id

    with new_col:
        st.button(
            "新建会话",
            key=f"new_{session_type}_session",
            on_click=create_and_select_chat_session,
            args=(session_type,),
        )

    with delete_col:
        st.button(
            "删除会话",
            key=f"delete_{session_type}_session",
            on_click=delete_and_select_next_chat_session,
            args=(session_type, active_id),
        )

    messages = get_chat_messages(active_id)
    st.caption(f"当前会话共 {len(messages)} 条消息")
    return active_id, messages


def get_file_summary_rows() -> List[Dict[str, Any]]:
    points = qdrant_scroll_points(limit=100000)
    rows_by_key = {}
    for point in points:
        payload = point_payload(point)
        metadata = dict(payload)
        metadata.pop("document", None)
        if not metadata:
            continue
        key = (
            metadata.get("doc_id", ""),
            metadata.get("file_name", ""),
            metadata.get("doc_category", ""),
        )
        row = rows_by_key.setdefault(
            key,
            {
                "文件": metadata.get("file_name", "未知文件"),
                "资料类型": metadata.get("doc_category_name", "未知类型"),
                "来源格式": metadata.get("source_type", "未知"),
                "SHA256": str(metadata.get("file_sha256", ""))[:16],
                "chunk 数": 0,
            },
        )
        row["chunk 数"] += 1
    return list(rows_by_key.values())


def create_vector_library_backup() -> bytes:
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        if os.path.exists(APP_DB_FILE):
            archive.write(APP_DB_FILE, APP_DB_FILE)
        if os.path.isdir(QDRANT_DIR):
            for root, _dirs, files in os.walk(QDRANT_DIR):
                for file_name in files:
                    file_path = os.path.join(root, file_name)
                    archive_name = os.path.relpath(file_path, ".")
                    archive.write(file_path, archive_name)
    return output.getvalue()


def safe_extract_backup(uploaded_backup) -> Tuple[str, bool]:
    temp_dir = tempfile.mkdtemp(prefix="ocr_rag_restore_")
    has_qdrant = False
    has_app_db = False
    try:
        uploaded_backup.seek(0)
    except Exception:
        pass
    with zipfile.ZipFile(uploaded_backup) as archive:
        for member in archive.infolist():
            normalized = member.filename.replace("\\", "/")
            if normalized.endswith("/"):
                continue
            if normalized == APP_DB_FILE:
                has_app_db = True
            elif normalized.startswith(f"{QDRANT_DIR}/"):
                has_qdrant = True
            else:
                continue
            target_path = os.path.abspath(os.path.join(temp_dir, normalized))
            if not target_path.startswith(os.path.abspath(temp_dir) + os.sep):
                raise ValueError("备份文件包含不安全路径，已拒绝导入。")
            os.makedirs(os.path.dirname(target_path), exist_ok=True)
            with archive.open(member) as source, open(target_path, "wb") as target:
                shutil.copyfileobj(source, target)
    if not has_qdrant and not has_app_db:
        raise ValueError("备份包中没有发现 app_state.sqlite3 或 qdrant_db 内容。")
    return temp_dir, has_qdrant


def restore_vector_library_backup(uploaded_backup) -> str:
    restore_dir, has_qdrant = safe_extract_backup(uploaded_backup)
    try:
        try:
            if hasattr(vector_client, "close"):
                vector_client.close()
            load_qdrant_client.clear()
        except Exception:
            pass
        restored_app_db = os.path.join(restore_dir, APP_DB_FILE)
        if os.path.exists(restored_app_db):
            shutil.copy2(restored_app_db, APP_DB_FILE)
        restored_qdrant_dir = os.path.join(restore_dir, QDRANT_DIR)
        if has_qdrant and os.path.isdir(restored_qdrant_dir):
            if os.path.isdir(QDRANT_DIR):
                shutil.rmtree(QDRANT_DIR)
            shutil.copytree(restored_qdrant_dir, QDRANT_DIR)
        return "备份已导入。建议重启 Streamlit，确保 Qdrant 和 SQLite 重新加载。"
    finally:
        shutil.rmtree(restore_dir, ignore_errors=True)


# =========================
# UI 主体
# =========================
tab_upload, tab_search, tab_compliance, tab_config, tab_models, tab_manage = st.tabs(
    ["上传入库", "检索问答", "合规分析", "模型配置", "模型状态", "文档库管理"]
)

with tab_upload:
    st.subheader("上传文件并写入向量库")
    upload_mode_options = ["文件 / 多文件", "文件夹（含子文件夹）"]
    saved_upload_mode = get_config_value("upload_mode_label", upload_mode_options[0])
    if saved_upload_mode not in upload_mode_options:
        saved_upload_mode = upload_mode_options[0]
    upload_mode_label = st.radio(
        "上传方式",
        upload_mode_options,
        index=upload_mode_options.index(saved_upload_mode),
        horizontal=True,
        key="upload_mode_label",
    )
    set_config_value("upload_mode_label", upload_mode_label)

    accept_multiple_files = "directory" if upload_mode_label == "文件夹（含子文件夹）" else True
    upload_button_label = translate_text("选择文件夹" if accept_multiple_files == "directory" else "选择文件")
    st.markdown(
        f"""
        <style>
        section[data-testid="stFileUploaderDropzone"] button [data-testid="stMarkdownContainer"],
        section[data-testid="stFileUploaderDropzone"] button [data-testid="stMarkdownContainer"] * {{
            font-size: 0 !important;
            line-height: 0 !important;
        }}
        section[data-testid="stFileUploaderDropzone"] button [data-testid="stMarkdownContainer"]::after {{
            content: "{upload_button_label}";
            font-size: 1rem;
            line-height: 1.5;
        }}
        section[data-testid="stFileUploaderDropzone"] [data-testid="stFileUploaderFileLimit"] {{
            display: none !important;
        }}
        </style>
        """,
        unsafe_allow_html=True,
    )
    uploaded_items = st.file_uploader(
        "上传文件或文件夹",
        accept_multiple_files=accept_multiple_files,
        help=f"支持格式：{SUPPORTED_TYPE_LABEL}。文件夹模式会包含子文件夹中的文件。",
    )
    uploaded_files = uploaded_items or []
    if not isinstance(uploaded_files, list):
        uploaded_files = [uploaded_files]
    st.caption(f"支持格式：{SUPPORTED_TYPE_LABEL}。不支持的文件会在批量处理结果中列出。")
    if any(uploaded_file and is_legacy_office_file(get_uploaded_relative_name(uploaded_file)) for uploaded_file in uploaded_files):
        soffice_binary = find_soffice_binary()
        if soffice_binary:
            st.caption(f"检测到老版 Office 文件，将使用 LibreOffice 转换：{soffice_binary}")
        else:
            st.warning("检测到老版 Office 文件。当前未找到 LibreOffice；DOC 在 macOS 可尝试用 textutil 转换，PPT/XLS 需要安装 LibreOffice。")

    category_options = list(DOC_CATEGORY_OPTIONS.keys())
    saved_category_label = get_config_value("upload_doc_category_label", category_options[0])
    if saved_category_label not in category_options:
        saved_category_label = category_options[0]
    category_label = st.radio(
        "资料类型",
        category_options,
        index=category_options.index(saved_category_label),
        horizontal=True,
        key="upload_doc_category_label",
    )
    set_config_value("upload_doc_category_label", category_label)
    doc_category = DOC_CATEGORY_OPTIONS[category_label]
    doc_label = st.text_input(
        "资料名称 / 备注",
        value="",
        placeholder="批量上传时留空会使用各自文件名",
    )
    ocr_enhance = st.checkbox(
        "启用 Office 内嵌图片 OCR",
        value=get_bool_config("upload_ocr_enhance", True),
        key="upload_ocr_enhance",
        help="DOCX/PPTX/XLSX 会提取内嵌图片并 OCR。PDF 的 OCR 策略请看下面的 PDF OCR 模式。",
    )
    set_bool_config("upload_ocr_enhance", ocr_enhance)
    auto_unload_models_after_ingest = st.checkbox(
        "入库完成后自动释放 OCR / BGE-M3 模型缓存",
        value=get_bool_config("auto_unload_models_after_ingest", False),
        key="auto_unload_models_after_ingest",
        help="适合内存较小的机器；下次入库或检索会重新加载模型。",
    )
    set_bool_config("auto_unload_models_after_ingest", auto_unload_models_after_ingest)
    replace_changed_same_name = st.checkbox(
        "同名文件变更时替换旧版本",
        value=get_bool_config("replace_changed_same_name", True),
        key="replace_changed_same_name_input",
        help="同一路径文件 SHA256 变化时，会先删除旧 chunk 和去重记录，再写入新版本。",
    )
    set_bool_config("replace_changed_same_name", replace_changed_same_name)
    background_ingest = st.checkbox(
        "后台入库队列",
        value=get_bool_config("background_ingest", DEFAULT_BACKGROUND_INGEST),
        key="background_ingest_input",
        help="提交后由应用内单 worker 后台处理，页面可继续操作，并支持暂停、继续和终止。",
    )
    set_bool_config("background_ingest", background_ingest)
    active_ocr_config = get_paddleocr_model_config()
    st.caption(
        "当前 PaddleOCR 模型："
        f"{get_paddleocr_model_label()}｜{active_ocr_config['det']} / {active_ocr_config['rec']}"
    )

    pdf_ocr_mode_labels = list(PDF_OCR_MODE_OPTIONS.keys())
    saved_pdf_ocr_mode_label = get_config_value("upload_pdf_ocr_mode_label", pdf_ocr_mode_labels[0])
    if saved_pdf_ocr_mode_label not in pdf_ocr_mode_labels:
        saved_pdf_ocr_mode_label = pdf_ocr_mode_labels[0]
    pdf_ocr_mode_label = st.selectbox(
        "PDF OCR 模式",
        pdf_ocr_mode_labels,
        index=pdf_ocr_mode_labels.index(saved_pdf_ocr_mode_label),
        key="upload_pdf_ocr_mode_label",
        help="智能 OCR 只处理文字很少的 PDF 页。强制每页 OCR 最完整，但大 PDF 可能非常慢且占内存。",
    )
    set_config_value("upload_pdf_ocr_mode_label", pdf_ocr_mode_label)
    pdf_ocr_mode = PDF_OCR_MODE_OPTIONS[pdf_ocr_mode_label]
    auto_install_libreoffice = st.checkbox(
        "缺少 LibreOffice 时自动下载安装转换工具",
        value=get_bool_config("auto_install_libreoffice", True),
        key="auto_install_libreoffice",
        help="用于自动转换 DOC/PPT/XLS 老版 Office 文件。Windows/Linux 可能需要管理员权限。",
    )
    set_bool_config("auto_install_libreoffice", auto_install_libreoffice)

    col1, col2 = st.columns(2)
    with col1:
        saved_chunk_size = min(max(get_int_config("upload_chunk_size", 600), 300), 1500)
        chunk_size = st.slider(
            "Chunk 大小",
            min_value=300,
            max_value=1500,
            value=saved_chunk_size,
            step=100,
            key="upload_chunk_size",
        )
        set_config_value("upload_chunk_size", chunk_size)
    with col2:
        overlap_max = min(300, chunk_size - 50)
        saved_overlap = min(max(get_int_config("upload_overlap", 100), 0), overlap_max)
        overlap = st.slider(
            "Chunk 重叠",
            min_value=0,
            max_value=overlap_max,
            value=saved_overlap,
            step=50,
            key="upload_overlap",
        )
        set_config_value("upload_overlap", overlap)

    recent_tasks = list_ingest_tasks(limit=5)
    if recent_tasks:
        with st.expander("最近入库任务", expanded=False):
            st.markdown("##### 任务控制")
            for task in recent_tasks:
                render_ingest_task_controls(task)
            st.markdown("##### 任务列表")
            render_result_dataframe(
                [
                    {
                        "状态": format_task_status(task["status"]),
                        "进度": f"{task['processed_files']}/{task['total_files']}",
                        "成功": task["success_count"],
                        "重复": task["duplicate_count"],
                        "跳过": task["skipped_count"],
                        "失败": task["failed_count"],
                        "当前文件": task["current_file"],
                        "说明": task["message"],
                        "更新时间": time.strftime(
                            "%Y-%m-%d %H:%M:%S",
                            time.localtime(float(task["updated_at"])),
                        ),
                    }
                    for task in recent_tasks
                ]
            )
            latest_task_items = list_ingest_task_items(recent_tasks[0]["id"], limit=100)
            if latest_task_items:
                st.markdown("##### 最近任务文件明细")
                render_result_dataframe(
                    [
                        {
                            "文件": item["file_name"],
                            "状态": item["status"],
                            "说明": item["message"],
                            "chunk 数": item["chunk_count"],
                            "SHA256": str(item["sha256"])[:16],
                            "更新时间": time.strftime(
                                "%Y-%m-%d %H:%M:%S",
                                time.localtime(float(item["updated_at"])),
                            ),
                        }
                        for item in latest_task_items
                    ],
                    max_rows=100,
                )

    if uploaded_files:
        processable_count = sum(1 for item in uploaded_files if is_processable_upload(get_uploaded_relative_name(item)))
        pending_count = len(uploaded_files) - processable_count
        st.info(
            f"已选择 {len(uploaded_files)} 个文件，其中当前可处理 {processable_count} 个，"
            f"待跳过/不支持 {pending_count} 个。开始处理后会按 SHA256 自动跳过重复文件。"
        )

        if st.button("开始导入文件", type="primary", key="upload_ingest"):
            has_legacy_office = any(is_legacy_office_file(get_uploaded_relative_name(item)) for item in uploaded_files)
            if has_legacy_office and not find_soffice_binary() and auto_install_libreoffice:
                with st.status("未检测到 LibreOffice，正在尝试自动下载安装...", expanded=True) as status:
                    install_plan = get_libreoffice_install_plan()
                    st.write(f"当前系统：{install_plan.get('platform', '未知系统')}")
                    st.write(f"安装命令：{install_plan.get('manual', '无可用命令')}")
                    install_ok, install_message = install_libreoffice_automatically()
                    st.write(install_message)
                    if install_ok:
                        status.update(label="LibreOffice 已可用", state="complete")
                    else:
                        status.update(label="LibreOffice 自动安装失败", state="error")

            if background_ingest:
                batch_id = uuid.uuid4().hex
                task_id = create_ingest_task(len(uploaded_files))
                file_specs = []
                for uploaded_file in uploaded_files:
                    relative_name = get_uploaded_relative_name(uploaded_file)
                    file_specs.append(
                        {
                            "relative_name": relative_name,
                            "file_path": save_uploaded_file(uploaded_file, batch_id=batch_id),
                            "sha256": calculate_uploaded_file_sha256(uploaded_file),
                        }
                    )
                runtime = load_ingest_executor()
                future = runtime["executor"].submit(
                    run_background_ingest_task,
                    task_id,
                    file_specs,
                    {
                        "doc_category": doc_category,
                        "doc_label": doc_label,
                        "chunk_size": chunk_size,
                        "overlap": overlap,
                        "ocr_enhance": ocr_enhance,
                        "pdf_ocr_mode": pdf_ocr_mode,
                        "replace_changed_same_name": replace_changed_same_name,
                        "auto_unload_models_after_ingest": auto_unload_models_after_ingest,
                    },
                )
                runtime["futures"][task_id] = future
                update_ingest_task(task_id, message="后台任务已提交")
                st.success("后台入库任务已提交。可在“最近入库任务”里查看进度，页面可以继续操作。")
                st.stop()

            batch_id = uuid.uuid4().hex
            success_rows = []
            unsupported_rows = []
            skipped_rows = []
            failed_rows = []
            duplicate_rows = []
            progress_bar = st.progress(0, text="准备处理文件...")
            status_box = st.empty()
            total_files = len(uploaded_files)
            task_id = create_ingest_task(total_files)

            def sync_ingest_task(processed_files: int, current_file: str, message: str, status: str = "running") -> None:
                update_ingest_task(
                    task_id,
                    status=status,
                    processed_files=processed_files,
                    success_count=len(success_rows),
                    duplicate_count=len(duplicate_rows),
                    skipped_count=len(skipped_rows) + len(unsupported_rows),
                    failed_count=len(failed_rows),
                    current_file=current_file,
                    message=message,
                )

            for file_index, uploaded_file in enumerate(uploaded_files, start=1):
                relative_name = get_uploaded_relative_name(uploaded_file)
                status_box.info(f"正在处理 {file_index}/{total_files}：{relative_name}")
                sync_ingest_task(file_index - 1, relative_name, f"正在处理：{relative_name}")
                progress_bar.progress(
                    (file_index - 1) / total_files,
                    text=f"正在处理 {file_index}/{total_files}：{relative_name}",
                )

                if not is_supported_upload(relative_name):
                    ext = Path(relative_name).suffix.lower() or "无扩展名"
                    unsupported_rows.append(
                        {
                            "文件": relative_name,
                            "原因": f"不支持的文件类型：{ext}",
                        }
                    )
                    progress_bar.progress(file_index / total_files, text=f"已跳过：{relative_name}")
                    sync_ingest_task(file_index, relative_name, f"已跳过不支持文件：{relative_name}")
                    record_ingest_task_item(task_id, relative_name, "unsupported", f"不支持的文件类型：{ext}")
                    continue

                file_sha256 = calculate_uploaded_file_sha256(uploaded_file)
                existing_file = get_duplicate_ingested_file(file_sha256)
                if existing_file:
                    duplicate_rows.append(
                        {
                            "文件": relative_name,
                            "已入库文件": existing_file.get("file_name", ""),
                            "资料类型": DOC_CATEGORY_NAMES.get(
                                existing_file.get("doc_category", ""),
                                existing_file.get("doc_category", ""),
                            ),
                            "chunk 数": existing_file.get("chunk_count", 0),
                            "SHA256": file_sha256[:16],
                        }
                    )
                    progress_bar.progress(file_index / total_files, text=f"已跳过重复文件：{relative_name}")
                    sync_ingest_task(file_index, relative_name, f"已跳过重复文件：{relative_name}")
                    record_ingest_task_item(
                        task_id,
                        relative_name,
                        "duplicate",
                        "已跳过重复文件",
                        chunk_count=int(existing_file.get("chunk_count", 0) or 0),
                        file_sha256=file_sha256,
                    )
                    continue

                if is_legacy_office_file(relative_name):
                    ext = get_file_extension_from_name(relative_name)
                    can_convert, conversion_message = get_legacy_conversion_status(ext)
                    if not can_convert:
                        unsupported_rows.append(
                            {
                                "文件": relative_name,
                                "原因": conversion_message,
                            }
                        )
                        progress_bar.progress(file_index / total_files, text=f"已跳过：{relative_name}")
                        sync_ingest_task(file_index, relative_name, f"已跳过无法转换文件：{relative_name}")
                        record_ingest_task_item(task_id, relative_name, "unsupported", conversion_message, file_sha256=file_sha256)
                        continue

                try:
                    deleted_chunks = 0
                    file_path = save_uploaded_file(uploaded_file, batch_id=batch_id)

                    def extraction_progress(message: str) -> None:
                        status_box.info(f"{relative_name}：{message}")
                        sync_ingest_task(file_index - 1, relative_name, message)

                    sections = extract_document_sections(
                        file_path,
                        ocr_enhance=ocr_enhance,
                        pdf_ocr_mode=pdf_ocr_mode,
                        progress_callback=extraction_progress,
                    )

                    if not sections_have_text(sections):
                        skipped_rows.append(
                            {
                                "文件": relative_name,
                                "原因": "没有解析到有效文字",
                            }
                        )
                        record_ingest_task_item(task_id, relative_name, "skipped", "没有解析到有效文字", file_sha256=file_sha256)
                    else:
                        deleted_chunks, _deleted_records = replace_existing_same_name_if_needed(
                            relative_name,
                            file_sha256,
                            enabled=replace_changed_same_name,
                        )
                        if deleted_chunks:
                            status_box.info(f"已替换旧版本：{relative_name}（删除 {deleted_chunks} 个旧 chunk）")
                        chunk_count = add_document_to_vector_store(
                            file_name=relative_name,
                            file_sha256=file_sha256,
                            sections=sections,
                            chunk_size=chunk_size,
                            overlap=overlap,
                            doc_category=doc_category,
                            doc_label=doc_label or relative_name,
                        )
                        if chunk_count > 0:
                            record_ingested_file(
                                file_sha256=file_sha256,
                                file_name=relative_name,
                                doc_category=doc_category,
                                doc_label=doc_label or relative_name,
                                chunk_count=chunk_count,
                            )
                            success_rows.append(
                                {
                                    "文件": relative_name,
                                    "chunk 数": chunk_count,
                                    "资料类型": category_label,
                                    "SHA256": file_sha256[:16],
                                    "替换旧 chunk": deleted_chunks,
                                }
                            )
                            record_ingest_task_item(
                                task_id,
                                relative_name,
                                "success",
                                f"入库成功，写入 {chunk_count} 个 chunk",
                                chunk_count=chunk_count,
                                file_sha256=file_sha256,
                            )
                        else:
                            skipped_rows.append(
                                {
                                    "文件": relative_name,
                                    "原因": "解析成功但没有可入库 chunk",
                                }
                            )
                            record_ingest_task_item(
                                task_id,
                                relative_name,
                                "skipped",
                                "解析成功但没有可入库 chunk",
                                file_sha256=file_sha256,
                            )
                except Exception as e:
                    failed_rows.append(
                        {
                            "文件": relative_name,
                            "原因": str(e),
                        }
                    )
                    record_ingest_task_item(task_id, relative_name, "failed", str(e), file_sha256=file_sha256)

                progress_bar.progress(file_index / total_files, text=f"已完成：{relative_name}")
                sync_ingest_task(file_index, relative_name, f"已完成：{relative_name}")
                release_memory_after_file()

            status_box.success("批量处理完成")
            progress_bar.progress(1.0, text="批量处理完成")
            sync_ingest_task(total_files, "", "批量处理完成", status="completed")
            if auto_unload_models_after_ingest:
                load_ocr_model.clear()
                load_embedding_model.clear()
                release_memory_after_file()
                st.info("已按设置释放 OCR / BGE-M3 模型缓存。")

            if success_rows:
                st.success(f"成功入库 {len(success_rows)} 个文件。")
                render_result_dataframe(success_rows)
            else:
                st.warning("本次没有文件成功入库。")

            if unsupported_rows:
                st.warning(f"不支持 {len(unsupported_rows)} 个文件。")
                render_result_dataframe(unsupported_rows)

            if duplicate_rows:
                st.info(f"重复文件 {len(duplicate_rows)} 个，已跳过。")
                render_result_dataframe(duplicate_rows)

            if skipped_rows:
                st.info(f"跳过 {len(skipped_rows)} 个文件。")
                render_result_dataframe(skipped_rows)

            if failed_rows:
                st.error(f"处理失败 {len(failed_rows)} 个文件。")
                render_result_dataframe(failed_rows)

with tab_search:
    st.subheader("多轮检索问答")
    render_library_summary()
    rag_session_id, rag_messages = render_chat_session_controls("rag")

    with st.expander("对话与检索设置", expanded=False):
        setting_col1, setting_col2, setting_col3, setting_col4 = st.columns([1.2, 1, 1, 1])
        with setting_col1:
            scope_options = ["全部资料", *DOC_CATEGORY_OPTIONS.keys()]
            saved_scope = get_config_value("rag_search_scope_label", "全部资料")
            if saved_scope not in scope_options:
                saved_scope = "全部资料"
            search_scope_label = st.selectbox(
                "检索范围",
                scope_options,
                index=scope_options.index(saved_scope),
                key="rag_search_scope_label",
            )
            set_config_value("rag_search_scope_label", search_scope_label)
            search_category = None if search_scope_label == "全部资料" else DOC_CATEGORY_OPTIONS[search_scope_label]
        with setting_col2:
            top_k = st.slider(
                "召回片段数量",
                min_value=1,
                max_value=20,
                value=get_int_config("rag_top_k", DEFAULT_RAG_TOP_K),
                step=1,
                key="search_top_k",
                help="普通问答建议 3-5，太大会把弱相关片段带进上下文。",
            )
            set_config_value("rag_top_k", top_k)
        with setting_col3:
            mode_options = list(LLM_MODE_OPTIONS.keys())
            saved_mode_label = get_config_value("rag_mode_label", "快速")
            if saved_mode_label not in mode_options:
                saved_mode_label = "快速"
            mode_label = st.radio(
                "回答模式",
                mode_options,
                index=mode_options.index(saved_mode_label),
                horizontal=True,
                key="chat_llm_mode_label",
            )
            set_config_value("rag_mode_label", mode_label)
            chat_mode = LLM_MODE_OPTIONS[mode_label]
        with setting_col4:
            rag_context_turns = st.slider(
                "上下文轮数",
                min_value=0,
                max_value=20,
                value=get_int_config("rag_context_turns", DEFAULT_CONTEXT_TURNS),
                step=1,
                key="rag_context_turns_input",
                help="只把最近 N 轮对话放进模型上下文；历史仍会保存在数据库。",
            )
            set_config_value("rag_context_turns", rag_context_turns)

        allow_general_fallback = st.checkbox(
            "未检索到资料时使用本地大模型普通聊天",
            value=get_bool_config("rag_allow_general_fallback", True),
            key="rag_allow_general_fallback_input",
            help="当向量库无匹配结果时，允许模型按通用对话方式回复；涉及本地文档的问题仍以入库检索结果为准。",
        )
        set_bool_config("rag_allow_general_fallback", allow_general_fallback)

        query_col, threshold_col, distance_col = st.columns([1.2, 1.2, 1])
        with query_col:
            rag_query_rewrite = st.checkbox(
                "追问补全成完整检索问题",
                value=get_bool_config("rag_query_rewrite", True),
                key="rag_query_rewrite_input",
                help="多轮对话中把追问补全成完整问题后再检索。",
            )
            set_bool_config("rag_query_rewrite", rag_query_rewrite)
            rag_query_decompose = st.checkbox(
                "复杂问题拆解后分别检索",
                value=get_bool_config("rag_query_decompose", DEFAULT_QUERY_DECOMPOSE),
                key="rag_query_decompose_input",
                help="问题包含多个事项时，会拆成子问题分别检索后合并证据。",
            )
            set_bool_config("rag_query_decompose", rag_query_decompose)
        with threshold_col:
            rag_use_distance_threshold = st.checkbox(
                "启用向量距离阈值过滤",
                value=get_bool_config("rag_use_distance_threshold", True),
                key="rag_use_distance_threshold_input",
                help="过滤距离过大的弱相关片段；如果经常召回不到，可调大右侧阈值。",
            )
            set_bool_config("rag_use_distance_threshold", rag_use_distance_threshold)
        with distance_col:
            rag_max_distance = st.slider(
                "最大距离",
                min_value=0.20,
                max_value=2.00,
                value=min(max(get_float_config("rag_max_distance", DEFAULT_VECTOR_MAX_DISTANCE), 0.20), 2.00),
                step=0.05,
                key="rag_max_distance_input",
                disabled=not rag_use_distance_threshold,
            )
            set_config_value("rag_max_distance", rag_max_distance)

        strategy_col1, strategy_col2, strategy_col3 = st.columns([1, 1, 1])
        with strategy_col1:
            rag_use_hybrid = st.checkbox(
                "启用混合检索",
                value=get_bool_config("rag_use_hybrid", DEFAULT_USE_HYBRID_SEARCH),
                key="rag_use_hybrid_input",
                help="同时使用向量语义检索和关键词检索，适合制度编号、部门名称、文件名等精确命中。",
            )
            set_bool_config("rag_use_hybrid", rag_use_hybrid)
        with strategy_col2:
            rag_use_reranker = st.checkbox(
                "启用重排模型",
                value=get_bool_config("rag_use_reranker", DEFAULT_USE_RERANKER),
                key="rag_use_reranker_input",
                help="先多召回，再用 BGE reranker 重排；更准但会增加内存和耗时。",
            )
            set_bool_config("rag_use_reranker", rag_use_reranker)
        with strategy_col3:
            rag_fetch_k = st.slider(
                "候选召回数",
                min_value=top_k,
                max_value=50,
                value=min(max(get_int_config("rag_fetch_k", DEFAULT_RETRIEVAL_FETCH_K), top_k), 50),
                step=1,
                key="rag_fetch_k_input",
                help="用于混合检索和重排的候选数量。",
            )
            set_config_value("rag_fetch_k", rag_fetch_k)

        if st.button("清空当前对话", key="clear_rag_chat"):
            clear_chat_session(rag_session_id)
            st.rerun()

    render_rag_chat_panel(rag_messages)

    with st.form("rag_chat_form", clear_on_submit=True):
        question = st.text_area(
            "输入问题",
            placeholder="输入问题，点击发送",
            height=88,
            key="rag_chat_text",
        )
        send_question = st.form_submit_button("发送", type="primary")

    if send_question:
        question = question.strip()
        if not question:
            st.warning("请输入问题。")
            st.stop()

        chat_history = get_chat_messages(rag_session_id)
        append_chat_message(rag_session_id, "user", question)
        if not chat_history:
            maybe_update_session_title(rag_session_id, question)

        retrieval_query = question
        query_rewrite_error = None
        retrieval_sub_queries = []
        query_decompose_error = None
        try:
            with st.spinner("正在生成回答..."):
                current_count = count_chunks(search_category)
                search_results = []
                if allow_general_fallback and is_likely_general_chat_question(question):
                    answer = ask_llm_general(
                        question,
                        chat_history=chat_history,
                        mode=chat_mode,
                        context_turns=rag_context_turns,
                    )
                elif current_count == 0:
                    if allow_general_fallback:
                        answer = ask_llm_general(
                            question,
                            chat_history=chat_history,
                            mode=chat_mode,
                            context_turns=rag_context_turns,
                        )
                    else:
                        answer = "当前检索范围没有任何入库 chunk，请先上传资料。"
                else:
                    retrieval_query, query_rewrite_error = maybe_rewrite_retrieval_query(
                        question,
                        chat_history=chat_history,
                        enabled=rag_query_rewrite,
                        mode=chat_mode,
                        context_turns=rag_context_turns,
                        purpose="rag",
                    )
                    retrieval_sub_queries, query_decompose_error = decompose_retrieval_query(
                        retrieval_query,
                        enabled=rag_query_decompose,
                        mode=chat_mode,
                        purpose="rag",
                    )
                    search_results = search_vector_store_multi_query(
                        retrieval_sub_queries,
                        top_k=top_k,
                        doc_category=search_category,
                        max_distance=rag_max_distance if rag_use_distance_threshold else None,
                        fetch_k=rag_fetch_k,
                        use_hybrid=rag_use_hybrid,
                        use_reranker=rag_use_reranker,
                    )
                    if not search_results:
                        if allow_general_fallback:
                            answer = ask_llm_general(
                                question,
                                chat_history=chat_history,
                                mode=chat_mode,
                                context_turns=rag_context_turns,
                            )
                        elif rag_use_distance_threshold:
                            answer = "没有检索到满足当前距离阈值的相关内容。可以在检索设置里调大“最大距离”或暂时关闭阈值过滤。"
                        else:
                            answer = "没有检索到相关内容。"
                    else:
                        answer = ask_llm(
                            question,
                            search_results,
                            chat_history=chat_history,
                            mode=chat_mode,
                            context_turns=rag_context_turns,
                        )

            append_chat_message(
                rag_session_id,
                "assistant",
                answer,
                {
                    "search_results": search_results,
                    "retrieval_query": retrieval_query,
                    "retrieval_sub_queries": retrieval_sub_queries,
                    "query_rewrite_error": query_rewrite_error,
                    "query_decompose_error": query_decompose_error,
                    "mode": chat_mode,
                    "mode_label": mode_label,
                },
            )
            st.rerun()
        except Exception as e:
            append_chat_message(
                rag_session_id,
                "assistant",
                f"检索或回答失败：{e}",
                {
                    "search_results": [],
                    "retrieval_query": retrieval_query,
                    "retrieval_sub_queries": retrieval_sub_queries,
                    "query_rewrite_error": query_rewrite_error,
                    "query_decompose_error": query_decompose_error,
                    "mode": chat_mode,
                    "mode_label": mode_label,
                },
            )
            st.rerun()

with tab_compliance:
    st.subheader("多轮合规差距分析")
    render_library_summary()
    compliance_session_id, compliance_messages = render_chat_session_controls("compliance")

    with st.expander("分析与检索设置", expanded=False):
        col_reg, col_ent, col_mode, col_context = st.columns([1, 1, 1, 1])
        with col_reg:
            regulation_top_k = st.slider(
                "召回监管 / 规章片段数量",
                min_value=1,
                max_value=20,
                value=get_int_config("compliance_regulation_top_k", DEFAULT_COMPLIANCE_REGULATION_TOP_K),
                step=1,
                key="regulation_top_k",
            )
            set_config_value("compliance_regulation_top_k", regulation_top_k)
        with col_ent:
            enterprise_top_k = st.slider(
                "召回企业资料片段数量",
                min_value=1,
                max_value=20,
                value=get_int_config("compliance_enterprise_top_k", DEFAULT_COMPLIANCE_ENTERPRISE_TOP_K),
                step=1,
                key="enterprise_top_k",
            )
            set_config_value("compliance_enterprise_top_k", enterprise_top_k)
        with col_mode:
            compliance_mode_options = list(LLM_MODE_OPTIONS.keys())
            saved_compliance_mode_label = get_config_value("compliance_mode_label", "快速")
            if saved_compliance_mode_label not in compliance_mode_options:
                saved_compliance_mode_label = "快速"
            compliance_mode_label = st.radio(
                "分析模式",
                compliance_mode_options,
                index=compliance_mode_options.index(saved_compliance_mode_label),
                horizontal=True,
                key="compliance_llm_mode_label",
            )
            set_config_value("compliance_mode_label", compliance_mode_label)
            compliance_mode = LLM_MODE_OPTIONS[compliance_mode_label]
        with col_context:
            compliance_context_turns = st.slider(
                "上下文轮数",
                min_value=0,
                max_value=20,
                value=get_int_config("compliance_context_turns", DEFAULT_CONTEXT_TURNS),
                step=1,
                key="compliance_context_turns_input",
                help="只把最近 N 轮合规分析对话放进模型上下文；历史仍会保存在数据库。",
            )
            set_config_value("compliance_context_turns", compliance_context_turns)

        query_col, threshold_col, distance_col = st.columns([1.2, 1.2, 1])
        with query_col:
            compliance_query_rewrite = st.checkbox(
                "追问补全成完整检索问题",
                value=get_bool_config("compliance_query_rewrite", True),
                key="compliance_query_rewrite_input",
                help="合规多轮分析中把追问补全成完整检索问题。",
            )
            set_bool_config("compliance_query_rewrite", compliance_query_rewrite)
            compliance_query_decompose = st.checkbox(
                "复杂问题拆解后分别检索",
                value=get_bool_config("compliance_query_decompose", DEFAULT_QUERY_DECOMPOSE),
                key="compliance_query_decompose_input",
                help="复杂合规问题会拆成多个子问题分别检索，再合并监管和企业证据。",
            )
            set_bool_config("compliance_query_decompose", compliance_query_decompose)
        with threshold_col:
            compliance_use_distance_threshold = st.checkbox(
                "启用向量距离阈值过滤",
                value=get_bool_config("compliance_use_distance_threshold", True),
                key="compliance_use_distance_threshold_input",
                help="分别过滤监管资料和企业资料里的弱相关片段。",
            )
            set_bool_config("compliance_use_distance_threshold", compliance_use_distance_threshold)
        with distance_col:
            compliance_max_distance = st.slider(
                "最大距离",
                min_value=0.20,
                max_value=2.00,
                value=min(
                    max(get_float_config("compliance_max_distance", DEFAULT_VECTOR_MAX_DISTANCE), 0.20),
                    2.00,
                ),
                step=0.05,
                key="compliance_max_distance_input",
                disabled=not compliance_use_distance_threshold,
            )
            set_config_value("compliance_max_distance", compliance_max_distance)

        strategy_col1, strategy_col2, strategy_col3 = st.columns([1, 1, 1])
        with strategy_col1:
            compliance_use_hybrid = st.checkbox(
                "启用混合检索",
                value=get_bool_config("compliance_use_hybrid", DEFAULT_USE_HYBRID_SEARCH),
                key="compliance_use_hybrid_input",
                help="同时使用向量语义检索和关键词检索，适合条款号、制度名称、部门名称。",
            )
            set_bool_config("compliance_use_hybrid", compliance_use_hybrid)
        with strategy_col2:
            compliance_use_reranker = st.checkbox(
                "启用重排模型",
                value=get_bool_config("compliance_use_reranker", DEFAULT_USE_RERANKER),
                key="compliance_use_reranker_input",
                help="分别对监管资料和企业资料做候选重排；更准但会增加内存和耗时。",
            )
            set_bool_config("compliance_use_reranker", compliance_use_reranker)
        with strategy_col3:
            max_selected_top_k = max(regulation_top_k, enterprise_top_k)
            compliance_fetch_k = st.slider(
                "候选召回数",
                min_value=max_selected_top_k,
                max_value=50,
                value=min(
                    max(get_int_config("compliance_fetch_k", DEFAULT_RETRIEVAL_FETCH_K), max_selected_top_k),
                    50,
                ),
                step=1,
                key="compliance_fetch_k_input",
                help="用于混合检索和重排的候选数量。",
            )
            set_config_value("compliance_fetch_k", compliance_fetch_k)

        coverage_col1, coverage_col2, coverage_col3, coverage_col4 = st.columns([1, 1, 1, 1])
        with coverage_col1:
            compliance_min_regulation_evidence = st.slider(
                "监管最少证据数",
                min_value=0,
                max_value=regulation_top_k,
                value=min(
                    max(
                        get_int_config(
                            "compliance_min_regulation_evidence",
                            DEFAULT_COMPLIANCE_MIN_REGULATION_EVIDENCE,
                        ),
                        0,
                    ),
                    regulation_top_k,
                ),
                step=1,
                key="compliance_min_regulation_evidence_input",
                help="合规分析会尽量保证监管证据不少于该数量；不足时会放宽距离阈值补齐。",
            )
            set_config_value("compliance_min_regulation_evidence", compliance_min_regulation_evidence)
        with coverage_col2:
            compliance_min_enterprise_evidence = st.slider(
                "企业最少证据数",
                min_value=0,
                max_value=enterprise_top_k,
                value=min(
                    max(
                        get_int_config(
                            "compliance_min_enterprise_evidence",
                            DEFAULT_COMPLIANCE_MIN_ENTERPRISE_EVIDENCE,
                        ),
                        0,
                    ),
                    enterprise_top_k,
                ),
                step=1,
                key="compliance_min_enterprise_evidence_input",
                help="合规分析会尽量保证企业资料证据不少于该数量；不足时会放宽距离阈值补齐。",
            )
            set_config_value("compliance_min_enterprise_evidence", compliance_min_enterprise_evidence)
        with coverage_col3:
            compliance_clause_by_clause = st.checkbox(
                "按监管条款逐条对照",
                value=get_bool_config("compliance_clause_by_clause", False),
                key="compliance_clause_by_clause_input",
                help="适合监管条款较清晰的场景，模型会尽量一条监管要求对应一行分析。",
            )
            set_bool_config("compliance_clause_by_clause", compliance_clause_by_clause)
        with coverage_col4:
            compliance_include_missing_list = st.checkbox(
                "输出资料不足清单",
                value=get_bool_config("compliance_include_missing_list", True),
                key="compliance_include_missing_list_input",
                help="要求模型列出还需要补充哪些企业资料，并可一起导出 Excel。",
            )
            set_bool_config("compliance_include_missing_list", compliance_include_missing_list)

        if st.button("清空合规分析对话", key="clear_compliance_chat"):
            clear_chat_session(compliance_session_id)
            st.rerun()

    render_compliance_chat_panel(compliance_messages)

    with st.form("compliance_chat_form", clear_on_submit=True):
        topic = st.text_area(
            "输入合规分析问题",
            placeholder="例如：数据安全管理制度是否满足监管要求？供应商准入流程有什么合规缺口？",
            height=88,
            key="compliance_chat_text",
        )
        send_compliance = st.form_submit_button("发送", type="primary")

    if send_compliance:
        topic = topic.strip()
        if not topic:
            st.warning("请输入分析主题。")
            st.stop()

        chat_history = get_chat_messages(compliance_session_id)
        append_chat_message(compliance_session_id, "user", topic)
        if not chat_history:
            maybe_update_session_title(compliance_session_id, topic)

        regulation_results = []
        enterprise_results = []
        answer = ""
        structured_rows = []
        retrieval_query = topic
        query_rewrite_error = None
        retrieval_sub_queries = []
        query_decompose_error = None

        if count_chunks("regulation") == 0:
            answer = "请先上传并入库监管要求或规章制度。"
        elif count_chunks("enterprise") == 0:
            answer = "请先上传并入库企业资料。"
        else:
            try:
                with st.spinner("正在检索并生成合规分析..."):
                    retrieval_query, query_rewrite_error = maybe_rewrite_retrieval_query(
                        topic,
                        chat_history=chat_history,
                        enabled=compliance_query_rewrite,
                        mode=compliance_mode,
                        context_turns=compliance_context_turns,
                        purpose="compliance",
                    )
                    retrieval_sub_queries, query_decompose_error = decompose_retrieval_query(
                        retrieval_query,
                        enabled=compliance_query_decompose,
                        mode=compliance_mode,
                        purpose="compliance",
                    )
                    regulation_results = search_with_min_coverage(
                        retrieval_sub_queries,
                        top_k=regulation_top_k,
                        min_results=compliance_min_regulation_evidence,
                        doc_category="regulation",
                        max_distance=compliance_max_distance if compliance_use_distance_threshold else None,
                        fetch_k=compliance_fetch_k,
                        use_hybrid=compliance_use_hybrid,
                        use_reranker=compliance_use_reranker,
                    )
                    enterprise_results = search_with_min_coverage(
                        retrieval_sub_queries,
                        top_k=enterprise_top_k,
                        min_results=compliance_min_enterprise_evidence,
                        doc_category="enterprise",
                        max_distance=compliance_max_distance if compliance_use_distance_threshold else None,
                        fetch_k=compliance_fetch_k,
                        use_hybrid=compliance_use_hybrid,
                        use_reranker=compliance_use_reranker,
                    )

                    if not regulation_results:
                        if compliance_use_distance_threshold:
                            answer = "没有检索到满足当前距离阈值的相关监管要求。可以在分析设置里调大“最大距离”或暂时关闭阈值过滤。"
                        else:
                            answer = "没有检索到相关监管要求。"
                    elif not enterprise_results:
                        if compliance_use_distance_threshold:
                            answer = "没有检索到满足当前距离阈值的相关企业资料。可以在分析设置里调大“最大距离”或暂时关闭阈值过滤。"
                        else:
                            answer = "没有检索到相关企业资料。"
                    else:
                        answer = ask_llm_compliance(
                            topic,
                            regulation_results,
                            enterprise_results,
                            chat_history=chat_history,
                            mode=compliance_mode,
                            context_turns=compliance_context_turns,
                            clause_by_clause=compliance_clause_by_clause,
                            include_missing_list=compliance_include_missing_list,
                        )
                        structured_rows = parse_markdown_table(answer)
            except Exception as e:
                answer = f"合规分析失败：{e}"

        append_chat_message(
            compliance_session_id,
            "assistant",
            answer,
            {
                "regulation_results": regulation_results,
                "enterprise_results": enterprise_results,
                "structured_rows": structured_rows,
                "retrieval_query": retrieval_query,
                "retrieval_sub_queries": retrieval_sub_queries,
                "query_rewrite_error": query_rewrite_error,
                "query_decompose_error": query_decompose_error,
                "mode": compliance_mode,
                "mode_label": compliance_mode_label,
            },
        )
        st.rerun()

with tab_config:
    st.subheader("配置中心")
    reset_notice = st.session_state.pop("app_reset_notice", "")
    if reset_notice:
        st.success(reset_notice)

    config_language_tab, config_model_tab, config_llm_tab, config_reset_tab = st.tabs(
        ["语言设置", "模型与路径", "本地大模型", "初始化"]
    )

    with config_language_tab:
        with st.container(border=True):
            st.markdown("#### 界面语言")
            current_language_code = get_ui_language()
            current_language_label = LANGUAGE_LABEL_BY_CODE.get(current_language_code, "English")
            selected_language_label = st.selectbox(
                "选择界面语言",
                list(LANGUAGE_OPTIONS.keys()),
                index=list(LANGUAGE_OPTIONS.keys()).index(current_language_label),
                key="ui_language_selector",
            )
            st.caption("语言偏好会保存到 app_state.sqlite3，下次打开会自动生效。")
            selected_language_code = LANGUAGE_OPTIONS[selected_language_label]
            if selected_language_code != current_language_code:
                set_config_value("ui_language", selected_language_code)
                st.success("语言设置已保存。")
                st.rerun()

    with config_model_tab:
        ocr_col, path_col = st.columns([1, 2])
        with ocr_col:
            with st.container(border=True):
                st.markdown("#### OCR 模型")
                paddleocr_model_labels = list(PADDLEOCR_MODEL_OPTIONS.keys())
                current_paddleocr_label = get_paddleocr_model_label()
                selected_paddleocr_label = st.selectbox(
                    "PaddleOCR 模型",
                    paddleocr_model_labels,
                    index=paddleocr_model_labels.index(current_paddleocr_label),
                    key="config_paddleocr_model_label",
                    help="Server 精度更高但占用更大；Mobile 占用更低、速度更快。",
                )
                selected_config = PADDLEOCR_MODEL_OPTIONS[selected_paddleocr_label]
                st.caption(f"检测模型：{selected_config['det']}")
                st.caption(f"识别模型：{selected_config['rec']}")
                if selected_paddleocr_label != current_paddleocr_label:
                    save_paddleocr_model_label(selected_paddleocr_label)
                    st.success("PaddleOCR 模型配置已保存，OCR 缓存已清理。")
                    st.rerun()

        with path_col:
            with st.container(border=True):
                st.markdown("#### 模型保存路径")
                st.caption("默认会把 PaddleOCR、BGE-M3 和 Reranker 缓存在项目的 model_cache 目录下；单独指定路径后会优先读取指定路径。")
                with st.form("model_cache_config_form"):
                    model_cache_root = st.text_input(
                        "默认模型根目录",
                        value=get_model_cache_root(),
                        key="model_cache_root_input",
                        help="未单独指定模型路径时，会在该目录下创建各模型缓存目录。",
                    )
                    default_root_for_form = normalize_local_path(model_cache_root, DEFAULT_MODEL_CACHE_ROOT)
                    path_input_col1, path_input_col2 = st.columns(2)
                    with path_input_col1:
                        paddleocr_cache_dir = st.text_input(
                            "PaddleOCR 模型目录",
                            value=get_config_value("paddleocr_cache_dir", ""),
                            key="paddleocr_cache_dir_input",
                            placeholder=os.path.join(default_root_for_form, "paddlex"),
                            help="对应 PADDLE_PDX_CACHE_HOME，PaddleOCR 官方模型会缓存到该目录下的 official_models。",
                        )
                        bge_cache_dir = st.text_input(
                            "BAAI/bge-m3 模型目录",
                            value=get_config_value("bge_cache_dir", ""),
                            key="bge_cache_dir_input",
                            placeholder=os.path.join(default_root_for_form, "bge-m3"),
                            help="BGE-M3 的 sentence-transformers/HuggingFace 缓存目录。",
                        )
                    with path_input_col2:
                        reranker_cache_dir = st.text_input(
                            "BAAI/bge-reranker-v2-m3 模型目录",
                            value=get_config_value("reranker_cache_dir", ""),
                            key="reranker_cache_dir_input",
                            placeholder=os.path.join(default_root_for_form, "bge-reranker-v2-m3"),
                            help="Reranker 的 sentence-transformers/HuggingFace 缓存目录。",
                        )
                        soffice_binary_path = st.text_input(
                            "LibreOffice / soffice 路径",
                            value=get_configured_soffice_path(),
                            key="soffice_binary_path_input",
                            placeholder="留空则自动搜索系统路径",
                            help="LibreOffice 不是模型，这里配置的是 soffice 可执行文件或安装目录。",
                        )

                    path_col_save, path_col_reset = st.columns([1, 1])
                    with path_col_save:
                        save_model_paths = st.form_submit_button("保存模型路径", type="primary")
                    with path_col_reset:
                        reset_model_paths = st.form_submit_button("恢复默认路径")

                if save_model_paths:
                    try:
                        save_model_cache_config(
                            model_cache_root,
                            paddleocr_cache_dir,
                            bge_cache_dir,
                            reranker_cache_dir,
                            soffice_binary_path,
                        )
                        if soffice_binary_path and not find_soffice_binary():
                            st.warning("模型路径已保存，但当前 LibreOffice / soffice 路径未检测到可执行文件。")
                        else:
                            st.success("模型路径已保存，已清理已加载模型缓存；下次加载会使用新路径。")
                    except Exception as e:
                        st.error(f"保存模型路径失败：{e}")

                if reset_model_paths:
                    try:
                        save_model_cache_config(DEFAULT_MODEL_CACHE_ROOT, "", "", "", DEFAULT_SOFFICE_BINARY_PATH)
                        st.success("已恢复默认模型路径。")
                        st.rerun()
                    except Exception as e:
                        st.error(f"恢复默认模型路径失败：{e}")

                with st.expander("当前路径", expanded=False):
                    st.json(
                        {
                            "model_cache_root": get_model_cache_root(),
                            "paddleocr_cache_dir": get_paddleocr_cache_dir(),
                            "bge_cache_dir": get_bge_cache_dir(),
                            "reranker_cache_dir": get_reranker_cache_dir(),
                            "soffice_binary_path": get_configured_soffice_path() or "自动搜索系统路径",
                            "detected_soffice": find_soffice_binary() or "未检测到",
                        }
                    )

    with config_llm_tab:
        current_config = get_llm_config()
        with st.form("llm_config_form"):
            endpoint_col, mode_col = st.columns([1, 1])
            with endpoint_col:
                st.markdown("#### 接口")
                base_url = st.text_input(
                    "OpenAI 兼容接口 Base URL",
                    value=current_config["base_url"],
                    placeholder="例如：http://127.0.0.1:27292/v1",
                )
                api_key = st.text_input(
                    "API Key",
                    value=current_config["api_key"],
                    type="password",
                    placeholder="没有鉴权时可填 EMPTY",
                )
                model = st.text_input(
                    "默认模型名称",
                    value=current_config["model"],
                    placeholder="填写 OLMX / Ollama / LM Studio 中显示的真实模型名",
                )
            with mode_col:
                st.markdown("#### 模式")
                fast_model = st.text_input(
                    "快速模式模型名",
                    value=current_config.get("fast_model", current_config["model"]),
                    placeholder="留空则使用默认模型名称",
                )
                thinking_model = st.text_input(
                    "思考模式模型名",
                    value=current_config.get("thinking_model", current_config["model"]),
                    placeholder="如果后端用不同模型区分快慢，可在这里填思考模型名",
                )
                st.caption("如果后端通过请求参数控制思考模式，可在 extra_body 中配置；不支持时保持 {}。")

            extra_col1, extra_col2 = st.columns(2)
            with extra_col1:
                fast_extra_body = st.text_area(
                    "快速模式 extra_body JSON",
                    value=current_config.get("fast_extra_body", DEFAULT_LLM_EXTRA_BODY),
                    height=90,
                    placeholder='例如：{"enable_thinking": false}',
                )
            with extra_col2:
                thinking_extra_body = st.text_area(
                    "思考模式 extra_body JSON",
                    value=current_config.get("thinking_extra_body", DEFAULT_LLM_EXTRA_BODY),
                    height=90,
                    placeholder='例如：{"enable_thinking": true}',
                )

            col_save, col_reset = st.columns([1, 1])
            with col_save:
                save_config = st.form_submit_button("保存配置", type="primary")
            with col_reset:
                reset_config = st.form_submit_button("恢复默认值")

        if save_config:
            try:
                save_llm_config(
                    base_url,
                    api_key,
                    model,
                    fast_model=fast_model,
                    thinking_model=thinking_model,
                    fast_extra_body=fast_extra_body,
                    thinking_extra_body=thinking_extra_body,
                )
                st.success("配置已保存到本地数据库，并已在当前会话生效。")
            except Exception as e:
                st.error(f"保存失败：{e}")

        if reset_config:
            try:
                save_llm_config(
                    DEFAULT_LLM_BASE_URL,
                    DEFAULT_LLM_API_KEY,
                    DEFAULT_LLM_MODEL,
                    fast_model=DEFAULT_LLM_MODEL,
                    thinking_model=DEFAULT_LLM_MODEL,
                    fast_extra_body=DEFAULT_LLM_EXTRA_BODY,
                    thinking_extra_body=DEFAULT_LLM_EXTRA_BODY,
                )
                st.success("已恢复默认配置。")
                st.rerun()
            except Exception as e:
                st.error(f"恢复失败：{e}")

        active_config = get_llm_config()
        test_col, active_col = st.columns([1, 2])
        with test_col:
            config_test_mode_options = list(LLM_MODE_OPTIONS.keys())
            saved_config_test_mode = get_config_value("config_test_mode_label", "快速")
            if saved_config_test_mode not in config_test_mode_options:
                saved_config_test_mode = "快速"
            test_mode_label = st.radio(
                "测试模式",
                config_test_mode_options,
                index=config_test_mode_options.index(saved_config_test_mode),
                horizontal=True,
                key="config_test_llm_mode_label",
            )
            set_config_value("config_test_mode_label", test_mode_label)
            test_mode = LLM_MODE_OPTIONS[test_mode_label]
            if st.button("测试当前配置", key="test_config_llm"):
                try:
                    with st.status("正在测试本地大模型接口...", expanded=True) as status:
                        test_model, test_extra_body = get_llm_mode_config(test_mode)
                        st.write(f"接口地址：{active_config['base_url']}")
                        st.write(f"模型名称：{test_model}")
                        st.write(f"extra_body：{test_extra_body or {}}")
                        reply = test_llm_connection(mode=test_mode)
                        record_model_event("本地大模型", "完成", reply)
                        status.update(label=f"接口可用：{reply}", state="complete")
                except Exception as e:
                    record_model_event("本地大模型", "失败", str(e))
                    st.error(f"接口测试失败：{e}")
        with active_col:
            with st.expander("当前生效配置", expanded=False):
                st.json(
                    {
                        "base_url": active_config["base_url"],
                        "api_key": "******" if active_config["api_key"] else "",
                        "model": active_config["model"],
                        "fast_model": active_config.get("fast_model", active_config["model"]),
                        "thinking_model": active_config.get("thinking_model", active_config["model"]),
                        "fast_extra_body": active_config.get("fast_extra_body", DEFAULT_LLM_EXTRA_BODY),
                        "thinking_extra_body": active_config.get("thinking_extra_body", DEFAULT_LLM_EXTRA_BODY),
                    }
                )

    with config_reset_tab:
        with st.container(border=True):
            st.markdown("#### 初始化可配置数据")
            st.warning("这里会清空模型配置、界面设置和所有历史会话；不会删除 Qdrant 文档库、上传文件和已入库向量。")
            confirm_reset = st.checkbox("我确认要初始化可配置数据和历史会话", key="confirm_reset_app_state")
            if st.button("初始化可配置数据", type="primary", disabled=not confirm_reset):
                reset_app_state_database()
                st.session_state["app_reset_notice"] = "已初始化配置和历史会话。"
                st.rerun()

with tab_models:
    st.subheader("模型状态 / 下载查询")
    llm_config = get_llm_config()
    active_ocr_config = get_paddleocr_model_config()
    soffice_binary = find_soffice_binary()
    libreoffice_plan = get_libreoffice_install_plan()
    model_rows = [
        {
            "组件": "PaddleOCR",
            "状态": get_paddle_cache_status(),
            "用途": (
                f"图片和扫描 PDF OCR；当前模型："
                f"{active_ocr_config['det']} / {active_ocr_config['rec']}"
            ),
        },
        {
            "组件": EMBEDDING_MODEL_NAME,
            "状态": get_bge_cache_status(),
            "用途": "文本向量化和向量检索",
        },
        {
            "组件": RERANKER_MODEL_NAME,
            "状态": get_reranker_cache_status(),
            "用途": "候选片段重排，提升检索精度",
        },
        {
            "组件": "LibreOffice / soffice",
            "状态": (
                f"已安装：{soffice_binary}"
                if soffice_binary
                else f"未检测到；自动安装方案：{libreoffice_plan.get('manual', '无')}"
            ),
            "用途": "DOC/PPT/XLS 老版 Office 文件转换",
        },
        {
            "组件": f"本地大模型接口：{llm_config['model']}",
            "状态": (
                f"地址：{llm_config['base_url']}；"
                f"快速：{llm_config.get('fast_model', llm_config['model'])}；"
                f"思考：{llm_config.get('thinking_model', llm_config['model'])}"
            ),
            "用途": "问答和合规分析生成",
        },
    ]
    st.dataframe(model_rows, width="stretch")

    st.markdown("### 操作")

    with st.container(border=True):
        st.markdown("#### 内存管理")
        st.caption("处理大文件后可释放本地模型缓存；下次 OCR、向量化或重排会重新加载模型。")
        if st.button("释放 OCR / BGE-M3 / Reranker 模型缓存", key="clear_model_cache"):
            try:
                load_ocr_model.clear()
                load_embedding_model.clear()
                load_reranker_model.clear()
                release_memory_after_file()
                record_model_event("模型缓存", "完成", "已清理 OCR、BGE-M3 和 Reranker 缓存")
                st.success("已释放 OCR / BGE-M3 / Reranker 模型缓存。")
            except Exception as e:
                record_model_event("模型缓存", "失败", str(e))
                st.error(f"释放模型缓存失败：{e}")

    with st.container(border=True):
        st.markdown("#### PaddleOCR")
        if st.button("预加载 PaddleOCR", key="preload_ocr"):
            try:
                with st.status("正在加载 PaddleOCR...", expanded=True) as status:
                    st.write("检查本地缓存...")
                    st.write(get_paddle_cache_status())
                    st.write("加载 OCR 模型...")
                    load_ocr_model()
                    record_model_event("PaddleOCR", "完成", "OCR 模型已可用")
                    status.update(label="PaddleOCR 已可用", state="complete")
            except Exception as e:
                record_model_event("PaddleOCR", "失败", str(e))
                st.error(f"PaddleOCR 加载失败：{e}")

    with st.container(border=True):
        st.markdown("#### BGE-M3")
        if st.button("预加载 BGE-M3", key="preload_bge"):
            try:
                with st.status("正在加载 BGE-M3...", expanded=True) as status:
                    st.write("检查本地缓存...")
                    st.write(get_bge_cache_status())
                    st.write("加载 embedding 模型...")
                    load_embedding_model()
                    st.write("执行一次测试向量化...")
                    embed_texts(["模型预加载测试"])
                    record_model_event(EMBEDDING_MODEL_NAME, "完成", "BGE-M3 已可用")
                    status.update(label="BGE-M3 已可用", state="complete")
            except Exception as e:
                record_model_event(EMBEDDING_MODEL_NAME, "失败", str(e))
                st.error(f"BGE-M3 加载失败：{e}")

    with st.container(border=True):
        st.markdown("#### BGE Reranker")
        st.caption("仅在检索设置里启用重排模型时才会使用；首次加载可能下载模型并占用额外内存。")
        if st.button("预加载 Reranker", key="preload_reranker"):
            try:
                with st.status("正在加载 Reranker...", expanded=True) as status:
                    st.write("检查本地缓存...")
                    st.write(get_reranker_cache_status())
                    st.write("加载重排模型...")
                    reranker = load_reranker_model()
                    st.write("执行一次测试重排...")
                    reranker.predict([["测试问题", "测试片段"]], show_progress_bar=False)
                    record_model_event(RERANKER_MODEL_NAME, "完成", "Reranker 已可用")
                    status.update(label="Reranker 已可用", state="complete")
            except Exception as e:
                record_model_event(RERANKER_MODEL_NAME, "失败", str(e))
                st.error(f"Reranker 加载失败：{e}")

    with st.container(border=True):
        st.markdown("#### Office 老格式转换")
        soffice_binary = find_soffice_binary()
        if soffice_binary:
            st.success(f"LibreOffice 已可用：{soffice_binary}")
            if st.button("测试 LibreOffice", key="test_libreoffice"):
                try:
                    with st.status("正在测试 LibreOffice...", expanded=True) as status:
                        result = run_subprocess([soffice_binary, "--version"], timeout=60)
                        if result.returncode == 0:
                            version_text = (result.stdout or result.stderr or "").strip()
                            record_model_event("LibreOffice", "完成", version_text)
                            status.update(label=f"LibreOffice 可用：{version_text}", state="complete")
                        else:
                            detail = (result.stderr or result.stdout or "").strip()
                            raise RuntimeError(detail or "未知错误")
                except Exception as e:
                    record_model_event("LibreOffice", "失败", str(e))
                    st.error(f"LibreOffice 测试失败：{e}")
        else:
            st.warning("未检测到 LibreOffice。上传 DOC/PPT/XLS 时会尝试自动安装，也可以在这里手动触发。")
            st.caption(f"当前系统安装方案：{libreoffice_plan.get('manual', '无可用方案')}")
            if st.button("自动安装 LibreOffice", key="install_libreoffice"):
                with st.status("正在自动安装 LibreOffice...", expanded=True) as status:
                    st.write(f"当前系统：{libreoffice_plan.get('platform', '未知系统')}")
                    st.write(f"安装命令：{libreoffice_plan.get('manual', '无可用命令')}")
                    install_ok, install_message = install_libreoffice_automatically()
                    st.write(install_message)
                    if install_ok:
                        record_model_event("LibreOffice", "完成", install_message)
                        status.update(label="LibreOffice 已可用", state="complete")
                    else:
                        record_model_event("LibreOffice", "失败", install_message)
                        status.update(label="LibreOffice 自动安装失败", state="error")

    with st.container(border=True):
        st.markdown("#### 本地大模型")
        saved_model_status_test_mode = get_config_value("model_status_test_mode_label", "快速")
        if saved_model_status_test_mode not in LLM_MODE_OPTIONS:
            saved_model_status_test_mode = "快速"
        model_test_mode_label = st.radio(
            "测试模式",
            list(LLM_MODE_OPTIONS.keys()),
            index=list(LLM_MODE_OPTIONS.keys()).index(saved_model_status_test_mode),
            horizontal=False,
            key="model_status_test_llm_mode_label",
        )
        set_config_value("model_status_test_mode_label", model_test_mode_label)
        model_test_mode = LLM_MODE_OPTIONS[model_test_mode_label]
        if st.button("测试本地大模型", key="test_llm"):
            try:
                with st.status("正在测试本地大模型接口...", expanded=True) as status:
                    active_config = get_llm_config()
                    test_model, test_extra_body = get_llm_mode_config(model_test_mode)
                    st.write(f"接口地址：{active_config['base_url']}")
                    st.write(f"模型名称：{test_model}")
                    st.write(f"extra_body：{test_extra_body or {}}")
                    reply = test_llm_connection(mode=model_test_mode)
                    record_model_event("本地大模型", "完成", reply)
                    status.update(label=f"本地大模型可用：{reply}", state="complete")
            except Exception as e:
                record_model_event("本地大模型", "失败", str(e))
                st.error(f"本地大模型测试失败：{e}")

    if st.session_state["model_events"]:
        st.markdown("### 最近模型事件")
        st.dataframe(st.session_state["model_events"][-20:], width="stretch")

with tab_manage:
    st.subheader("文档库管理")
    st.write(translate_text("当前 Qdrant Collection："), COLLECTION_NAME)
    render_library_summary()

    if st.button("刷新文件摘要", key="refresh_summary"):
        st.rerun()

    summary_rows = get_file_summary_rows()
    if summary_rows:
        st.dataframe(summary_rows, width="stretch")
    else:
        st.info("当前文档库为空。")

    ingested_rows = [
        {
            "文件": item["file_name"],
            "资料类型": DOC_CATEGORY_NAMES.get(item["doc_category"], item["doc_category"]),
            "资料名称": item["doc_label"],
            "chunk 数": item["chunk_count"],
            "SHA256": item["sha256"][:16],
            "入库时间": time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(float(item["created_at"]))),
        }
        for item in list_ingested_files()
    ]
    with st.expander("查看去重记录", expanded=False):
        if ingested_rows:
            st.dataframe(ingested_rows, width="stretch")
        else:
            st.info("暂无去重记录。")

    with st.expander("向量库备份 / 导入 / 导出", expanded=False):
        st.caption("导出会包含 Qdrant 向量库和 app_state.sqlite3；不包含 uploads 原始文件。")
        try:
            backup_bytes = create_vector_library_backup()
            st.download_button(
                "导出文档库备份",
                data=backup_bytes,
                file_name=f"ocr_rag_backup_{time.strftime('%Y%m%d_%H%M%S')}.zip",
                mime="application/zip",
                key="download_vector_backup",
            )
        except Exception as e:
            st.error(f"生成备份失败：{e}")

        backup_file = st.file_uploader("导入备份 ZIP", type=["zip"], key="restore_backup_zip")
        confirm_restore = st.checkbox("我确认导入备份并覆盖当前文档库", key="confirm_restore_backup")
        if st.button("导入并覆盖当前文档库", disabled=not backup_file or not confirm_restore, key="restore_vector_backup"):
            try:
                message = restore_vector_library_backup(backup_file)
                st.success(message)
            except Exception as e:
                st.error(f"导入备份失败：{e}")

    st.warning("清空向量库会删除所有已入库 chunk 和 SHA256 去重记录，原始上传文件不会删除。")
    if st.button("清空 Qdrant 向量库", key="clear_qdrant"):
        try:
            deleted_count = count_chunks()
            if deleted_count:
                recreate_qdrant_collection()
                delete_all_ingested_file_records()
                st.success(f"已删除 {deleted_count} 个 chunk。")
            else:
                delete_all_ingested_file_records()
                st.info("当前向量库为空。")
        except Exception as e:
            st.error(f"清空失败：{e}")
